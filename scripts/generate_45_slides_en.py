"""Comprehensive 45-Slide Graduation Defense Presentation in English for VlegalAI.
Version with LARGER FONTS, LARGER DIAGRAMS, AND ENHANCED READABILITY.
Sourced directly from Codebase and Final Capstone Report.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(r"f:\VlegalAI")
OUT_SLIDES = ROOT / "slides.pptx"
OUT_SLDIES = ROOT / "sldies.pptx"
PNG_DIR = ROOT / "diagramv2" / "png"
SYS_DESIGN_PNG = ROOT / "01_system_design.png"
EVAL_PNG = ROOT / "evaluation" / "benchmarks" / "ragas-gemini-100" / "architecture_comparison.png"

# High-Contrast Colors
DARK_BG = RGBColor(15, 23, 42)       # #0F172A Slate Dark
CARD_DARK = RGBColor(30, 41, 59)     # #1E293B Slate Navy
TEXT_LIGHT = RGBColor(248, 250, 252) # #F8FAFC
TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A

EMERALD = RGBColor(16, 185, 129)     # #10B981 Vibrant Green
DARK_EMERALD = RGBColor(15, 76, 58)  # #0F4C3A Primary Green
TEAL = RGBColor(14, 165, 233)        # #0EA5E9 Sky Blue
PURPLE = RGBColor(139, 92, 246)      # #8B5CF6 Purple
AMBER = RGBColor(217, 119, 6)        # #D97706 Amber
ROSE = RGBColor(225, 29, 72)         # #E11D48 Crimson Rose

BG_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC
WHITE = RGBColor(255, 255, 255)
BORDER_LIGHT = RGBColor(203, 213, 225) # #CBD5E1
SOFT_GREEN = RGBColor(236, 253, 245)   # #ECFDF5
SOFT_BLUE = RGBColor(240, 249, 255)    # #F0F9FF
SOFT_AMBER = RGBColor(254, 243, 199)   # #FEF3C7
SOFT_PURPLE = RGBColor(245, 243, 255)  # #F5F3FF

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, section: str, title: str, page_num: int, total_pages: int = 45):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.fill.background()

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.48))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = SOFT_GREEN
        top_bar.line.fill.background()

        p_sec = top_bar.text_frame.paragraphs[0]
        p_sec.text = f"  VLEGALAI | {section.upper()}"
        p_sec.font.size = Pt(12)
        p_sec.font.bold = True
        p_sec.font.color.rgb = DARK_EMERALD

        tb_page = slide.shapes.add_textbox(Inches(11.3), Inches(0.04), Inches(1.7), Inches(0.4))
        p_pg = tb_page.text_frame.paragraphs[0]
        p_pg.text = f"{page_num:02d} / {total_pages:02d}"
        p_pg.alignment = PP_ALIGN.RIGHT
        p_pg.font.size = Pt(12)
        p_pg.font.bold = True
        p_pg.font.color.rgb = TEXT_MUTED

        tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12.133), Inches(0.65))
        p_t = tb_title.text_frame.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(21)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.22), Inches(1.5), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = EMERALD
        line.line.fill.background()

        fline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.02), Inches(12.133), Inches(0.02))
        fline.fill.solid()
        fline.fill.fore_color.rgb = BORDER_LIGHT
        fline.line.fill.background()

        tb_f = slide.shapes.add_textbox(Inches(0.6), Inches(7.06), Inches(8.5), Inches(0.35))
        pf = tb_f.text_frame.paragraphs[0]
        pf.text = "VlegalAI: Vietnamese Labor-Law GraphRAG System | Capstone Defense"
        pf.font.size = Pt(10)
        pf.font.color.rgb = TEXT_MUTED

        tb_fr = slide.shapes.add_textbox(Inches(9.5), Inches(7.06), Inches(3.2), Inches(0.35))
        pfr = tb_fr.text_frame.paragraphs[0]
        pfr.text = f"Slide {page_num:02d}"
        pfr.alignment = PP_ALIGN.RIGHT
        pfr.font.size = Pt(10)
        pfr.font.bold = True
        pfr.font.color.rgb = TEXT_MUTED

    def add_card(slide, title: str, text: str, left, top, width, height, accent_color=DARK_EMERALD, bg_color=WHITE, title_size=13.5, body_size=11.5):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = BORDER_LIGHT
        box.line.width = Pt(1)

        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.1), Inches(height))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent_color
        stripe.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.1), Inches(width - 0.28), Inches(height - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(title_size)
        p0.font.bold = True
        p0.font.color.rgb = accent_color

        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.size = Pt(body_size)
        p1.font.color.rgb = TEXT_DARK

    def add_stat_box(slide, value: str, label: str, left, top, width, height=0.95, color=DARK_EMERALD, val_size=21, lbl_size=10):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER_LIGHT
        box.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.06), Inches(width), Inches(height - 0.12))
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = value
        p0.alignment = PP_ALIGN.CENTER
        p0.font.size = Pt(val_size)
        p0.font.bold = True
        p0.font.color.rgb = color

        p1 = tf.add_paragraph()
        p1.text = label
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(lbl_size)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_MUTED

    def add_image_safe(slide, img_path: Path, left, top, width, height):
        if img_path.exists():
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = BORDER_LIGHT
            box.line.width = Pt(1.5)
            pad = 0.06
            slide.shapes.add_picture(str(img_path), Inches(left + pad), Inches(top + pad), Inches(width - 2*pad), Inches(height - 2*pad))
        else:
            add_card(slide, "Diagram Placeholder", f"Image file not found: {img_path.name}", left, top, width, height, ROSE)

    # =========================================================================
    # 45 SLIDES - LARGE FONTS & LARGE VISUALS
    # =========================================================================

    # Slide 1: Cover
    s = prs.slides.add_slide(blank_layout)
    bg1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EMERALD
    bar.line.fill.background()

    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.65), Inches(5.2), Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_DARK
    pill.line.color.rgb = EMERALD
    p_pill = pill.text_frame.paragraphs[0]
    p_pill.text = "CAPSTONE DEFENSE 2026 - SE / AI"
    p_pill.alignment = PP_ALIGN.CENTER
    p_pill.font.size = Pt(12)
    p_pill.font.bold = True
    p_pill.font.color.rgb = EMERALD

    tb_c1 = s.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.8), Inches(1.8))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    p_c1 = tf_c1.paragraphs[0]
    p_c1.text = "VLEGALAI: EVIDENCE-GROUNDED LABOR LAW AI"
    p_c1.font.size = Pt(32)
    p_c1.font.bold = True
    p_c1.font.color.rgb = WHITE
    p_c2 = tf_c1.add_paragraph()
    p_c2.text = "A Production-Grade Vietnamese Labor-Law GraphRAG System with Adaptive Hybrid Retrieval"
    p_c2.font.size = Pt(19)
    p_c2.font.color.rgb = TEAL

    box_sv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.4), Inches(5.7), Inches(2.5))
    box_sv.fill.solid()
    box_sv.fill.fore_color.rgb = CARD_DARK
    box_sv.line.color.rgb = BORDER_LIGHT
    tf_sv = box_sv.text_frame
    tf_sv.word_wrap = True
    p_sv0 = tf_sv.paragraphs[0]
    p_sv0.text = "STUDENT TEAM:"
    p_sv0.font.size = Pt(13)
    p_sv0.font.bold = True
    p_sv0.font.color.rgb = EMERALD
    p_sv1 = tf_sv.add_paragraph()
    p_sv1.text = "• Tran Tuan Kiet (Lead / AI Backend) - QE180152\n• Le Thanh Dat (Product / Frontend) - QE170186\n• Phan Bao Khanh (Quality / Evaluation) - DE170648"
    p_sv1.font.size = Pt(13.5)
    p_sv1.font.color.rgb = TEXT_LIGHT

    box_gv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.4), Inches(5.7), Inches(2.5))
    box_gv.fill.solid()
    box_gv.fill.fore_color.rgb = CARD_DARK
    box_gv.line.color.rgb = BORDER_LIGHT
    tf_gv = box_gv.text_frame
    tf_gv.word_wrap = True
    p_gv0 = tf_gv.paragraphs[0]
    p_gv0.text = "PROJECT SUPERVISORS:"
    p_gv0.font.size = Pt(13)
    p_gv0.font.bold = True
    p_gv0.font.color.rgb = AMBER
    p_gv1 = tf_gv.add_paragraph()
    p_gv1.text = "• Supervisor: MSc. Le Trung Hieu\n• Co-supervisor: MSc. Truong Ngoc Hung\n• Institution: FPT University Quy Nhon (AIP491)"
    p_gv1.font.size = Pt(13.5)
    p_gv1.font.color.rgb = TEXT_LIGHT

    tb_foot = s.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.5))
    p_foot = tb_foot.text_frame.paragraphs[0]
    p_foot.text = "August 2026 | Major: Artificial Intelligence & Software Engineering"
    p_foot.font.size = Pt(11.5)
    p_foot.font.color.rgb = TEXT_MUTED

    # Slide 2: Team Organization
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Team Organization", "Team Responsibility Assignment Matrix (RAM) & Governance", 2)
    add_card(s, "1. Tran Tuan Kiet (Lead / AI Backend)",
             "• System architecture, deterministic parsing, hierarchy-aware chunking & multi-store indexing.\n"
             "• Implementation of pgvector HNSW, Neo4j GraphRAG, weighted RRF & adaptive routing.\n"
             "• Cloud Run service orchestration, FastAPI async backend, Docker & PostgreSQL.",
             0.6, 1.4, 3.85, 4.5, DARK_EMERALD, WHITE, 14, 11.5)
    add_card(s, "2. Le Thanh Dat (Product / Frontend)",
             "• React 18 + Vite SPA design, responsive UI/UX & WebSocket/SSE chat streaming.\n"
             "• Interactive contract workspace: Drafting, Risk Review, and Version Comparison.\n"
             "• Citation Drawer UI [S1-Sn], fulltext provision preview & user session management.",
             4.75, 1.4, 3.85, 4.5, TEAL, WHITE, 14, 11.5)
    add_card(s, "3. Phan Bao Khanh (Quality & Evaluation)",
             "• RAGAS benchmark design across 100 realistic legal questions (Single-hop / Multi-hop).\n"
             "• Test automation suite: 492 backend unit tests, 16 frontend tests, coverage analysis.\n"
             "• Paired bootstrap 95% CI statistical analysis, error taxonomy & capstone documentation.",
             8.9, 1.4, 3.85, 4.5, PURPLE, WHITE, 14, 11.5)
    add_card(s, "Equal Contribution & Quality Commitment",
             "Demonstrated through clear module ownership, automated CI/CD gating, and 100% reproducible evidence artifacts.",
             0.6, 6.05, 12.1, 0.85, DARK_EMERALD, SOFT_GREEN, 12, 10.5)

    # Slide 3: Scope & Compliance
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Scope & Compliance", "Project Scope Baseline & Capstone Rubric Traceability", 3)
    add_stat_box(s, "74", "LEGAL INSTRUMENTS", 0.6, 1.4, 2.2, 0.95, DARK_EMERALD, 24, 10.5)
    add_stat_box(s, "32,334", "RETRIEVAL CHUNKS", 3.05, 1.4, 2.2, 0.95, TEAL, 22, 10.5)
    add_stat_box(s, "29,575", "GRAPH NODES", 5.5, 1.4, 2.2, 0.95, PURPLE, 22, 10.5)
    add_stat_box(s, "108,368", "TYPED EDGES", 7.95, 1.4, 2.3, 0.95, AMBER, 22, 10.5)
    add_stat_box(s, "492 / 492", "BACKEND TESTS PASS", 10.5, 1.4, 2.2, 0.95, ROSE, 22, 10.5)
    add_card(s, "Primary Legal Corpus Scope",
             "• Labor Code 45/2019/QH14 and relevant statutory acts (Social Insurance, Employment, OSH).\n"
             "• Implementing Decrees (Decree 145/2020/ND-CP, Decree 293/2025/ND-CP regional wages).\n"
             "• Ministerial Circulars (BLDTBXH) and consolidated legal instruments (VBHN).",
             0.6, 2.5, 5.9, 3.4, DARK_EMERALD, WHITE, 14, 11.5)
    add_card(s, "Capstone Checklist Traceability (AI / SE Standards)",
             "• AI1-AI5: Structured legal knowledge graph, vector index, and RAGAS comparative evaluation.\n"
             "• SE1-SE5: 4-tier layered architecture, 23 SQL tables, 18 Alembic migrations, OIDC PKCE security.\n"
             "• PE1-PE3: Complete reproducible release bundle, test catalogue, runbooks, and residual risk disclosure.",
             6.8, 2.5, 5.9, 3.4, TEAL, WHITE, 14, 11.5)
    add_card(s, "System Boundary Disclaimer",
             "VlegalAI is an intelligent decision-support assistant and does not substitute for certified legal counsel.",
             0.6, 6.05, 12.1, 0.85, AMBER, SOFT_AMBER, 12, 10.5)

    # Slide 4: Vietnamese Legal Information Challenges
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Problem Framing", "Four Fundamental Challenges in Vietnamese Legal Question Answering", 4)
    add_card(s, "1. Statutory Structural Hierarchy",
             "Answers are governed by strict hierarchy (Instrument > Chapter > Section > Article > Clause > Point). Arbitrary character chunking breaks conditions from governing actors.",
             0.6, 1.4, 5.9, 2.15, ROSE, WHITE, 13.5, 11)
    add_card(s, "2. Cross-Referencing & Exception Rules",
             "Legal rules span multiple articles (e.g., 'Except as provided in Clause 2 Article 156...'). Standard LLMs fail to follow statutory references.",
             6.8, 1.4, 5.9, 2.15, AMBER, WHITE, 13.5, 11)
    add_card(s, "3. Temporal Validity & Supersession",
             "Legislation evolves frequently. Answers based on superseded decrees cause severe compliance risks if version provenance is not tracked.",
             0.6, 3.7, 5.9, 2.15, PURPLE, WHITE, 13.5, 11)
    add_card(s, "4. Noisy & Compound Queries",
             "Users write colloquial Vietnamese, abbreviations (HĐLĐ, BHXH), teencode, or compound multi-issue questions requiring facet decomposition.",
             6.8, 3.7, 5.9, 2.15, TEAL, WHITE, 13.5, 11)
    add_card(s, "Core Design Insight",
             "Legal QA requires Grounded Retrieval and Constrained Generation rather than unconstrained language model completion.",
             0.6, 6.0, 12.1, 0.9, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    # Slide 5: System Context & Trust Boundary Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "System Boundary", "System Context, Trust Boundaries & External Integrations", 5)
    add_card(s, "System Context & Trust Boundaries",
             "• Authenticated User: Web SPA client with Google OAuth2 OIDC + PKCE.\n"
             "• Application Core: Cloud Run FastAPI handling chat, contract tools & legal search.\n"
             "• Data Stores: Cloud SQL PostgreSQL (pgvector + GIN) & Neo4j Graph.\n"
             "• Hosted AI: Google Vertex AI (Gemini 2.5 Flash & gemini-embedding-001).\n"
             "• Confidentiality: Private attachments are AES-GCM encrypted and never indexed in public stores.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "01-system-context.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 6: Research Objectives & Formulation
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Research Objectives", "Research Objectives & Design-Science Formulation", 6)
    add_card(s, "Formal Problem Statement",
             "Given an authenticated user's Vietnamese question and optional private attachment:\n"
             "1. Retrieve the smallest sufficient set of governing legal evidence.\n"
             "2. Answer every requested facet without omitting minority conditions.\n"
             "3. Expose verifiable citations [S1-Sn] grounded in active legislation.\n"
             "4. Protect user confidentiality and maintain low latency.",
             0.6, 1.4, 5.9, 4.4, DARK_EMERALD, WHITE, 14, 11.5)
    add_card(s, "Measurable Engineering Objectives",
             "• O1 (Grounded Quality): Faithfulness >= 0.90 across evaluation benchmarks.\n"
             "• O2 (Latency & Efficiency): Single-hop retrieval < 20ms; adaptive graph activation.\n"
             "• O3 (Provenance & Safety): 100% stable citation paths from Answer -> Evidence -> Chunk -> Article.\n"
             "• O4 (Reproducibility): Software-level test reproducibility (492/492 passed).",
             6.8, 1.4, 5.9, 4.4, TEAL, WHITE, 14, 11.5)
    add_card(s, "Methodological Cycle",
             "Design-Science Research: Problem Analysis -> Requirement Specification -> Artifact Construction -> Empirical Benchmark -> Release Packaging.",
             0.6, 5.95, 12.1, 0.95, PURPLE, SOFT_PURPLE, 12, 10.5)

    # Slide 7: Research Questions & Hypotheses
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Research Questions", "Four Core Research Questions (RQ1 - RQ4) & Experimental Hypotheses", 7)
    add_card(s, "RQ1: Architecture Quality Comparison",
             "How do Dense RAG, LightRAG, and GraphRAG differ in Faithfulness, Answer Relevancy, and Factual Correctness on Vietnamese labor law QA?",
             0.6, 1.4, 5.9, 2.15, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "RQ2: Latency & Computational Trade-offs",
             "What is the measured retrieval latency overhead of graph expansion relative to dense vector retrieval, and does it justify universal use?",
             6.8, 1.4, 5.9, 2.15, TEAL, WHITE, 13.5, 11)
    add_card(s, "RQ3: Impact of Reasoning Depth",
             "As question complexity scales from Single-hop to Multi-hop Specific and Abstract, which architecture maintains information coverage best?",
             0.6, 3.7, 5.9, 2.15, PURPLE, WHITE, 13.5, 11)
    add_card(s, "RQ4: Reproducibility & Integrity",
             "How can legal RAG systems be validated transparently without fabricated scores or leakage between test and indexing snapshots?",
             6.8, 3.7, 5.9, 2.15, AMBER, WHITE, 13.5, 11)
    add_card(s, "Resulting Architecture Principle",
             "Empirical evidence for RQ1-RQ4 directly motivated Adaptive Routing: Fast Hybrid RAG for single-hop and Graph expansion only for multi-hop queries.",
             0.6, 6.0, 12.1, 0.9, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    # Slide 8: Related Work & Research Lineage
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Literature Review", "Related Work & Adopted Design Decisions", 8)
    add_card(s, "1. Legal QA & Benchmarks (LegalBench, LexGLUE)",
             "• Guha et al. (2023), Chalkidis et al. (2022): Unassisted LLMs struggle with statutory rule adherence.\n"
             "• Adopted: Treat legal QA primarily as a grounding, provenance, and citation verification problem.",
             0.6, 1.4, 5.9, 2.6, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. Dense Retrieval & Hybrid Search (DPR, BM25, RRF)",
             "• Karpukhin et al. (2020), Cormack et al. (2009): Dense vectors miss exact law codes and numeric thresholds.\n"
             "• Adopted: Weighted Reciprocal Rank Fusion combining pgvector (0.55) and GIN BM25 (0.45) with K=60.",
             6.8, 1.4, 5.9, 2.6, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. GraphRAG & Vietnamese Legal KG (Microsoft GraphRAG, LightRAG, Vuong et al.)",
             "• Edge et al. (2024), Guo et al. (2024), Vuong et al. (2023): Graph structures improve relational recall.\n"
             "• Adopted: 10-Layer Legal Knowledge Graph with 2-hop bounded expansion terminating at verified LegalChunks.",
             0.6, 4.15, 12.1, 2.7, PURPLE, SOFT_BLUE, 13.5, 11)

    # Slide 9: 10-Phase Engineering Evolution
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Engineering Evolution", "Ten-Phase Engineering Evolution & Design-Science Loop", 9)
    p1 = [
        ("Phase 1: Raw LLM Evaluation", "Observed hallucinations -> Reframed as evidence grounding problem."),
        ("Phase 2: Fine-Tuning Study", "High update cost, loss of provenance -> Externalized knowledge in DB."),
        ("Phase 3: Dense RAG Baseline", "Missed article numbers and money amounts -> Added BM25/GIN branch."),
        ("Phase 4: Ranking Calibration", "Incompatible score scales -> Implemented Weighted RRF (K=60)."),
        ("Phase 5: Statutory Chunking", "Arbitrary splits broke clauses -> Created Hierarchy-Aware Chunking.")
    ]
    p2 = [
        ("Phase 6: Graph Latency Trade-off", "Neo4j overhead on simple queries -> Built Adaptive Complexity Routing."),
        ("Phase 7: Query Noise Handling", "Teencode & compound queries -> Gated rewrite & Facet Planning."),
        ("Phase 8: Failure Transparency", "Distinguish quota vs missing law -> Implemented Evidence Gate."),
        ("Phase 9: Production Governance", "Data drift & privacy -> Added Cloud Run, WIF, 18 Migrations, Cache."),
        ("Phase 10: Defensible Release", "Avoid vanity metrics -> 100-Question RAGAS Benchmark & 492 Tests.")
    ]
    y_pos = 1.4
    for title, desc in p1:
        add_card(s, title, desc, 0.6, y_pos, 5.9, 0.98, DARK_EMERALD, WHITE, 12, 10)
        y_pos += 1.08
    y_pos = 1.4
    for title, desc in p2:
        add_card(s, title, desc, 6.8, y_pos, 5.9, 0.98, TEAL, WHITE, 12, 10)
        y_pos += 1.08

    # Slide 10: End-to-End User Journey & Workflows Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "User Experience", "End-to-End User Journey & Core Product Workflows", 10)
    add_card(s, "Core User Capabilities",
             "• Legal Q&A: Natural language chat with verifiable [S1-Sn] citations.\n"
             "• Contract Drafting: Standard labor contract generation with auto regional wages.\n"
             "• Risk Review: Automated scanning of uploaded contracts for adverse clauses.\n"
             "• Version Compare: Side-by-side clause alignment and delta risk scoring.\n"
             "• HITL Feedback: Users submit GOOD/BAD ratings and trigger answer regenerations.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "08-end-to-end-user-journey.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 11: 01_system_design.png (System Architecture Overview - Large Display)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "System Architecture", "System Design & End-to-End Architecture Overview", 11)
    add_card(s, "System Architecture Core Principles",
             "• Client SPA (React 18 + Vite)  |  API Gateway (FastAPI Cloud Run)  |  Storage (Cloud SQL pgvector + Neo4j)  |  AI Engine (Vertex AI Gemini 2.5 Flash).\n"
             "• Strict separation between immutable inference models and versioned external legal knowledge stores.",
             0.6, 1.35, 12.1, 1.05, DARK_EMERALD, SOFT_GREEN, 13, 11)
    add_image_safe(s, SYS_DESIGN_PNG, 0.6, 2.5, 12.1, 4.4)

    # Slide 12: Stakeholders & Business Pain Points
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Target Users", "Target Stakeholders & Solved Pain Points", 12)
    add_card(s, "1. Employees (Laborers)",
             "• Pain: Inability to decipher complex legal terminology, unawareness of statutory benefits (maternity, severance, overtime).\n"
             "• Solution: Free natural language QA with transparent citation links to verify statutory rights.",
             0.6, 1.4, 5.9, 2.6, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. HR Specialists & SMB Employers",
             "• Pain: High risk of unlawful contract terms, compliance penalties due to outdated wage scale knowledge.\n"
             "• Solution: Automated contract drafting, risk scanning, and instant regional minimum wage validation.",
             6.8, 1.4, 5.9, 2.6, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. Legal Compliance Officers",
             "• Pain: Time-consuming cross-referencing between primary laws and implementing decrees.\n"
             "• Solution: 10-Layer knowledge graph visualizing cross-references and version lineage in seconds.",
             0.6, 4.15, 12.1, 2.7, PURPLE, SOFT_BLUE, 13.5, 11)

    # Slide 13: 10-Layer Knowledge Base Model Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Knowledge Base", "Ten-Layer Legal Knowledge Base Architecture (L0 - L9)", 13)
    add_card(s, "Layered Architecture",
             "• L0-L1 (Foundation): Instrument provenance, validity intervals, and document hierarchy.\n"
             "• L2-L7 (Domain Semantics): Defined terms, wages, actors, procedures, timelines, sanctions.\n"
             "• L8-L9 (Cross-Document): Contract lifecycle and case law interpretations.\n"
             "• LegalChunk Bridge: Every graph path must resolve to a valid LegalChunk.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "25-knowledge-base-layers.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 14: Data Pipeline Diagram (Ingestion & Sync)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Data Pipeline", "Legal Data Ingestion, Embedding & Multi-Store Synchronization", 14)
    add_card(s, "Data Ingestion Pipeline",
             "1. Deterministic Extraction: Ordered DOCX parser extracting paragraphs and tables.\n"
             "2. Hierarchy-Aware Chunking: Bounded 360-word windows with 70-word overlap.\n"
             "3. Vector Embedding: gemini-embedding-001 (1024D) with SHA-256 caching.\n"
             "4. Atomic Activation: Simultaneous commit to SQLite, PostgreSQL, and Neo4j.\n"
             "5. Checkpointing (640 chunks/batch) enables seamless resumption upon provider quota limits.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "11-legal-data-pipeline.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 15: Chat Request Sequence Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Runtime Execution", "Chat Request Lifecycle, Evidence Gating & Citation Validation", 15)
    add_card(s, "Runtime Sequence Flow",
             "1. User Query & Session Validation.\n"
             "2. Gated Query Normalization & Facet Planning.\n"
             "3. Hybrid Retrieval: pgvector HNSW + GIN BM25.\n"
             "4. Conditional Neo4j 2-Hop Expansion.\n"
             "5. Evidence Gating & Deduplication.\n"
             "6. Vertex AI Gemini 2.5 Flash Generation.\n"
             "7. Citation Post-Validation & Persistence.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "09-chat-request-sequence.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 16: Private Attachment Workflow & Deployment Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Attachment Processing", "Private Attachment Workflow & GCP Production Deployment", 16)
    add_card(s, "Confidentiality Architecture",
             "• User-scoped temporary tokens.\n"
             "• AES-GCM application encryption.\n"
             "• Segregated from public RAG indexes.\n"
             "• Cloud Run ingress with auto-scaling.\n"
             "• Serverless containers (8 vCPU, 16GB RAM) backed by Cloud SQL PostgreSQL.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "15-attachment-question-workflow.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 17: Application Component Responsibilities Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Component Design", "Application Component Responsibilities & Interface Map", 17)
    add_card(s, "Component Roles",
             "• API Gateway: Routing, auth, rate limiting.\n"
             "• Retrieval Service: Hybrid search & RRF.\n"
             "• Graph Service: Neo4j Cypher traversals.\n"
             "• Contract Service: Draft, review, compare.\n"
             "• Worker Service: Async Celery background tasks.\n"
             "• Strict interface boundaries ensure decoupled testability.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "03-application-components.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 18: Cloud SQL Complete Schema Diagram (Large Display)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Physical Database", "Cloud SQL PostgreSQL Complete Physical Schema (23 Tables + 1 MV)", 18)
    add_stat_box(s, "23", "BASE TABLES", 0.6, 1.35, 2.2, 0.9, DARK_EMERALD, 22, 10)
    add_stat_box(s, "1", "MATERIALIZED VIEW", 3.05, 1.35, 2.2, 0.9, TEAL, 22, 10)
    add_stat_box(s, "18", "ALEMBIC MIGRATIONS", 5.5, 1.35, 2.2, 0.9, PURPLE, 22, 10)
    add_stat_box(s, "1024-D", "PGVECTOR COLUMN", 7.95, 1.35, 2.3, 0.9, DARK_EMERALD, 22, 10)
    add_stat_box(s, "0018", "ALEMBIC HEAD", 10.5, 1.35, 2.2, 0.9, AMBER, 22, 10)
    add_image_safe(s, PNG_DIR / "23-cloud-sql-complete-schema.png", 0.6, 2.35, 12.1, 4.55)

    # Slide 19: Database ERD: Identity, Chat & HITL Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Database ERD", "ERD: Identity, Conversation, Chat History & HITL Feedback (7 Tables)", 19)
    add_card(s, "Identity & Chat Model",
             "• app_user: Core UUID identity linked 1-N to sso_identity (Google OAuth).\n"
             "• conversation: Owned by app_user; stores title and timestamps.\n"
             "• chat_message: Stores role, message text, and citation metadata.\n"
             "• chat_answer_feedback: Captures GOOD/BAD ratings and reason codes.\n"
             "• Referential integrity with cascade delete on user conversations.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "04-postgres-erd-identity-chat.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 20: Database ERD: Content, Catalog & GraphRAG Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Database ERD", "ERD: Legal Catalog, GraphRAG Chunks & Runtime Artifacts (10 Tables)", 20)
    add_card(s, "Storage & Serving Model",
             "• graphrag_chunk: Stores 32,334 chunks, vector(1024), tsvector, doc_id, node_id.\n"
             "• graphrag_embedding_checkpoint: Tracks batch embedding progress for resumption.\n"
             "• graphrag_law_version: Version validity intervals and replacement lineage.\n"
             "• legal_answer_cache: Semantic question caching for cost and latency reduction.\n"
             "• 100% parity between SQL and Neo4j via stable natural IDs.",
             0.6, 1.4, 4.8, 5.45, PURPLE, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "05-postgres-erd-content-runtime.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 21: GraphRAG Physical Storage Model Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Storage Architecture", "GraphRAG Multi-Store Storage Model & Synchronization Contract", 21)
    add_card(s, "Multi-Store Projections",
             "• Local SQLite: Single-file offline build artifact + FTS5.\n"
             "• Cloud SQL: Serving store with pgvector HNSW (M=16, ef=64) + GIN.\n"
             "• Neo4j: Graph topology store for multi-hop relation traversal.\n"
             "• Atomic Swap: Only promoted when chunk counts and SHA-256 hashes agree.\n"
             "• Zero embedding drift across model revisions.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "06-graphrag-storage-model.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 22: Physical Database Guarantees & Constraints
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Database Guarantees", "Technical Guarantees & Integrity Constraints of the Database", 22)
    add_card(s, "1. Natural Stable Identifiers",
             "• Uses deterministic slug-based IDs (dieu:bo-luat-45-2019-qh14:91) instead of auto-incrementing integers.\n"
             "• Enables seamless cross-database joins between SQLite, PostgreSQL, Neo4j, and UI citations.",
             0.6, 1.4, 5.9, 2.55, DARK_EMERALD, WHITE, 14, 11.5)
    add_card(s, "2. Indexing Performance Guarantees",
             "• pgvector HNSW: Cosine distance indexing (vector_cosine_ops) for sub-20ms approximate nearest neighbor search.\n"
             "• GIN Indexing: Lexical indexing on generated tsvector combining title, citation, and body text.",
             6.8, 1.4, 5.9, 2.55, TEAL, WHITE, 14, 11.5)
    add_card(s, "3. Zero Schema Drift (Alembic Migrations)",
             "• 18 versioned migration scripts tracking every table, index, and constraint evolution.\n"
             "• Automated migration validation in CI/CD pipeline before production container promotion.",
             0.6, 4.15, 12.1, 2.7, PURPLE, SOFT_GREEN, 13.5, 11)

    # Slide 23: Neo4j Knowledge Graph & Entity Relationships Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Neo4j Topology", "Neo4j Knowledge Graph Structure & Typed Edge Topology", 23)
    add_card(s, "Graph Statistics",
             "• 29,575 Nodes: 13.7k Clauses, 9.8k Points, 4.2k Articles, 434 Chapters.\n"
             "• 108,368 Edges: 28.7k THUỘC_VỀ, 19.2k QUY_ĐỊNH_TẠI, 9.9k DẪN_CHIẾU_ĐẾN.\n"
             "• 41 Node Types & 43 Relation Types covering legal semantics.\n"
             "• Unique constraint on LegalChunk(chunk_id).\n"
             "• Bounded 2-hop traversal follows statutory dependencies.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "07-neo4j-knowledge-graph.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 24: 10-Layer Legal Knowledge Graph Breakdown (L0-L4)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Ontology Breakdown", "Detailed Legal Ontology Specification: Layers L0 to L4", 24)
    add_card(s, "L0: Provenance & Validity",
             "Legal instrument metadata, issuing bodies (National Assembly, Government, Ministries), effective/expiry dates, supersession and amendment relations.",
             0.6, 1.4, 12.1, 0.98, DARK_EMERALD, WHITE, 12.5, 10.5)
    add_card(s, "L1: Document Structure",
             "Full statutory hierarchy: Document > Chapter > Section > Article > Clause > Point. Encodes HAS_PART and cross-article reference edges.",
             0.6, 2.5, 12.1, 0.98, TEAL, WHITE, 12.5, 10.5)
    add_card(s, "L2: Terminology & Topics",
             "Extracted legal definitions ('Employee is...', 'Labor contract is...'), synonyms, and topic hubs connecting user colloquialisms to formal statutes.",
             0.6, 3.6, 12.1, 0.98, PURPLE, WHITE, 12.5, 10.5)
    add_card(s, "L3: Wages & Rewards",
             "Income components, 4 regional minimum wage brackets, overtime pay rates (150%, 200%, 300%), wage deduction caps, and statutory calculation formulas.",
             0.6, 4.7, 12.1, 0.98, AMBER, WHITE, 12.5, 10.5)
    add_card(s, "L4: Actors & Labor Relations",
             "Legal entities: Employees, Employers, Trade Unions, Labor Inspectors; resolves who holds rights and who bears obligations.",
             0.6, 5.8, 12.1, 0.98, ROSE, WHITE, 12.5, 10.5)

    # Slide 25: 10-Layer Legal Knowledge Graph Breakdown (L5-L9)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Ontology Breakdown", "Detailed Legal Ontology Specification: Layers L5 to L9", 25)
    add_card(s, "L5: Procedures & Documentation",
             "Administrative procedures, required dossiers, notification timelines, competent dispute resolution bodies, and procedural prerequisites.",
             0.6, 1.4, 12.1, 0.98, DARK_EMERALD, WHITE, 12.5, 10.5)
    add_card(s, "L6: Time & Limitation Periods",
             "Statutory notice periods (45/30/3 days for contract termination), probation duration caps (180/60/30/6 days), and statute of limitations for labor claims.",
             0.6, 2.5, 12.1, 0.98, TEAL, WHITE, 12.5, 10.5)
    add_card(s, "L7: Sanctions & Compliance Risks",
             "Unlawful conduct, administrative fine brackets (Decree 12/2022/ND-CP), remedial orders, and adverse legal conditions.",
             0.6, 3.6, 12.1, 0.98, ROSE, WHITE, 12.5, 10.5)
    add_card(s, "L8: Contract Lifecycle",
             "Sequential stages: Formation -> Performance -> Modification -> Suspension -> Unilateral/Mutual Termination -> Post-termination obligations.",
             0.6, 4.7, 12.1, 0.98, PURPLE, WHITE, 12.5, 10.5)
    add_card(s, "L9: Case Law & Judicial Precedent",
             "Court judgments, legal holdings, and precedent interpretations (ready for case-law expansion).",
             0.6, 5.8, 12.1, 0.98, AMBER, WHITE, 12.5, 10.5)

    # Slide 26: Adaptive Retrieval & Routing Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Adaptive Routing", "Adaptive Retrieval Routing: Single-Hop vs Multi-Hop Execution", 26)
    add_card(s, "Complexity-Aware Routing",
             "• Single-Hop Path: Fast Hybrid Search (pgvector + GIN) in ~15.8ms.\n"
             "• Multi-Hop Path: Hybrid seed retrieval + Neo4j 2-hop graph expansion in ~146.7ms.\n"
             "• Facet Planning: Splits multi-issue questions into independent sub-queries.\n"
             "• Evidence Gate: Eliminates irrelevant chunks before prompt construction.\n"
             "• Saves 90%+ latency on routine lookups.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "10-adaptive-retrieval-routing.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 27: Measured Data Landscape Diagram (Large Display)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Corpus Metrics", "Measured Legal Data Landscape & Chunk Distribution Baseline", 27)
    add_stat_box(s, "74", "DOCUMENTS", 0.6, 1.35, 2.85, 0.9, DARK_EMERALD, 22, 10)
    add_stat_box(s, "32,334", "CHUNKS", 3.68, 1.35, 2.85, 0.9, TEAL, 22, 10)
    add_stat_box(s, "3.31 M", "TOTAL TOKENS", 6.76, 1.35, 2.85, 0.9, PURPLE, 22, 10)
    add_stat_box(s, "102.35", "AVG TOKENS / CHUNK", 9.84, 1.35, 2.85, 0.9, AMBER, 22, 10)
    add_image_safe(s, PNG_DIR / "21-measured-data-landscape.png", 0.6, 2.35, 12.1, 4.55)

    # Slide 28: Document & Chunk Distribution Analysis
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Chunk Taxonomy", "Breakdown of the 9 Specialized Chunk Types in VlegalAI", 28)
    chunk_rows = [
        ("Clause Chunks", "13,700 (42.37%)", "Primary retrieval unit containing operational rules, rights, obligations, and statutory conditions."),
        ("Point Chunks", "9,762 (30.19%)", "Fine-grained enumerated items, specific compliance duties, and itemized penalty triggers."),
        ("Article Chunks", "4,192 (12.96%)", "Complete or first-window article text for article-level overview and primary citation anchors."),
        ("Table Chunks", "1,248 (3.86%)", "Preserved table rows (regional wage scales, social insurance rates, compensation schedules)."),
        ("Structural / Sliding", "3,432 (10.62%)", "Chapter/Section navigation chunks and sliding-window tails for long articles exceeding 440 words.")
    ]
    y_pos = 1.4
    for title, count, desc in chunk_rows:
        add_card(s, f"{title} - {count}", desc, 0.6, y_pos, 12.1, 0.98, DARK_EMERALD, WHITE, 12.5, 10.5)
        y_pos += 1.08

    # Slide 29: [KEY QUESTION 5] Deterministic Document Parsing & Envelope Schema
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 5: PARSING", "Deterministic Document Parsing & JSON Envelope Schema 1.0", 29)
    add_card(s, "1. Why Deterministic Parsing (No LLM)?",
             "• Parsing statutes via LLM suffers from non-deterministic omissions, hallucinated numbers, and token costs.\n"
             "• VlegalAI implements a deterministic State Machine: σ_i = (d, c_i, s_i, a_i, q_i, p_i) walking ordered DOCX blocks.\n"
             "• Recognizes Chapter (Chương), Section (Mục), Article (Điều), Clause (1.), Point (a)), and nested Tables.",
             0.6, 1.4, 5.9, 2.6, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. Header Regex Specification",
             "• Chapter: ^Chương\\s+([IVXLCDM]+|\\d+)\n"
             "• Section: ^Mục\\s+([IVXLCDM]+|\\d+)\n"
             "• Article: ^Điều\\s+(\\d+[a-zA-Z]?)\\s*[\\.:]\\s*(.+)\n"
             "• Clause: ^(\\d{1,3})\\.\\s+(.+)\n"
             "• Point: ^([a-zđ](\\d+)?)\\)\\s+(.+)\n"
             "• Tables: Materialized beneath governing clause; table rows preserved as formatted lines.",
             6.8, 1.4, 5.9, 2.6, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. Output Document Envelope Schema (JSON 1.0 Contract)",
             "• source: { path, filename, size_bytes, sha256 } -> Immutable identity for build verification.\n"
             "• document: { doc_id, filename, title, code (45/2019/QH14), doc_type, issuer, text } -> Fulltext record.\n"
             "• nodes[]: Typed hierarchy records { node_id, label, parent_id, path_label, text, ordinal, child_count }.\n"
             "• edges[]: Directed relations { edge_id, source_id, target_id, relation, evidence }.\n"
             "• chunks[]: Bounded retrieval units { chunk_id, doc_id, node_id, chunk_type, title, citation, text, token_count }.",
             0.6, 4.15, 12.1, 2.7, PURPLE, SOFT_BLUE, 13.5, 11)

    # Slide 30: [KEY QUESTION 1] Hierarchy-Aware Chunking & String Splitting
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 1: CHUNKING", "Hierarchy-Aware Chunking & String Splitting Algorithm", 30)
    add_stat_box(s, "W = 360", "WINDOW WORDS", 0.6, 1.4, 2.85, 0.95, DARK_EMERALD, 24, 10.5)
    add_stat_box(s, "O = 70", "OVERLAP WORDS", 3.68, 1.4, 2.85, 0.95, TEAL, 24, 10.5)
    add_stat_box(s, "Δ = 290", "STRIDE STEP", 6.76, 1.4, 2.85, 0.95, PURPLE, 24, 10.5)
    add_stat_box(s, "N ≤ 440", "KEEP-WHOLE THRESHOLD", 9.84, 1.4, 2.85, 0.95, AMBER, 24, 10.5)

    add_card(s, "1. Statutory Boundary Segmentation",
             "• Primary segmentation is legal structure (Article, Clause, Point, Table) rather than fixed characters.\n"
             "• Each chunk carries explicit node_id, parent_id, and citation path: 'Labor Code 2019 > Chapter VI > Article 91 > Clause 1'.",
             0.6, 2.5, 5.9, 2.15, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. Sliding Window Fallback for Long Nodes",
             "• Vietnamese word-regex token counter: T = [0-9A-Za-zÀ-ỹĐđ]+.\n"
             "• If N(x) <= 440 words: Retained as single chunk.\n"
             "• If N(x) > 440 words: Sliding windows C_j = x[j*Δ : j*Δ + W] with step Δ = 290 words. Discard tail < 80 words.",
             6.8, 2.5, 5.9, 2.15, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. Text Contract for Embedding",
             "Embedding text format: X_c = title(c) + '\\n' + path_label(c) + '\\n' + chunk.text\n"
             "-> Ensures short clauses retain complete statutory instrument and chapter context during dense vector matching!",
             0.6, 4.8, 12.1, 2.1, PURPLE, SOFT_GREEN, 13.5, 11)

    # Slide 31: [KEY QUESTION 4] Data Processing, Vector Embedding & Caching
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 4: EMBEDDING", "Data Processing, Vector Embedding (1024D) & SHA-256 Caching", 31)
    add_card(s, "1. Embedding Generation & L2 Normalization",
             "• Model: Vertex AI gemini-embedding-001 (1,024 dimensions, float32).\n"
             "• Separate task types: RETRIEVAL_DOCUMENT for chunks, RETRIEVAL_QUERY for user queries.\n"
             "• Client-side L2 Normalization: e(x) = z(x) / ||z(x)||_2.\n"
             "• Cosine distance in pgvector simplifies to dot product: sim(q, c) = e(q)^T * e(c).",
             0.6, 1.4, 5.9, 2.6, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. Storage Payload & Batch Checkpointing",
             "• 1,024 float32 values = 4,096 bytes per vector.\n"
             "• 32,334 chunks ≈ 126.3 MiB raw vector payload.\n"
             "• Checkpointing in batches of 640 pending chunks into PostgreSQL: Allows resuming interrupted indexing jobs without re-embedding.",
             6.8, 1.4, 5.9, 2.6, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. Content Hash Invalidation (SHA-256 Caching Contract)",
             "• Cached vector identity: h_c = SHA256(X_c).\n"
             "• Vectors are reused ONLY when chunk_id, h_c, provider, model, revision, task_type, and dimension agree.\n"
             "• Any modification to legal text changes h_c and invalidates the cached vector automatically, preventing stale index reuse.",
             0.6, 4.15, 12.1, 2.7, PURPLE, SOFT_BLUE, 13.5, 11)

    # Slide 32: [KEY QUESTION 1] Multi-Store Indexing & Index Synchronization Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 1: INDEXING", "Multi-Store Indexing & Stage-and-Activate Synchronization Contract", 32)
    add_card(s, "Index Architecture",
             "• pgvector: HNSW cosine index (M=16, ef_construction=64).\n"
             "• PostgreSQL GIN: Lexical search on generated tsvector.\n"
             "• Neo4j: Graph index on LegalChunk(chunk_id).\n"
             "• Weighted RRF: S_RRF = 0.55 * Vector + 0.45 * BM25 (K=60).\n"
             "• Atomic Activation: Promoted only when chunk counts, dimensions, and SHA-256 fingerprints align 100%.",
             0.6, 1.4, 4.8, 5.45, DARK_EMERALD, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "12-index-synchronization.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 33: [KEY QUESTION 3] Model Strategy: Why VlegalAI Does NOT Fine-Tune LLM
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 3: MODEL", "Model Strategy: Why VlegalAI Does NOT Fine-Tune the Generation LLM", 33)
    box_no = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(3.2), Inches(5.45))
    box_no.fill.solid()
    box_no.fill.fore_color.rgb = DARK_BG
    box_no.line.fill.background()
    tf_no = box_no.text_frame
    tf_no.word_wrap = True
    p_no0 = tf_no.paragraphs[0]
    p_no0.text = "NO"
    p_no0.alignment = PP_ALIGN.CENTER
    p_no0.font.size = Pt(44)
    p_no0.font.bold = True
    p_no0.font.color.rgb = ROSE
    p_no1 = tf_no.add_paragraph()
    p_no1.text = "LLM Weight Fine-Tuning\n(Model weights fixed: θ' = θ)"
    p_no1.alignment = PP_ALIGN.CENTER
    p_no1.font.size = Pt(13)
    p_no1.font.bold = True
    p_no1.font.color.rgb = WHITE
    p_no2 = tf_no.add_paragraph()
    p_no2.text = "\n• No Pretraining\n• No SFT\n• No LoRA / QLoRA\n• Hosted Gemini 2.5 Flash via Vertex AI"
    p_no2.font.size = Pt(11.5)
    p_no2.font.color.rgb = TEXT_LIGHT

    add_card(s, "1. Factual Provenance & Citation Entailment",
             "Fine-tuned LLMs store knowledge in opaque parametric weights (θ), losing exact provision provenance. RAG guarantees transparent evidence trails: Answer -> Evidence -> Chunk -> Provision -> Instrument.",
             4.0, 1.4, 4.3, 2.6, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. Legal Freshness & Zero-Retraining Updates",
             "When legislation changes, fine-tuning requires costly dataset collection and retraining. With GraphRAG, updating laws only requires re-indexing the modified document in minutes.",
             8.5, 1.4, 4.2, 2.6, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. Hallucination Elimination",
             "In legal domains, hallucinating a single exception rule creates severe liability. Enforcing an Evidence-First Prompt Contract constrains generation strictly to retrieved context.",
             4.0, 4.15, 4.3, 2.7, PURPLE, WHITE, 13.5, 11)
    add_card(s, "4. Focused Engineering Optimization",
             "Engineering effort was directed where it matters: 10-layer ontology, hierarchy-aware chunking, hybrid RRF, and deterministic citation validation.",
             8.5, 4.15, 4.2, 2.7, AMBER, WHITE, 13.5, 11)

    # Slide 34: Experimental Architecture Definitions
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Experimental Design", "Four Compared Retrieval-Augmented Generation Architectures", 34)
    add_card(s, "1. Dense RAG (Vector Baseline)",
             "• Dense semantic retrieval using gemini-embedding-001 (1024D) on pgvector.\n"
             "• Retrieves Top-k = 4 chunks by cosine similarity; feeds directly to Gemini 2.5 Flash.\n"
             "• Baseline for dense semantic representation speed and accuracy.",
             0.6, 1.4, 5.9, 2.6, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. LightRAG (Hybrid + 1-Hop Graph)",
             "• Combines Dense Vector + SQLite FTS5 BM25 via Reciprocal Rank Fusion.\n"
             "• Expands 1-hop graph neighbors around top retrieved entities.\n"
             "• Lightweight hybrid baseline balancing lexical matching and local relations.",
             6.8, 1.4, 5.9, 2.6, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. GraphRAG (Hybrid Seed + 2-Hop Graph)",
             "• Seeds candidates via Hybrid Retrieval, then performs 2-hop Neo4j graph expansion.\n"
             "• Traverses DẪN_CHIẾU_ĐẾN, QUY_ĐỊNH_TẠI, and CÓ_KHOẢN relations.\n"
             "• Deep relational reasoning architecture for complex statutory dependencies.",
             0.6, 4.15, 5.9, 2.7, PURPLE, WHITE, 13.5, 11)
    add_card(s, "4. RAG + GraphRAG (Combined Pipeline)",
             "• Combines direct RAG seed evidence with graph-derived candidate paths.\n"
             "• Evaluated across 100 requests with complete latency and ID recall telemetry.",
             6.8, 4.15, 5.9, 2.7, AMBER, WHITE, 13.5, 11)

    # Slide 35: [KEY QUESTION 2] Evaluation Methodology: Accuracy as an 8-D Vector
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 2: EVALUATION", "Accuracy Methodology: 8-Dimensional Quality Vector & RAGAS Benchmark", 35)
    add_card(s, "Accuracy in Legal QA is an 8-Dimensional Vector",
             "Q = (Q_ret, Q_faith, Q_rel, Q_fact, Q_cite, Q_facet, Q_temp, Q_lat)\n"
             "Exact-string matching is inadequate for legal QA: fluent text can cite superseded laws, while two different wordings can state the exact same statutory rule.",
             0.6, 1.4, 12.1, 1.15, DARK_EMERALD, SOFT_GREEN, 13.5, 11)

    add_card(s, "1. FAITHFULNESS",
             "Proportion of atomic claims in the answer supported by retrieved context (|S_q| / |A_q|). Measures hallucination elimination.",
             0.6, 2.75, 2.85, 2.2, DARK_EMERALD, WHITE, 13, 10.5)
    add_card(s, "2. ANSWER RELEVANCY",
             "Degree to which the generated answer directly addresses the user's inquiry without extraneous digressions.",
             3.68, 2.75, 2.85, 2.2, TEAL, WHITE, 13, 10.5)
    add_card(s, "3. CONTEXT PRECISION & RECALL",
             "Precision and Recall of retrieved chunks compared to ground-truth reference evidence.",
             6.76, 2.75, 2.85, 2.2, PURPLE, WHITE, 13, 10.5)
    add_card(s, "4. FACTUAL CORRECTNESS",
             "Semantic and factual agreement between the answer and expert ground-truth answers.",
             9.84, 2.75, 2.85, 2.2, AMBER, WHITE, 13, 10.5)

    add_card(s, "Context ID Precision & Recall",
             "Measures exact statutory identifier overlap: R_ID = |R_k ∩ G| / |G| and P_ID = |R_k ∩ G| / |R_k|.",
             0.6, 5.1, 5.9, 1.75, DARK_EMERALD, WHITE, 13, 10.5)
    add_card(s, "100-Question RAGAS Benchmark Dataset",
             "Stratified into 50 Single-hop, 25 Multi-hop Specific, and 25 Multi-hop Abstract questions evaluated across all 4 architectures.",
             6.8, 5.1, 5.9, 1.75, TEAL, WHITE, 13, 10.5)

    # Slide 36: [KEY QUESTION 2] Aggregate Quality Results & Evaluation Chart
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 2: RESULTS", "Aggregate Quality Results across 100 Legal Questions (RAGAS Benchmark)", 36)
    table_shape = s.shapes.add_table(5, 9, Inches(0.6), Inches(1.4), Inches(7.5), Inches(2.3))
    table = table_shape.table
    table.columns[0].width = Inches(1.5)
    for i in range(1, 9):
        table.columns[i].width = Inches(0.75)

    headers = ["Arch", "Faith", "Relev", "CtxPrec", "CtxRec", "Fact", "IDRec", "Lat(ms)", "Overall"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE

    rows_data = [
        ["Dense RAG", "0.9111", "0.9126", "0.7783", "0.8154", "0.3985", "0.6800", "15.8", "0.7619*"],
        ["LightRAG", "0.9160", "0.8951", "0.6950", "0.8124", "0.4271", "0.6650", "148.8", "0.7481"],
        ["GraphRAG", "0.9213*", "0.8641", "0.5558", "0.7820", "0.4003", "0.5750", "146.7", "0.7041"],
        ["RAG+Graph", "0.9403#", "--", "--", "--", "--", "0.6850*", "162.6", "Checkpoint"]
    ]
    for row_idx, r in enumerate(rows_data):
        for col_idx, val in enumerate(r):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT_GREEN if row_idx == 0 else (SOFT_BLUE if row_idx == 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.bold = (col_idx == 0 or "*" in val)
            p.font.color.rgb = DARK_EMERALD if "*" in val else TEXT_DARK

    add_image_safe(s, EVAL_PNG, 8.3, 1.4, 4.4, 5.45)

    add_card(s, "Key Empirical Findings",
             "• GraphRAG achieves the highest Faithfulness (0.9213), eliminating hallucinations via relational context.\n"
             "• Dense RAG achieves the highest Context Precision (0.7783) and fastest retrieval (15.8ms vs 146.7ms).\n"
             "• Validates Adaptive Routing: Fast Hybrid RAG for single-hop, GraphRAG for complex queries.",
             0.6, 3.85, 7.5, 3.0, DARK_EMERALD, SOFT_GREEN, 13, 11)

    # Slide 37: [KEY QUESTION 2] Performance Stratification by Question Complexity
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "KEY FOCUS 2: COMPLEXITY", "Performance Stratification by Question Complexity (Single-hop vs Multi-hop)", 37)
    add_card(s, "1. Single-Hop Specific (50 Questions)",
             "RAGAS Overall: Dense 0.8114 | Light 0.8039 | Graph 0.7458\n\n"
             "• Direct lookup of specific provisions, wage rates, or notice periods.\n"
             "• Dense RAG dominates speed and precision without multi-hop noise.",
             0.6, 1.4, 3.85, 3.7, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. Multi-Hop Specific (25 Questions)",
             "RAGAS Overall: Dense 0.7625 | Light 0.7283 | Graph 0.6914\n\n"
             "• Requires combining 2-3 provisions (e.g. eligibility condition + employer obligations).\n"
             "• Knowledge graph captures statutory cross-reference pathways.",
             4.75, 1.4, 3.85, 3.7, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. Multi-Hop Abstract (25 Questions)",
             "RAGAS Overall: Dense 0.6624 | Light 0.6563 | Graph 0.6334\n\n"
             "• Synthesis and comparative scenarios across employee categories.\n"
             "• GraphRAG shows smallest degradation (-0.112 vs -0.149 for Dense).",
             8.9, 1.4, 3.85, 3.7, PURPLE, WHITE, 13.5, 11)
    add_card(s, "Trade-Off Conclusion",
             "Graph reasoning adds ~130ms retrieval latency; pairing Dense seeds with bounded 2-hop graph expansion yields the optimal production balance.",
             0.6, 5.25, 12.1, 1.6, AMBER, SOFT_AMBER, 12.5, 10.5)

    # Slide 38: Request-Level Latency Baseline
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Latency Evaluation", "Empirical Latency Profile: Retrieval vs Generation Overhead", 38)
    add_stat_box(s, "15.8 ms", "DENSE RETRIEVAL (MEAN)", 0.6, 1.4, 2.85, 0.95, DARK_EMERALD, 24, 10.5)
    add_stat_box(s, "146.7 ms", "GRAPH RETRIEVAL (MEAN)", 3.68, 1.4, 2.85, 0.95, TEAL, 24, 10.5)
    add_stat_box(s, "5.72 s", "DENSE E2E (MEAN)", 6.76, 1.4, 2.85, 0.95, PURPLE, 24, 10.5)
    add_stat_box(s, "6.13 s", "GRAPH E2E (MEAN)", 9.84, 1.4, 2.85, 0.95, AMBER, 24, 10.5)

    add_card(s, "Latency Breakdown & Analysis",
             "• Retrieval Phase: Dense RAG (15.8ms) is ~9.3x faster than GraphRAG (146.7ms) due to avoiding remote Neo4j graph traversals.\n"
             "• Generation Phase: Dominates 95-98% of total End-to-End latency (5.7s - 8.3s), driven by Vertex AI cloud token generation.\n"
             "• P95 Tail Latency: Reaches 25.4s - 48.7s on long multi-part legal scenarios.\n"
             "• Architectural Takeaway: Optimizing graph queries alone cannot eliminate the p95 tail; prompt budgeting and streaming are essential.",
             0.6, 2.5, 12.1, 4.35, DARK_EMERALD, WHITE, 13.5, 11.5)

    # Slide 39: Paired Bootstrap 95% Confidence Intervals
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Statistical Analysis", "Paired Bootstrap 95% Confidence Intervals on Overall RAGAS", 39)
    add_card(s, "1. Dense RAG vs LightRAG",
             "• Mean Difference: +0.0138\n• 95% CI: [-0.0144, +0.0423]\n• Conclusion: Interval spans zero -> Difference is not statistically significant.",
             0.6, 1.4, 3.85, 3.1, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "2. Dense RAG vs GraphRAG",
             "• Mean Difference: +0.0578\n• 95% CI: [+0.0270, +0.0916]\n• Conclusion: Strictly positive -> Dense RAG is statistically significantly higher overall.",
             4.75, 1.4, 3.85, 3.1, TEAL, WHITE, 13.5, 11)
    add_card(s, "3. LightRAG vs GraphRAG",
             "• Mean Difference: +0.0440\n• 95% CI: [+0.0092, +0.0816]\n• Conclusion: Strictly positive -> LightRAG is statistically significantly higher overall.",
             8.9, 1.4, 3.85, 3.1, PURPLE, WHITE, 13.5, 11)
    add_card(s, "Methodological Rigor",
             "Paired bootstrap resamples question indices (not independent rows), preserving query pairing and avoiding over-optimistic significance claims.",
             0.6, 4.7, 12.1, 2.15, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    # Slide 40: Checkpoint Completeness & Missing Value Handling
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Evaluation Integrity", "Dataset Integrity & Transparent Handling of Checkpoint Values", 40)
    add_card(s, "Honest Reporting Principles (No Imputed or Synthetic Data)",
             "• 3 Complete Architectures (Dense, Light, Graph): 100/100 complete evaluations across all dimensions.\n"
             "• RAG+GraphRAG Checkpoint: Contains 100 valid latency and ID recall/precision rows, but only 45 valid Faithfulness rows due to evaluator timeout.\n"
             "• Strict Protocol: Missing quality fields are NOT imputed with zeros or synthetic approximations. The partial branch is excluded from overall ranking.",
             0.6, 1.4, 12.1, 5.45, DARK_EMERALD, SOFT_GREEN, 14, 12)

    # Slide 41: Benchmark Error Analysis & Failure Modes Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Error Analysis", "Empirical Error Analysis, Failure Taxonomy & Recovery Controls", 41)
    add_card(s, "Observed Error Patterns",
             "• Low Identifier Precision (0.21-0.24): Retrieved context includes helpful background provisions beyond minimal gold set.\n"
             "• Graph Context Precision Drop (0.5558): 2-hop traversal introduces distant neighbors on simple single-hop queries.\n"
             "• Multi-Abstract Degradation: Abstract questions require cross-decree synthesis, increasing context dilution risk.\n"
             "• Engineered Quality Gates: Evidence gates and retry budgets prevent provider timeouts from cascading.",
             0.6, 1.4, 4.8, 5.45, ROSE, WHITE, 14, 11.5)
    add_image_safe(s, PNG_DIR / "20-observability-recovery.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 42: CI/CD & Automated System Verification Diagram (Large Display)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "System Verification", "CI/CD Pipeline, Automated Test Pyramid & Verification Evidence", 42)
    add_stat_box(s, "492 / 492", "BACKEND TESTS PASS", 0.6, 1.35, 2.85, 0.9, DARK_EMERALD, 22, 10)
    add_stat_box(s, "16 / 16", "FRONTEND TESTS PASS", 3.68, 1.35, 2.85, 0.9, TEAL, 22, 10)
    add_stat_box(s, "18", "ALEMBIC MIGRATIONS", 6.76, 1.35, 2.85, 0.9, PURPLE, 22, 10)
    add_stat_box(s, "100%", "IMMUTABLE CI/CD BUILD", 9.84, 1.35, 2.85, 0.9, AMBER, 22, 10)
    add_image_safe(s, PNG_DIR / "18-cicd-release-workflow.png", 0.6, 2.35, 12.1, 4.55)

    # Slide 43: Answers to Research Questions (RQ1 - RQ4)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Research Conclusions", "Summary of Empirical Answers to Research Questions (RQ1 - RQ4)", 43)
    add_card(s, "RQ1 (Answer Quality)",
             "GraphRAG achieves highest Faithfulness (0.9213) for hallucination suppression. Dense RAG achieves highest Overall score (0.7619) due to superior Context Precision on single-hop queries.",
             0.6, 1.4, 5.9, 2.6, DARK_EMERALD, WHITE, 13.5, 11)
    add_card(s, "RQ2 (Latency Trade-off)",
             "Dense retrieval is ~9.3x faster than graph expansion (15.8ms vs 146.7ms). Confirms that graph traversal should be selectively routed rather than universally executed.",
             6.8, 1.4, 5.9, 2.6, TEAL, WHITE, 13.5, 11)
    add_card(s, "RQ3 (Reasoning Depth)",
             "GraphRAG demonstrates smallest relative drop (-0.112) when scaling from single-hop to multi-abstract queries, proving its strength in synthesizing cross-statutory provisions.",
             0.6, 4.15, 5.9, 2.7, PURPLE, WHITE, 13.5, 11)
    add_card(s, "RQ4 (Reproducibility)",
             "Achieved 100% software reproducibility: 1 Git commit SHA, 492 passed tests, 1 Alembic head, immutable Docker image, and open benchmark logs.",
             6.8, 4.15, 5.9, 2.7, AMBER, WHITE, 13.5, 11)

    # Slide 44: Threats to Validity, Limitations & Roadmap Diagram
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Limitations & Roadmap", "Threats to Validity, Current Limitations & Future Roadmap", 44)
    add_card(s, "Current Limitations & Roadmap",
             "• P95 tail latency is dominated by cloud LLM token generation.\n"
             "• Evaluation benchmark limited to 100 questions; needs broader multi-expert adjudication.\n"
             "• Layer 9 (Case Law) is structurally designed but awaiting bulk judgment corpus ingestion.\n\n"
             "Future Roadmap:\n"
             "1. Train local cross-encoder reranker to reduce latency.\n"
             "2. Ingest judicial case law into Layer 9.\n"
             "3. Expand corpus to Corporate & Tax Law.",
             0.6, 1.4, 4.8, 5.45, AMBER, WHITE, 13.5, 11)
    add_image_safe(s, PNG_DIR / "13-freshness-reindex-workflow.png", 5.6, 1.4, 7.1, 5.45)

    # Slide 45: Defense Conclusion & Technical Q&A Cheatsheet
    s = prs.slides.add_slide(blank_layout)
    bg45 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg45.fill.solid()
    bg45.fill.fore_color.rgb = DARK_BG
    bg45.line.fill.background()

    tb_end = s.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.3))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True
    p_end0 = tf_end.paragraphs[0]
    p_end0.text = "THANK YOU TO THE DEFENSE COMMITTEE"
    p_end0.alignment = PP_ALIGN.CENTER
    p_end0.font.size = Pt(28)
    p_end0.font.bold = True
    p_end0.font.color.rgb = WHITE
    p_end1 = tf_end.add_paragraph()
    p_end1.text = "VLEGALAI IS READY FOR QUESTIONS & FEEDBACK"
    p_end1.alignment = PP_ALIGN.CENTER
    p_end1.font.size = Pt(15)
    p_end1.font.bold = True
    p_end1.font.color.rgb = TEAL

    qa_list = [
        ("1. Chunking & Splitting", "Hierarchy-Aware by Article/Clause/Point. Sliding window W=360, Overlap=70 words; Vietnamese regex counter."),
        ("2. LLM Training?", "NO fine-tuning. Hosted Gemini 2.5 Flash + RAG/GraphRAG preserves 100% provenance and instant statutory freshness."),
        ("3. Accuracy Evaluation", "8-D quality vector via RAGAS (100 legal questions): GraphRAG highest Faithfulness (0.9213), Dense fastest (15.8ms)."),
        ("4. Embedding Pipeline", "Text contract: Title + Path + Text. gemini-embedding-001 (1024D), L2 normalized, SHA-256 content caching."),
        ("5. Document Parsing", "Deterministic State Machine parsing DOCX into JSON 1.0 Envelope: source, document, nodes, edges, chunks.")
    ]

    x_positions = [0.8, 4.8, 8.8, 2.8, 6.8]
    y_positions = [2.1, 2.1, 2.1, 4.25, 4.25]
    for idx, (q_title, q_desc) in enumerate(qa_list):
        box_qa = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(y_positions[idx]), Inches(3.7), Inches(1.95))
        box_qa.fill.solid()
        box_qa.fill.fore_color.rgb = CARD_DARK
        box_qa.line.color.rgb = EMERALD
        tf_q = box_qa.text_frame
        tf_q.word_wrap = True
        tf_q.margin_left = tf_q.margin_top = tf_q.margin_right = tf_q.margin_bottom = Inches(0.12)
        pq0 = tf_q.paragraphs[0]
        pq0.text = q_title
        pq0.font.size = Pt(13)
        pq0.font.bold = True
        pq0.font.color.rgb = AMBER
        pq1 = tf_q.add_paragraph()
        pq1.text = q_desc
        pq1.font.size = Pt(10.5)
        pq1.font.color.rgb = TEXT_LIGHT

    tb_foot = s.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4))
    pf = tb_foot.text_frame.paragraphs[0]
    pf.text = "VlegalAI: Vietnamese Labor-Law GraphRAG & AI Assistant | FPT University Quy Nhon 2026"
    pf.alignment = PP_ALIGN.CENTER
    pf.font.size = Pt(10.5)
    pf.font.color.rgb = TEXT_MUTED

    print(f"Saving {len(prs.slides)} slides to {OUT_SLIDES}...")
    prs.save(str(OUT_SLIDES))
    print(f"Saving {len(prs.slides)} slides to {OUT_SLDIES}...")
    prs.save(str(OUT_SLDIES))
    print("Successfully built English 45-slide presentation with LARGE fonts and LARGE diagrams!")

if __name__ == "__main__":
    build_presentation()
