"""Script to generate graduation defense presentation PPTX for VlegalAI.

Covers all 5 core technical questions:
1. Chunking & Indexing, string splitting
2. Accuracy & Evaluation methodology
3. LLM Training / Fine-tuning vs RAG
4. Data preprocessing & Vector embedding
5. Document structure when parsing
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Output paths
ROOT = Path(r"f:\VlegalAI")
OUT_SLIDES = ROOT / "slides.pptx"
OUT_SLDIES = ROOT / "sldies.pptx"

# Color Palette
DARK_BG = RGBColor(15, 23, 42)       # #0F172A Slate Dark
CARD_DARK = RGBColor(30, 41, 59)     # #1E293B Slate Navy
TEXT_LIGHT = RGBColor(248, 250, 252) # #F8FAFC
TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A

EMERALD = RGBColor(16, 185, 129)     # #10B981 Vibrant Green
DARK_EMERALD = RGBColor(15, 76, 58)  # #0F4C3A Primary Dark Green
TEAL = RGBColor(14, 165, 233)        # #0EA5E9 Sky Blue
PURPLE = RGBColor(139, 92, 246)      # #8B5CF6 Purple
AMBER = RGBColor(245, 158, 11)       # #F59E0B Amber Gold
ROSE = RGBColor(244, 63, 94)         # #F43F5E Crimson Rose

BG_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC
WHITE = RGBColor(255, 255, 255)
BORDER_LIGHT = RGBColor(226, 232, 240) # #E2E8F0
SOFT_GREEN = RGBColor(236, 253, 245)   # #ECFDF5
SOFT_BLUE = RGBColor(240, 249, 255)    # #F0F9FF
SOFT_AMBER = RGBColor(254, 243, 199)   # #FEF3C7

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, section: str, title: str, page_num: int, total_pages: int = 20):
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.fill.background()

        # Top Section Bar
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.45))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = SOFT_GREEN
        top_bar.line.fill.background()

        tf_top = top_bar.text_frame
        tf_top.word_wrap = True
        p_sec = tf_top.paragraphs[0]
        p_sec.text = f"  VLEGALAI | {section.upper()}"
        p_sec.font.size = Pt(11)
        p_sec.font.bold = True
        p_sec.font.color.rgb = DARK_EMERALD

        # Page Number Top Right
        tb_page = slide.shapes.add_textbox(Inches(11.5), Inches(0.05), Inches(1.5), Inches(0.35))
        p_pg = tb_page.text_frame.paragraphs[0]
        p_pg.text = f"{page_num:02d} / {total_pages:02d}"
        p_pg.alignment = PP_ALIGN.RIGHT
        p_pg.font.size = Pt(11)
        p_pg.font.bold = True
        p_pg.font.color.rgb = TEXT_MUTED

        # Main Title
        tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.133), Inches(0.65))
        p_t = tb_title.text_frame.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        # Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(1.2), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = EMERALD
        line.line.fill.background()

        # Footer Line
        fline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.05), Inches(12.133), Inches(0.02))
        fline.fill.solid()
        fline.fill.fore_color.rgb = BORDER_LIGHT
        fline.line.fill.background()

        # Footer Text
        tb_f = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(8.0), Inches(0.3))
        pf = tb_f.text_frame.paragraphs[0]
        pf.text = "VlegalAI: Vietnamese Labor Law GraphRAG & AI Assistant | Graduation Thesis Defense 2026"
        pf.font.size = Pt(9.5)
        pf.font.color.rgb = TEXT_MUTED

        tb_fr = slide.shapes.add_textbox(Inches(9.5), Inches(7.1), Inches(3.2), Inches(0.3))
        pfr = tb_fr.text_frame.paragraphs[0]
        pfr.text = f"Trang {page_num:02d}"
        pfr.alignment = PP_ALIGN.RIGHT
        pfr.font.size = Pt(9.5)
        pfr.font.color.rgb = TEXT_MUTED

    def add_card(slide, title: str, text: str, left, top, width, height, accent_color=EMERALD, bg_color=WHITE, title_size=13, body_size=10.5):
        # Card Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = BORDER_LIGHT
        box.line.width = Pt(1)

        # Left Accent Stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent_color
        stripe.line.fill.background()

        # Content
        tb = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.3), Inches(height - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(title_size)
        p0.font.bold = True
        p0.font.color.rgb = accent_color

        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.size = Pt(body_size)
        p1.font.color.rgb = TEXT_DARK

    def add_stat_box(slide, value: str, label: str, left, top, width, height=0.95, color=DARK_EMERALD):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER_LIGHT
        box.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.08), Inches(width), Inches(height - 0.16))
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = value
        p0.alignment = PP_ALIGN.CENTER
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = color

        p1 = tf.add_paragraph()
        p1.text = label
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: COVER SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()

    # Left decorative bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EMERALD
    bar.line.fill.background()

    # Badge Pill
    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.7), Inches(4.5), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_DARK
    pill.line.color.rgb = EMERALD
    p_pill = pill.text_frame.paragraphs[0]
    p_pill.text = "HỘI ĐỒNG BẢO VỆ ĐỒ ÁN TỐT NGHIỆP AI - 2026"
    p_pill.alignment = PP_ALIGN.CENTER
    p_pill.font.size = Pt(11)
    p_pill.font.bold = True
    p_pill.font.color.rgb = EMERALD

    # Main Title
    tb_c1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.8))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    p_c1 = tf_c1.paragraphs[0]
    p_c1.text = "VLEGALAI: HỆ THỐNG TRỢ LÝ PHÁP LÝ LAO ĐỘNG"
    p_c1.font.size = Pt(32)
    p_c1.font.bold = True
    p_c1.font.color.rgb = WHITE

    p_c2 = tf_c1.add_paragraph()
    p_c2.text = "Ứng dụng GraphRAG, Đồ thị Tri thức 10 Tầng và Mô hình Ngôn ngữ Lớn"
    p_c2.font.size = Pt(20)
    p_c2.font.color.rgb = TEAL

    # Team & Supervisor Info Cards
    box_sv = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.4), Inches(5.6), Inches(2.3))
    box_sv.fill.solid()
    box_sv.fill.fore_color.rgb = CARD_DARK
    box_sv.line.color.rgb = BORDER_LIGHT
    tf_sv = box_sv.text_frame
    tf_sv.word_wrap = True
    p_sv0 = tf_sv.paragraphs[0]
    p_sv0.text = "SINH VIÊN THỰC HIỆN:"
    p_sv0.font.size = Pt(12)
    p_sv0.font.bold = True
    p_sv0.font.color.rgb = EMERALD
    p_sv1 = tf_sv.add_paragraph()
    p_sv1.text = "• Trần Tuấn Kiệt (Trưởng nhóm) - QE180152\n• Lê Thanh Đạt - QE170186\n• Phan Bảo Khánh - DE170648"
    p_sv1.font.size = Pt(13)
    p_sv1.font.color.rgb = TEXT_LIGHT

    box_gv = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.4), Inches(5.6), Inches(2.3))
    box_gv.fill.solid()
    box_gv.fill.fore_color.rgb = CARD_DARK
    box_gv.line.color.rgb = BORDER_LIGHT
    tf_gv = box_gv.text_frame
    tf_gv.word_wrap = True
    p_gv0 = tf_gv.paragraphs[0]
    p_gv0.text = "GIẢNG VIÊN HƯỚNG DẪN:"
    p_gv0.font.size = Pt(12)
    p_gv0.font.bold = True
    p_gv0.font.color.rgb = AMBER
    p_gv1 = tf_gv.add_paragraph()
    p_gv1.text = "• GVHD: ThS. Lê Trung Hiếu\n• Đồng GVHD: ThS. Trương Ngọc Hùng\n• Mã đề tài: AIP491 - FPT University Quy Nhơn"
    p_gv1.font.size = Pt(13)
    p_gv1.font.color.rgb = TEXT_LIGHT

    # Footer note
    tb_foot = s1.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5))
    p_foot = tb_foot.text_frame.paragraphs[0]
    p_foot.text = "Tháng 08 / 2026 | Chuyên ngành Trí tuệ Nhân tạo (Artificial Intelligence)"
    p_foot.font.size = Pt(11)
    p_foot.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: BỐI CẢNH & TÍNH CẤP THIẾT CỦA ĐỀ TÀI
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Tổng quan Đề tài", "Bối cảnh Thực tiễn & Thách thức trong Tra cứu Pháp luật Lao động", 2)
    add_card(s2, "1. Thách thức Pháp luật Lao động Việt Nam",
             "• Cấu trúc thứ bậc nghiêm ngặt: Văn bản > Chương > Mục > Điều > Khoản > Điểm.\n"
             "• Chứa nhiều điều kiện loại trừ, dẫn chiếu chéo phức tạp giữa các văn bản quy phạm.\n"
             "• Thường xuyên sửa đổi, thay thế hiệu lực (Nghị định mới bãi bỏ hoặc thay thế Nghị định cũ).",
             0.6, 1.55, 5.9, 2.35, ROSE)

    add_card(s2, "2. Hạn chế Cốt tử của LLM Truyền thống",
             "• Ảo giác (Hallucination): Tự bịa số hiệu điều luật, sai mức lương tối thiểu vùng hoặc mức phạt.\n"
             "• Thiếu bằng chứng kiểm chứng: Không thể chứng minh nguồn gốc chính xác của câu trả lời.\n"
             "• Tri thức đóng băng: Không thể cập nhật luật mới nếu không trải qua quá trình train lại tốn kém.",
             6.8, 1.55, 5.9, 2.35, AMBER)

    add_card(s2, "3. Mục tiêu & Giải pháp Đột phá của VlegalAI",
             "Xây dựng hệ thống Trợ lý Pháp lý chuyên sâu cho Luật Lao động Việt Nam kết hợp Đồ thị Tri thức (Legal Knowledge Graph 10 Tầng) và Adaptive Hybrid RAG (pgvector + BM25 + Neo4j):\n"
             "  [✓] 100% câu trả lời có trích dẫn nguồn luật chính xác đến từng Điều, Khoản, Điểm (Grounding & Provenance).\n"
             "  [✓] Phân luồng thích ứng (Adaptive Routing): Tối ưu tốc độ (~15ms cho single-hop) và suy luận đa quan hệ cho multi-hop.\n"
             "  [✓] Cập nhật tri thức động qua External Stores mà không cần can thiệp trọng số LLM.",
             0.6, 4.1, 12.1, 2.7, DARK_EMERALD, SOFT_GREEN, 14, 11)

    # =========================================================================
    # SLIDE 3: KIẾN TRÚC TỔNG THỂ HỆ THỐNG (4 TẦNG)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Kiến trúc Hệ thống", "Kiến trúc 4 Tầng Toàn diện & Phân tách Trách nhiệm Rõ ràng", 3)

    add_card(s3, "1. TẦNG TRẢI NGHIỆM (UI/UX)",
             "• React 18 + Vite Single Page App\n"
             "• Trích dẫn S1-Sn nhấp chuột xem toàn văn\n"
             "• Chat đa phiên, upload hợp đồng DOCX/PDF/OCR\n"
             "• Thu thập phản hồi Human-in-the-Loop (HITL)",
             0.6, 1.55, 2.85, 4.2, DARK_EMERALD)

    add_card(s3, "2. TẦNG DỊCH VỤ (API/GATEWAY)",
             "• FastAPI async runtime + Gunicorn\n"
             "• OAuth 2.0 Google OIDC + PKCE\n"
             "• Phân luồng Adaptive Routing & Facet Planning\n"
             "• Cổng bằng chứng Evidence Gate & Citation Check",
             3.68, 1.55, 2.85, 4.2, TEAL)

    add_card(s3, "3. TẦNG LƯU TRỮ (STORAGE)",
             "• Cloud SQL PostgreSQL 18 (23 tables)\n"
             "• pgvector HNSW (vector 1,024D)\n"
             "• Neo4j Knowledge Graph (30k nodes, 108k edges)\n"
             "• SQLite FTS5 (Lexical build & local cache)",
             6.76, 1.55, 2.85, 4.2, PURPLE)

    add_card(s3, "4. TẦNG AI & KNOWLEDGE",
             "• Vertex AI Gemini 2.5 Flash (Generation)\n"
             "• gemini-embedding-001 (1024D vector)\n"
             "• 10-Layer Legal Knowledge Graph\n"
             "• Weighted Reciprocal Rank Fusion (RRF)",
             9.84, 1.55, 2.85, 4.2, AMBER)

    add_card(s3, "Quy tắc cốt lõi của Kiến trúc",
             "Phân tách tuyệt đối giữa Tầng suy luận (LLM cố định trọng số) và Tầng tri thức pháp lý (Vector + Đồ thị được cập nhật động theo thời gian thực).",
             0.6, 5.95, 12.1, 0.9, DARK_EMERALD, SOFT_GREEN, 11.5, 10)

    # =========================================================================
    # SLIDE 4: CẤU TRÚC DOCUMENT KHI PARSING (DETERMINISTIC PARSING)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Xử lý Dữ liệu", "Cấu trúc Document khi Parsing & Phương pháp Tất định (Deterministic Parsing)", 4)

    add_card(s4, "1. Tại sao Parsing Tất định (Không dùng LLM)?",
             "• Dùng LLM để parse dễ bị mất dòng, ảo giác số điều khoản và không bảo đảm tính nhất quán (reproducibility).\n"
             "• VlegalAI sử dụng State Machine kết hợp Regular Expressions đặc thù pháp lý Việt Nam: Duyệt tuần tự qua luồng khối (ordered block stream).\n"
             "• Trạng thái Parser: σ_i = (d, c_i, s_i, a_i, q_i, p_i) tương ứng Document > Chapter > Section > Article > Clause > Point.",
             0.6, 1.55, 5.9, 2.6, DARK_EMERALD)

    add_card(s4, "2. Quy chuẩn Nhận dạng Header Regex",
             "• Chương: ^Chương\\s+([IVXLCDM]+|\\d+)\n"
             "• Mục: ^Mục\\s+([IVXLCDM]+|\\d+)\n"
             "• Điều: ^Điều\\s+(\\d+[a-zA-Z]?)\\s*[\\.:]\\s*(.+)\n"
             "• Khoản: ^(\\d{1,3})\\.\\s+(.+)\n"
             "• Điểm: ^([a-zđ](\\d+)?)\\)\\s+(.+)\n"
             "• Bảng biểu (Table): Giữ nguyên vị trí ngay dưới điều khoản chứa nó (không bị dồn xuống cuối file).",
             6.8, 1.55, 5.9, 2.6, TEAL)

    add_card(s4, "3. Cấu trúc Document Envelope Output khi Parsing (JSON Schema 1.0)",
             "Mỗi văn bản được chuẩn hóa thành 1 Envelope chứa 4 mảng thực thể độc lập:\n"
             "• source: { path, filename, size_bytes, sha256 } -> Khóa bất biến ràng buộc toàn vẹn dữ liệu.\n"
             "• document: { doc_id, filename, title, code (45/2019/QH14), doc_type, issuer, text } -> Toàn văn và siêu dữ liệu.\n"
             "• nodes[]: Mảng các thực thể pháp lý { node_id, label, parent_id, path_label, text, ordinal, child_count }.\n"
             "• edges[]: Mảng quan hệ có hướng { edge_id, source_id, target_id, relation, evidence }.\n"
             "• chunks[]: Mảng đơn vị truy hồi { chunk_id, doc_id, node_id, chunk_type, title, citation, text, token_count }.",
             0.6, 4.3, 12.1, 2.5, PURPLE, SOFT_BLUE, 13, 10.5)

    # =========================================================================
    # SLIDE 5: MINH CHỨNG PARSING: BỘ LUẬT LAO ĐỘNG 45/2019/QH14
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Minh chứng Parsing", "Minh chứng Thực tế: Bóc tách Bộ Luật Lao Động 45/2019/QH14", 5)

    add_stat_box(s5, "1", "VĂN BẢN GỐC", 0.6, 1.55, 1.8, 0.95, DARK_EMERALD)
    add_stat_box(s5, "1,175", "NODES ĐỒ THỊ", 2.65, 1.55, 1.8, 0.95, TEAL)
    add_stat_box(s5, "2,347", "EDGES QUAN HỆ", 4.7, 1.55, 1.8, 0.95, PURPLE)
    add_stat_box(s5, "1,187", "CHUNKS TRUY HỒI", 6.75, 1.55, 1.9, 0.95, DARK_EMERALD)
    add_stat_box(s5, "3.85 MiB", "JSON SERIALIZED", 8.9, 1.55, 1.9, 0.95, AMBER)
    add_stat_box(s5, "0", "VECTORS (Parser Stage)", 11.05, 1.55, 1.7, 0.95, ROSE)

    add_card(s5, "Chi tiết Phân rã Node & Edge Pháp lý",
             "• 1 Document, 1 Cơ quan ban hành (Quốc hội).\n"
             "• 17 Chương, 24 Mục, 220 Điều, 640 Khoản, 272 Điểm.\n"
             "• 2,347 Edges quan hệ cấp bậc và dẫn chiếu có hướng (CÓ_ĐIỀU, CÓ_KHOẢN, CÓ_ĐIỂM, THUỘC_VỀ, DẪN_CHIẾU_ĐẾN...).\n"
             "• Lưu trữ trong diagramv2/examples/parsed-bo-luat-45-2019-qh14.json.",
             0.6, 2.65, 5.9, 2.3, DARK_EMERALD)

    add_card(s5, "Quy tắc Sinh Định danh Bất biến (Stable IDs)",
             "• Document: slug(filename_stem) -> 'bo-luat-45-2019-qh14'\n"
             "• Điều: dieu:bo-luat-45-2019-qh14:91\n"
             "• Khoản: khoan:bo-luat-45-2019-qh14:91:1\n"
             "• Điểm: diem:bo-luat-45-2019-qh14:91:1:a\n"
             "-> Cho phép join tất định giữa SQLite, PostgreSQL, Neo4j và Citation UI mà không phụ thuộc auto-increment DB ID.",
             6.8, 2.65, 5.9, 2.3, PURPLE)

    add_card(s5, "Ví dụ Cụ thể: Bóc tách Điều 91 (Mức lương tối thiểu)",
             "• Node Điều: dieu:bo-luat-45-2019-qh14:91 (Chương VI, 4 khoản trực thuộc).\n"
             "• Node Khoản 1: khoan:bo-luat-45-2019-qh14:91:1 (Parent: Điều 91, Path: Bộ Luật Lao Động > Chương VI > Điều 91 > Khoản 1).\n"
             "• Chuỗi trích dẫn hiển thị: 'Bộ Luật Lao Động 2019, Điều 91, Khoản 1' -> Đảm bảo tính minh bạch tuyệt đối.",
             0.6, 5.1, 12.1, 1.75, TEAL, SOFT_GREEN, 12.5, 10.5)

    # =========================================================================
    # SLIDE 6: CHIẾN LƯỢC CHUNKING DỮ LIỆU (HIERARCHY-AWARE CHUNKING)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Chunking Dữ liệu", "Chiến lược Chunking Phân cấp (Hierarchy-Aware Chunking) & 9 Loại Chunks", 6)

    add_card(s6, "Nguyên lý Phân đoạn theo Đơn vị Pháp lý",
             "• VlegalAI KHÔNG cắt văn bản theo số ký tự cố định hay generic token window (tránh làm đứt gãy câu, mất chủ thể hoặc mất điều kiện ngoại lệ).\n"
             "• Ranh giới phân đoạn đầu tiên là Đơn vị Cấu trúc Pháp lý (Legal Unit): Mỗi Khoản, Điểm hoặc Điều luật tự nhiên tạo thành một Retrieval Chunk độc lập.\n"
             "• Mỗi chunk luôn gắn liền với mã định danh node (node_id) và đường dẫn phân cấp (path_label) để phục vụ trích dẫn chính xác.\n\n"
             "Lợi ích cốt lõi:\n"
             "  [✓] Giữ trọn vẹn ngữ cảnh pháp lý và chủ thể thực hiện nghĩa vụ.\n"
             "  [✓] Khắc phục triệt để lỗi phân đoạn làm sai lệch ý nghĩa pháp lý.\n"
             "  [✓] Sẵn sàng trích dẫn trực tiếp S1-Sn đến từng Điều/Khoản/Điểm.",
             0.6, 1.55, 5.9, 5.3, DARK_EMERALD, WHITE, 13, 10.5)

    add_card(s6, "9 Loại Chunk Chuyên biệt trong Hệ thống",
             "1. article: Toàn văn Điều luật (dùng cho tra cứu mức điều & trích dẫn tổng quát).\n"
             "2. clause: Nội dung Khoản (đơn vị chứa quyền, nghĩa vụ, điều kiện cốt lõi).\n"
             "3. point: Nội dung Điểm (quy định chi tiết, danh sách hành vi vi phạm, mức phạt).\n"
             "4. table: Bảng lương tối thiểu vùng, phụ cấp, hệ số, tỷ lệ đóng bảo hiểm.\n"
             "5. structure: Tiêu đề & phạm vi Chương/Mục (phục vụ câu hỏi bao quát).\n"
             "6. document_intro: Lời mở đầu, căn cứ pháp lý, phạm vi điều chỉnh văn bản.\n"
             "7. semantic: Node Ontology (thuật ngữ, chủ thể, quy trình, mức phạt tiền).\n"
             "8. sliding: Cửa sổ trượt bổ trợ cho các điều/khoản dài vượt ngưỡng.\n"
             "9. document_structure: Thống kê số lượng chương/điều/khoản của văn bản.",
             6.8, 1.55, 5.9, 5.3, TEAL, SOFT_BLUE, 13, 10)

    # =========================================================================
    # SLIDE 7: CƠ CHẾ CẮT CHUỖI & CỬA SỔ TRƯỢT (SLIDING WINDOW)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Cắt chuỗi & Sliding Window", "Cơ chế Cắt Chuỗi (String Splitting) & Cửa sổ Trượt (Sliding Window Fallback)", 7)

    add_stat_box(s7, "W = 360", "ĐỘ DÀI CỬA SỔ (TỪ)", 0.6, 1.55, 2.85, 0.95, DARK_EMERALD)
    add_stat_box(s7, "O = 70", "ĐỘ GỐI ĐẦU OVERLAP", 3.68, 1.55, 2.85, 0.95, TEAL)
    add_stat_box(s7, "Δ = 290", "BƯỚC NHẢY STRIDE", 6.76, 1.55, 2.85, 0.95, PURPLE)
    add_stat_box(s7, "N ≤ 440", "NGƯỠNG GIỮ NGUYÊN 1 CHUNK", 9.84, 1.55, 2.85, 0.95, AMBER)

    add_card(s7, "1. Bộ đếm Token Tiếng Việt Chuẩn",
             "• Sử dụng biểu thức chính quy tiếng Việt: T = [0-9A-Za-zÀ-ỹĐđ]+\n"
             "• Đếm số từ: N(x) = len(VN_WORD_RE.findall(text))\n"
             "• Không bị lỗi cắt rời ký tự có dấu do phân tách UTF-8 multi-byte sai vị trí.",
             0.6, 2.65, 5.9, 2.0, DARK_EMERALD)

    add_card(s7, "2. Thuật toán Cửa sổ Trượt (app/legal_graphrag.py)",
             "• Nếu N(x) ≤ 440 từ: Giữ nguyên 1 chunk duy nhất, không cắt nhỏ.\n"
             "• Nếu N(x) > 440 từ: Tạo các cửa sổ trượt: C_j = x[j*Δ : j*Δ + W] với bước nhảy Δ = max(80, W - O) = 290 từ.\n"
             "• Tỷ lệ overlap danh định: ρ = 70 / 360 = 19.44%.",
             6.8, 2.65, 5.9, 2.0, TEAL)

    add_card(s7, "3. Quy tắc Kiểm soát Biên & Hợp đồng Văn bản Embedding",
             "• Cửa sổ đầu tiên giữ loại cấu trúc gốc (article, clause, point, table...); các cửa sổ tiếp theo được gắn nhãn sliding.\n"
             "• Phần đuôi cuối cùng nếu có số từ < 80 từ sẽ bị hủy bỏ (không tạo chunk vụn). Loại bỏ đoạn < 4 legal tokens.\n"
             "• Chuỗi đưa vào Embedding: X_c = title(c) + '\\n' + path_label(c) + '\\n' + chunk.text\n"
             "-> Đảm bảo một Khoản/Điểm ngắn vẫn mang đầy đủ ngữ cảnh của Điều luật và Tên văn bản chứa nó!",
             0.6, 4.8, 12.1, 2.05, PURPLE, SOFT_GREEN, 12.5, 10)

    # =========================================================================
    # SLIDE 8: QUÁ TRÌNH XỬ LÝ DỮ LIỆU & VECTOR EMBEDDING
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Vector Embedding", "Quy trình Xử lý Dữ liệu, Vector Embedding & Cơ chế Caching", 8)

    add_card(s8, "1. Chuẩn hóa Chuỗi Đầu vào Embedding",
             "Mỗi Chunk được đóng gói theo hợp đồng văn bản nghiêm ngặt:\n"
             "  X_c = title(c) || '\\n' || path(c) || '\\n' || text(c)\n"
             "Ví dụ:\n"
             "  'Mức lương tối thiểu'\n"
             "  'Bộ Luật Lao Động (45/2019/QH14) > Chương VI > Điều 91 > Khoản 1'\n"
             "  '1. Mức lương tối thiểu là mức lương thấp nhất được trả cho NLĐ...'",
             0.6, 1.55, 5.9, 2.5, DARK_EMERALD)

    add_card(s8, "2. Embedding Model & Task Type",
             "• Model: Vertex AI gemini-embedding-001\n"
             "• Chiều không gian: m = 1,024 chiều (float32).\n"
             "• Phân biệt Task Type chuyên biệt:\n"
             "  - Dữ liệu văn bản (Chunks): RETRIEVAL_DOCUMENT\n"
             "  - Câu hỏi người dùng (Queries): RETRIEVAL_QUERY\n"
             "• Chuẩn hóa L2: e(x) = z(x) / ||z(x)||_2.\n"
             "• Khoảng cách Cosine trên pgvector: d_cos = 1 - e(q)^T * e(c).",
             6.8, 1.55, 5.9, 2.5, TEAL)

    add_card(s8, "3. Quản lý Dung lượng & Checkpoint Batch",
             "• Mỗi vector float32 1024D = 4,096 bytes/chunk.\n"
             "• Tổng 32,334 chunks ≈ 126.3 MiB raw vector payload.\n"
             "• Checkpoint theo batch (640 chunks/batch) vào PostgreSQL: Khi gặp sự cố quota timeout có thể resume lại ngay lập tức mà không tính toán lại từ đầu.",
             0.6, 4.2, 5.9, 2.65, PURPLE)

    add_card(s8, "4. Caching theo Content Hash (SHA-256)",
             "• Hash nội dung: h_c = SHA256(X_c).\n"
             "• Vector chỉ được tái sử dụng khi khớp toàn bộ 7 yếu tố: chunk_id + h_c + provider + model + revision + task_type + dimension (1024).\n"
             "• Bất kỳ sửa đổi nào trong văn bản luật sẽ kích hoạt tính lại vector tự động, ngăn ngừa stale embedding.",
             6.8, 4.2, 5.9, 2.65, AMBER)

    # =========================================================================
    # SLIDE 9: CƠ CHẾ INDEXING ĐA TẦNG & HỢP NHẤT THỨ HẠNG RRF
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Indexing Đa tầng", "Hệ thống Indexing Đa Tầng Đồng bộ & Hợp nhất Thứ hạng RRF", 9)

    add_card(s9, "1. PGVECTOR (HNSW)",
             "• Cột vector(1024) trong bảng graphrag_chunk.\n"
             "• Thuật toán HNSW đồ thị tiệm cận phân tầng:\n"
             "  - vector_cosine_ops\n"
             "  - M = 16, ef_construction = 64\n"
             "• Truy hồi tương đồng ngữ nghĩa cực nhanh.",
             0.6, 1.55, 3.85, 2.9, TEAL)

    add_card(s9, "2. POSTGRESQL (GIN FTS)",
             "• Cột generated tsvector trên [title, citation, text].\n"
             "• GIN Index phục vụ tìm kiếm từ khóa chính xác.\n"
             "• Khớp số hiệu văn bản (45/2019/QH14), số điều luật, thuật ngữ chuyên ngành, mức tiền tệ.",
             4.75, 1.55, 3.85, 2.9, DARK_EMERALD)

    add_card(s9, "3. NEO4J GRAPH INDEX",
             "• 29,575 nodes & 108,368 edges.\n"
             "• Unique Constraint trên LegalChunk(chunk_id).\n"
             "• Fulltext Index trên [title, citation, text].\n"
             "• Mở rộng 2-hop theo các quan hệ pháp lý.",
             8.9, 1.55, 3.85, 2.9, PURPLE)

    add_card(s9, "Hợp nhất Thứ hạng RRF (Reciprocal Rank Fusion) tại Query Time",
             "• Kết hợp điểm số Vector và BM25 bằng công thức Weighted RRF (K = 60):\n"
             "  S_RRF(c) = w_v * (K + 1) / (K + r_v(c)) + w_b * (K + 1) / (K + r_b(c)) với trọng số thực nghiệm w_v = 0.55 (Vector) và w_b = 0.45 (BM25).\n"
             "• Điều chỉnh điểm số cơ sở theo mức độ bao phủ từ khóa: S_0(c) = S_RRF(c) / (r_f(c)^0.35) + 0.9 * m_c / min(|T_q|, 10) + B(c, q).\n"
             "• Atomic Activation: Chỉ kích hoạt phục vụ khi tổng số chunks và hash đồng bộ 100% giữa PostgreSQL và Neo4j.",
             0.6, 4.6, 12.1, 2.25, DARK_EMERALD, SOFT_GREEN, 12.5, 10)

    # =========================================================================
    # SLIDE 10: CHIẾN LƯỢC MÔ HÌNH: CÓ TRAIN / FINE-TUNE LLM HAY KHÔNG?
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Chiến lược Mô hình", "Có Train / Fine-Tune LLM hay không? Quyết định Kỹ thuật Cốt lõi (ADR-04)", 10)

    # Big NO Banner
    box_no = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.55), Inches(3.2), Inches(5.3))
    box_no.fill.solid()
    box_no.fill.fore_color.rgb = DARK_BG
    box_no.line.fill.background()
    tf_no = box_no.text_frame
    tf_no.word_wrap = True
    p_no0 = tf_no.paragraphs[0]
    p_no0.text = "KHÔNG"
    p_no0.alignment = PP_ALIGN.CENTER
    p_no0.font.size = Pt(40)
    p_no0.font.bold = True
    p_no0.font.color.rgb = ROSE

    p_no1 = tf_no.add_paragraph()
    p_no1.text = "Fine-tune hay Train lại Trọng số LLM\n(Trọng số mô hình giữ cố định: θ' = θ)"
    p_no1.alignment = PP_ALIGN.CENTER
    p_no1.font.size = Pt(13)
    p_no1.font.bold = True
    p_no1.font.color.rgb = WHITE

    p_no2 = tf_no.add_paragraph()
    p_no2.text = "\n• Không Pretraining\n• Không SFT (Supervised Fine-Tuning)\n• Không LoRA / QLoRA\n• Sử dụng Hosted Gemini 2.5 Flash qua Google Vertex AI"
    p_no2.font.size = Pt(11)
    p_no2.font.color.rgb = TEXT_LIGHT

    add_card(s10, "1. Khả năng Dẫn chứng & Truy xuất Nguồn gốc (Provenance)",
             "LLM fine-tuned lưu kiến thức vào 'hộp đen' trọng số, không thể bảo đảm trích dẫn chính xác tuyệt đối từng Điều/Khoản. RAG giữ nguyên vẹn chuỗi bằng chứng tường minh: Answer > Evidence > Chunk > Provision > Instrument.",
             4.0, 1.55, 4.3, 2.55, DARK_EMERALD)

    add_card(s10, "2. Tính Cập nhật & Xử lý Luật Hết hiệu lực (Freshness)",
             "Pháp luật thay đổi liên tục (Nghị định mới thay thế cũ). Nếu fine-tune, mỗi lần đổi luật phải thu thập dữ liệu & train lại rất tốn kém và dễ lẫn lộn. Với GraphRAG, chỉ cần re-index lại DB trong vài phút.",
             8.5, 1.55, 4.2, 2.55, TEAL)

    add_card(s10, "3. Kiểm soát Ảo giác Tuyệt đối (Hallucination Control)",
             "Lĩnh vực pháp lý đòi hỏi sự chính xác tuyệt đối. Việc ép mô hình chỉ sinh câu trả lời dựa trên tập bằng chứng được truy hồi (Evidence-first prompt contract) giúp triệt tiêu hoàn toàn nguy cơ bịa đặt số điều luật.",
             4.0, 4.3, 4.3, 2.55, PURPLE)

    add_card(s10, "4. Tối ưu hóa Trọng tâm Kỹ thuật",
             "Thay vì phân tán tài nguyên để train model, nhóm tập trung tối ưu: Cấu trúc đồ thị 10 tầng, Chiến lược Chunking phân cấp, Thuật toán Hybrid RRF, Prompt Engineering và Bộ lọc Kiểm chứng Trích dẫn (Evidence Gate).",
             8.5, 4.3, 4.2, 2.55, AMBER)

    # =========================================================================
    # SLIDE 11: ĐỒ THỊ TRI THỨC PHÁP LÝ 10 TẦNG (10-LAYER LEGAL KG)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Knowledge Graph", "Kiến trúc Đồ thị Tri thức Pháp lý 10 Tầng (10-Layer Legal Ontology)", 11)

    layers_left = [
        ("L0 - Nguồn & Hiệu lực", "Văn bản gốc, cơ quan ban hành, ngày hiệu lực, quan hệ sửa đổi/thay thế."),
        ("L1 - Cấu trúc Văn bản", "Chương > Mục > Điều > Khoản > Điểm, quan hệ HAS_PART, dẫn chiếu chéo."),
        ("L2 - Thuật ngữ & Chủ đề", "Định nghĩa pháp lý tự động, bản đồ chủ đề định hướng toàn bộ corpus."),
        ("L3 - Tiền lương & Thưởng", "Khoản thu nhập, kỳ hạn trả lương, lương tối thiểu vùng, công thức tính, tỷ lệ %."),
        ("L4 - Chủ thể & Quan hệ", "Người lao động, người sử dụng LĐ, tổ chức đại diện, quyền lợi & nghĩa vụ.")
    ]
    layers_right = [
        ("L5 - Quy trình & Thủ tục", "Thủ tục hành chính, hồ sơ yêu cầu, cơ quan giải quyết, thời hạn xử lý."),
        ("L6 - Thời gian & Thời hiệu", "Mốc thời gian luật định, thời hạn báo trước, thời hiệu khiếu nại."),
        ("L7 - Chế tài & Rủi ro", "Hành vi vi phạm, khung phạt tiền, biện pháp khắc phục hậu quả, mức độ rủi ro."),
        ("L8 - Vòng đời Quan hệ LĐ", "Giao kết > Thực hiện > Sửa đổi > Tạm hoãn > Chấm dứt hợp đồng lao động."),
        ("L9 - Án lệ & Thực tiễn", "Phán quyết tòa án, tình tiết cốt lõi (sẵn sàng mở rộng khi nạp án lệ).")
    ]

    y_pos = 1.55
    for title, desc in layers_left:
        add_card(s11, title, desc, 0.6, y_pos, 5.9, 0.95, DARK_EMERALD, WHITE, 11, 9)
        y_pos += 1.05

    y_pos = 1.55
    for title, desc in layers_right:
        add_card(s11, title, desc, 6.8, y_pos, 5.9, 0.95, TEAL, WHITE, 11, 9)
        y_pos += 1.05

    # =========================================================================
    # SLIDE 12: TRUY HỒI THÍCH ỨNG & EVIDENCE GATING
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "Truy hồi & Tạo sinh", "Phân Luồng Câu Hỏi Thích Ứng (Adaptive Routing) & Evidence Gating", 12)

    add_card(s12, "1. Phân luồng Độ phức tạp (Routing)",
             "• Single-hop (Đơn tầng): Tra cứu trực tiếp 1 điều khoản -> Chạy nhanh qua nhánh Hybrid RAG (pgvector + BM25) trong ~15ms.\n"
             "• Multi-hop (Đa tầng / Quan hệ): Câu hỏi so sánh, liên kết văn bản -> Kích hoạt duyệt đồ thị Neo4j 2-hop để gom đầy đủ căn cứ chéo.",
             0.6, 1.55, 5.9, 2.55, DARK_EMERALD)

    add_card(s12, "2. Xử lý Câu hỏi Phức & Nhiễu",
             "• Chuẩn hóa Teencode / Viết tắt: Chỉ chuẩn hóa khi phát hiện tín hiệu nhiễu, bảo toàn nguyên vẹn số tiền, ngày tháng.\n"
             "• Phân tách Khía cạnh (Facet Decomposition): Tách câu hỏi kép thành các vấn đề độc lập để truy hồi riêng biệt, tránh bỏ sót ý.",
             6.8, 1.55, 5.9, 2.55, TEAL)

    add_card(s12, "3. Cổng Bằng chứng (Evidence Gate)",
             "• Kiểm tra tính liên quan của Chunk trước khi đưa vào context prompt.\n"
             "• Nếu không tìm thấy căn cứ pháp lý phù hợp: Hệ thống chủ động trả về thông báo chưa đủ cơ sở thay vì để LLM tự suy diễn sai lệch.",
             0.6, 4.3, 5.9, 2.55, PURPLE)

    add_card(s12, "4. Xác thực Trích dẫn (Citation Validation)",
             "• Hệ thống gán nhãn trích dẫn [S1], [S2]... cho từng khẳng định.\n"
             "• Hậu kiểm tra (Post-validation): Kiểm tra từng [S_i] có thực sự chứa nội dung chứng minh cho câu trả lời hay không.\n"
             "• Loại bỏ trích dẫn mồ côi hoặc không liên quan.",
             6.8, 4.3, 5.9, 2.55, AMBER)

    # =========================================================================
    # SLIDE 13: PHƯƠNG PHÁP ĐÁNH GIÁ ĐỘ CHÍNH XÁC (EVALUATION METHODOLOGY)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "Đánh giá Độ chính xác", "Phương Pháp Đánh Giá Độ Chính Xác: Vector Chất Lượng 8 Chiều", 13)

    add_card(s13, "Độ chính xác trong Pháp lý là một Vector Đa chiều",
             "Q = (Q_ret, Q_faith, Q_rel, Q_fact, Q_cite, Q_facet, Q_temp, Q_lat)\n"
             "Exact match chuỗi không phù hợp cho pháp lý vì 2 câu trả lời khác từ ngữ vẫn có thể cùng biểu đạt đúng 1 quy định luật.",
             0.6, 1.55, 12.1, 1.1, DARK_EMERALD, SOFT_GREEN, 13, 10.5)

    add_card(s13, "1. FAITHFULNESS (Tính trung thực)",
             "Tỷ lệ các phát biểu trong câu trả lời được chứng minh trực tiếp bởi ngữ cảnh trích dẫn (|S_q| / |A_q|). Triệt tiêu ảo giác.",
             0.6, 2.8, 2.85, 2.3, DARK_EMERALD)

    add_card(s13, "2. ANSWER RELEVANCY",
             "Mức độ câu trả lời giải quyết đúng và trúng câu hỏi của người dùng, không trả lời lan man hoặc lạc đề.",
             3.68, 2.8, 2.85, 2.3, TEAL)

    add_card(s13, "3. CONTEXT PRECISION & RECALL",
             "Độ chính xác và độ bao phủ của tập chunk trích xuất so với tập bằng chứng chuẩn (Ground Truth).",
             6.76, 2.8, 2.85, 2.3, PURPLE)

    add_card(s13, "4. FACTUAL CORRECTNESS",
             "Độ chính xác về mặt sự thật và kết luận pháp lý so với câu trả lời mẫu của chuyên gia luật.",
             9.84, 2.8, 2.85, 2.3, AMBER)

    add_card(s13, "Context ID Precision & Recall",
             "Đo lường mức độ trùng khớp chính xác mã định danh điều khoản pháp lý: R_ID = |R_k ∩ G| / |G| và P_ID = |R_k ∩ G| / |R_k|.",
             0.6, 5.3, 5.9, 1.55, DARK_EMERALD)

    add_card(s13, "Bộ Benchmark RAGAS 100 Câu hỏi Thực tế",
             "100 câu hỏi luật lao động thực tế chia làm 3 nhóm: 50 Single-hop, 25 Multi-hop Specific, 25 Multi-hop Abstract. So sánh trực diện 4 kiến trúc: Dense RAG, LightRAG, GraphRAG, RAG+GraphRAG.",
             6.8, 5.3, 5.9, 1.55, TEAL)

    # =========================================================================
    # SLIDE 14: KẾT QUẢ THỰC NGHIỆM SO SÁNH 4 KIẾN TRÚC
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "Kết quả Thực nghiệm", "Bảng Kết Quả So Sánh Thực Nghiệm 4 Kiến Trúc (RAGAS Benchmark 100)", 14)

    # Benchmark Table
    table_shape = s14.shapes.add_table(5, 9, Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.3))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    for i in range(1, 9):
        table.columns[i].width = Inches(1.28)

    headers = ["Kiến trúc", "Faithfulness", "Ans. Rel.", "Ctx. Prec.", "Ctx. Rec.", "Factual", "ID Recall", "Latency", "Overall"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE

    rows_data = [
        ["Dense RAG", "0.9111", "0.9126", "0.7783", "0.8154", "0.3985", "0.6800", "15.8 ms", "0.7619*"],
        ["LightRAG", "0.9160", "0.8951", "0.6950", "0.8124", "0.4271", "0.6650", "148.8 ms", "0.7481"],
        ["GraphRAG", "0.9213*", "0.8641", "0.5558", "0.7820", "0.4003", "0.5750", "146.7 ms", "0.7041"],
        ["RAG+GraphRAG", "0.9403#", "--", "--", "--", "--", "0.6850*", "162.6 ms", "Checkpoint"]
    ]
    for row_idx, r in enumerate(rows_data):
        for col_idx, val in enumerate(r):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT_GREEN if row_idx == 0 else (SOFT_BLUE if row_idx == 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.bold = (col_idx == 0 or "*" in val)
            p.font.color.rgb = DARK_EMERALD if "*" in val else TEXT_DARK

    add_card(s14, "Nhận xét & Đánh giá Khoa học từ Dữ liệu Thực nghiệm",
             "• GraphRAG đạt độ trung thực Faithfulness cao nhất (0.9213) nhờ khả năng bao quát ngữ cảnh cấu trúc quan hệ, hạn chế tối đa ảo giác.\n"
             "• Dense RAG đạt Context Precision (0.7783) và Tốc độ truy hồi tốt nhất (15.8ms vs ~147ms của Graph) cho các câu hỏi tra cứu đơn tầng (Single-hop).\n"
             "• Minh chứng cho quyết định thiết kế Adaptive Routing: Dùng Hybrid Dense RAG làm mặc định cho câu hỏi đơn giản và chỉ kích hoạt GraphRAG cho câu hỏi phức tạp.",
             0.6, 4.1, 12.1, 2.7, DARK_EMERALD, SOFT_GREEN, 13, 11)

    # =========================================================================
    # SLIDE 15: HIỆU NĂNG THEO ĐỘ PHỨC TẠP CÂU HỎI
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    add_header(s15, "Phân tích Chuyên sâu", "Hiệu năng theo Phân loại Độ phức tạp Câu hỏi (Single-hop vs Multi-hop)", 15)

    add_card(s15, "1. Single-Hop Specific (50 câu)",
             "Điểm RAGAS: Dense 0.8114 | Light 0.8039 | Graph 0.7458\n\n"
             "• Tra cứu trực tiếp 1 điều luật, 1 mức lương tối thiểu, 1 thời hạn cụ thể.\n"
             "• Dense RAG chiếm ưu thế tuyệt đối về tốc độ và độ chính xác ngữ cảnh do không bị nhiễu bởi các node quan hệ xa.",
             0.6, 1.55, 3.85, 3.5, DARK_EMERALD)

    add_card(s15, "2. Multi-Hop Specific (25 câu)",
             "Điểm RAGAS: Dense 0.7625 | Light 0.7283 | Graph 0.6914\n\n"
             "• Câu hỏi kết hợp 2-3 điều luật (VD: Điều kiện hưởng lương + Trách nhiệm người sử dụng lao động).\n"
             "• Đồ thị hỗ trợ tìm kiếm đường dẫn chiếu chéo hiệu quả.",
             4.75, 1.55, 3.85, 3.5, TEAL)

    add_card(s15, "3. Multi-Hop Abstract (25 câu)",
             "Điểm RAGAS: Dense 0.6624 | Light 0.6563 | Graph 0.6334\n\n"
             "• Câu hỏi tình huống tổng hợp, so sánh quyền lợi giữa các nhóm lao động.\n"
             "• GraphRAG có mức độ suy giảm điểm số ít nhất (-0.112 vs -0.149 của Dense), chứng minh khả năng duy trì thông tin đa tầng.",
             8.9, 1.55, 3.85, 3.5, PURPLE)

    add_card(s15, "Kết luận về Bài toán Đánh đổi (Trade-Off): Độ chính xác vs Thời gian phản hồi",
             "• Đồ thị tri thức (Neo4j) giúp mở rộng quan hệ pháp lý rất tốt nhưng làm tăng độ trễ truy hồi từ 15.8ms lên ~147ms.\n"
             "• Thời gian sinh câu trả lời của LLM chiếm đa số (5.7s - 8.3s).\n"
             "• Thiết kế VlegalAI kết hợp: Seed Hybrid Retrieval + Bounded Graph Expansion (2-hop) là tối ưu nhất cho môi trường sản phẩm thực tế.",
             0.6, 5.25, 12.1, 1.6, AMBER, SOFT_AMBER, 12.5, 10.5)

    # =========================================================================
    # SLIDE 16: CƠ SỞ DỮ LIỆU & HẠ TẦNG GOOGLE CLOUD PLATFORM
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    add_header(s16, "Triển khai & Hạ tầng", "Mô hình Cơ Sở Dữ Liệu Vật Lý & Hạ Tầng Google Cloud Platform (GCP)", 16)

    add_stat_box(s16, "23", "POSTGRES BASE TABLES", 0.6, 1.55, 2.2, 0.95, DARK_EMERALD)
    add_stat_box(s16, "1", "MATERIALIZED VIEW", 3.05, 1.55, 2.2, 0.95, TEAL)
    add_stat_box(s16, "18", "ALEMBIC MIGRATIONS", 5.5, 1.55, 2.2, 0.95, PURPLE)
    add_stat_box(s16, "Cloud Run", "SERVERLESS INGRESS", 7.95, 1.55, 2.3, 0.95, DARK_EMERALD)
    add_stat_box(s16, "CI / CD", "WIF + GITHUB ACTIONS", 10.5, 1.55, 2.2, 0.95, AMBER)

    add_card(s16, "Identity & Chat (7 tables)",
             "app_user, sso_identity, user_feedback, conversation, chat_message, conversation_summary, chat_answer_feedback",
             0.6, 2.65, 2.85, 2.4, DARK_EMERALD, WHITE, 11, 9.5)

    add_card(s16, "Content & Catalog (6+1)",
             "article, artifact, signature_packet, legal_document, legal_chunk, legal_answer_cache + legal_catalog_corpus (MV)",
             3.68, 2.65, 2.85, 2.4, TEAL, WHITE, 11, 9.5)

    add_card(s16, "GraphRAG (4 tables)",
             "graphrag_chunk (vector 1024D), graphrag_embedding_checkpoint, graphrag_index_metadata, graphrag_law_version",
             6.76, 2.65, 2.85, 2.4, PURPLE, WHITE, 11, 9.5)

    add_card(s16, "Runtime & Queue (6 tables)",
             "guest_rate_limit, kombu_queue, kombu_message, celery_taskmeta, celery_tasksetmeta, alembic_version",
             9.84, 2.65, 2.85, 2.4, AMBER, WHITE, 11, 9.5)

    add_card(s16, "Quy trình CI/CD & Triển khai Bất biến (Immutable Container Image)",
             "• Triển khai qua GitHub Actions với Workload Identity Federation (WIF) - Không lưu trữ Service Account key trong repo.\n"
             "• 1 Container Image duy nhất cho cả Web API, Celery Worker, Scheduler và Reindex Job.\n"
             "• Tự động chạy Unit Test (492 backend tests, 16 frontend tests) và Database Migration trước khi release.",
             0.6, 5.25, 12.1, 1.6, DARK_EMERALD, SOFT_GREEN, 12.5, 10)

    # =========================================================================
    # SLIDE 17: TÍNH NĂNG SẢN PHẨM & TRẢI NGHIỆM NGƯỜI DÙNG
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    add_header(s17, "Tính năng Sản phẩm", "Bộ Tính Năng Toàn Diện Hỗ Trợ Doanh Nghiệp & Người Lao Động", 17)

    add_card(s17, "1. Tra cứu & Hỏi đáp Pháp lý Chuyên sâu",
             "• Hỏi đáp tự nhiên bằng tiếng Việt có dấu/không dấu.\n"
             "• Cung cấp trích dẫn nguồn luật chính xác [S1-Sn].\n"
             "• Nhấp chuột xem toàn văn điều khoản và hiệu lực.",
             0.6, 1.55, 5.9, 2.55, DARK_EMERALD)

    add_card(s17, "2. Soạn thảo Hợp đồng Lao động Chuẩn luật",
             "• Hỗ trợ tạo hợp đồng theo biểu mẫu chuẩn pháp luật.\n"
             "• Tự động điền mức lương tối thiểu vùng theo quy định mới nhất.\n"
             "• Xuất file DOCX quy chuẩn.",
             6.8, 1.55, 5.9, 2.55, TEAL)

    add_card(s17, "3. Rà soát & Đánh giá Rủi ro Hợp đồng",
             "• Upload hợp đồng (DOCX/PDF/Ảnh OCR) để quét điều khoản bất lợi.\n"
             "• Cảnh báo vi phạm luật lao động, thiếu quyền lợi bắt buộc.\n"
             "• Gợi ý điều khoản sửa đổi an toàn cho doanh nghiệp.",
             0.6, 4.3, 5.9, 2.55, PURPLE)

    add_card(s17, "4. So sánh Đối chiếu Hai Bản Hợp đồng",
             "• Tự động căn chỉnh các điều khoản tương ứng giữa 2 phiên bản.\n"
             "• Phát hiện nội dung thêm, bớt, sửa đổi chi tiết.\n"
             "• Đánh giá thay đổi mức độ rủi ro pháp lý giữa 2 bản thảo.",
             6.8, 4.3, 5.9, 2.55, AMBER)

    # =========================================================================
    # SLIDE 18: BẢO MẬT, QUYỀN RIÊNG TƯ & QUẢN TRỊ DỮ LIỆU
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    add_header(s18, "Bảo mật & Quản trị", "Bảo Mật Quyền Riêng Tư & Cơ Chế Quản Trị Dữ Liệu (Data Governance)", 18)

    add_card(s18, "1. Bảo vệ Tài liệu Tải lên Riêng tư",
             "• Hợp đồng và tệp đính kèm người dùng được mã hóa cấp ứng dụng (AES-GCM) bằng khóa bảo mật.\n"
             "• Cách ly hoàn toàn theo User Ownership: Không bao giờ đưa tài liệu riêng tư vào Vector DB hay Knowledge Graph dùng chung.\n"
             "• Token truy cập đính kèm có thời hạn (expiring attachment token).",
             0.6, 1.55, 5.9, 2.55, ROSE)

    add_card(s18, "2. Xác thực & Phân quyền Chuẩn OAuth",
             "• Đăng nhập qua Google OIDC với giao thức Authorization Code Flow kết hợp PKCE (Proof Key for Code Exchange).\n"
             "• Cookie phiên làm việc HttpOnly, Secure, SameSite bảo vệ chống tấn công XSS/CSRF.\n"
             "• Tuyệt đối không để lộ API Key của Gemini/GCP trên trình duyệt client.",
             6.8, 1.55, 5.9, 2.55, DARK_EMERALD)

    add_card(s18, "3. Chu kỳ Cập nhật Văn bản Pháp luật (Freshness)",
             "• Quản lý trạng thái văn bản: Đang có hiệu lực, Hết hiệu lực, Bị sửa đổi, Chưa có hiệu lực.\n"
             "• Scheduler kiểm tra định kỳ nguồn văn bản chính thống.\n"
             "• Tự động kích hoạt luồng Re-index cập nhật đồ thị tri thức khi có Nghị định mới ban hành.",
             0.6, 4.3, 5.9, 2.55, TEAL)

    add_card(s18, "4. Thu thập Phản hồi & Tái sinh Câu trả lời (HITL)",
             "• Ghi nhận đánh giá Hài lòng (GOOD) / Không hài lòng (BAD) của người dùng.\n"
             "• Cho phép người dùng yêu cầu tái sinh câu trả lời với ngữ cảnh được điều chỉnh.\n"
             "• Phục vụ cải tiến chất lượng hệ thống mà không làm lộ dữ liệu nhạy cảm.",
             6.8, 4.3, 5.9, 2.55, AMBER)

    # =========================================================================
    # SLIDE 19: TỔNG KẾT ĐÓNG GÓP & MINH CHỨNG ĐỘ TIN CẬY
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    add_header(s19, "Tổng kết Đồ án", "Tổng Kết Đóng Góp Cốt Lõi & Hướng Phát Triển Tương Lai", 19)

    add_card(s19, "Đóng Góp Khoa Học & Kỹ Thuật Cốt Lõi",
             "  [✓] Xây dựng thành công Đồ thị Tri thức Pháp luật Lao động 10 tầng với 29,575 nodes và 108,368 edges từ 74 văn bản quy chuẩn.\n"
             "  [✓] Đề xuất chiến lược Chunking theo cấp bậc pháp lý (Hierarchy-Aware Chunking) và cửa sổ trượt W=360, O=70 bảo toàn 100% ngữ cảnh điều luật.\n"
             "  [✓] Hiện thực hóa kiến trúc Adaptive Hybrid RAG (pgvector + BM25 + Neo4j) kết hợp Reciprocal Rank Fusion và Evidence Gate chống ảo giác.\n"
             "  [✓] Xây dựng bộ Benchmark chuẩn RAGAS 100 câu hỏi và đánh giá thực nghiệm toàn diện trên 8 tiêu chí.",
             0.6, 1.55, 5.9, 3.5, DARK_EMERALD, SOFT_GREEN, 13, 10.5)

    add_card(s19, "Hạn Chế Hiện Tại & Hướng Phát Triển",
             "• Thời gian sinh câu trả lời của LLM còn độ trễ đuôi (p95 latency) do phụ thuộc vào cloud API của Vertex AI.\n"
             "• Bộ dữ liệu đánh giá cần mở rộng thêm chuyên gia thẩm định chéo (Inter-rater Agreement).\n"
             "• Kế hoạch tương lai:\n"
             "  + Huấn luyện mô hình Reranker cục bộ để giảm độ trễ.\n"
             "  + Nạp thêm kho Án lệ Tòa án (Layer 9).\n"
             "  + Mở rộng sang Luật Doanh nghiệp và Luật Thuế.",
             6.8, 1.55, 5.9, 3.5, AMBER, WHITE, 13, 10.5)

    add_card(s19, "Minh chứng Chất lượng & Độ tin cậy Kỹ thuật",
             "• 492 Backend Unit Tests Passed (100% pass)  |  16 Frontend Tests Passed  |  18 Alembic Migrations\n"
             "• Toàn bộ mã nguồn, dữ liệu đánh giá, database migrations và tài liệu báo cáo đều có thể tái lập (Reproducible) 100%.",
             0.6, 5.25, 12.1, 1.6, TEAL, SOFT_BLUE, 12.5, 11)

    # =========================================================================
    # SLIDE 20: LỜI CẢM ƠN & DEFENSE Q&A
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    bg20 = s20.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg20.fill.solid()
    bg20.fill.fore_color.rgb = DARK_BG
    bg20.line.fill.background()

    tb_end = s20.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.5))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True
    p_end0 = tf_end.paragraphs[0]
    p_end0.text = "CẢM ƠN QUÝ THẦY CÔ TRONG HỘI ĐỒNG"
    p_end0.alignment = PP_ALIGN.CENTER
    p_end0.font.size = Pt(28)
    p_end0.font.bold = True
    p_end0.font.color.rgb = WHITE

    p_end1 = tf_end.add_paragraph()
    p_end1.text = "NHÓM VLEGALAI SẴN SÀNG LẮNG NGHE Ý KIẾN ĐÓNG GÓP & TRẢ LỜI CÂU HỎI"
    p_end1.alignment = PP_ALIGN.CENTER
    p_end1.font.size = Pt(14)
    p_end1.font.bold = True
    p_end1.font.color.rgb = TEAL

    # 5 Key Technical Recap Cards for Defense Q&A
    qa_list = [
        ("1. Chunking & Cắt chuỗi", "Hierarchy-Aware theo Điều/Khoản/Điểm. Cửa sổ trượt W=360, Overlap=70 từ; đếm từ bằng regex tiếng Việt."),
        ("2. Có Train LLM Không?", "KHÔNG train/fine-tune weight. Dùng Hosted Gemini 2.5 Flash + RAG/GraphRAG để giữ 100% provenance & cập nhật luật mới."),
        ("3. Đánh Giá Độ Chính Xác", "Vector 8 chiều qua RAGAS (100 câu hỏi luật thực tế): Faithfulness đạt 0.9213, ID Recall đạt 0.6800."),
        ("4. Xử Lý & Vector Embedding", "Input contract: Title + Path + Text. Model gemini-embedding-001 (1024D), chuẩn hóa L2, caching SHA-256."),
        ("5. Cấu Trúc Document Parsing", "Deterministic State Machine phân tích cú pháp DOCX thành Envelope 4 mảng: source, document, nodes, edges, chunks.")
    ]

    x_positions = [0.8, 4.8, 8.8, 2.8, 6.8]
    y_positions = [2.8, 2.8, 2.8, 4.8, 4.8]
    for idx, (q_title, q_desc) in enumerate(qa_list):
        box_qa = s20.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(y_positions[idx]), Inches(3.7), Inches(1.7))
        box_qa.fill.solid()
        box_qa.fill.fore_color.rgb = CARD_DARK
        box_qa.line.color.rgb = EMERALD
        tf_q = box_qa.text_frame
        tf_q.word_wrap = True
        tf_q.margin_left = tf_q.margin_top = tf_q.margin_right = tf_q.margin_bottom = Inches(0.12)
        pq0 = tf_q.paragraphs[0]
        pq0.text = q_title
        pq0.font.size = Pt(12)
        pq0.font.bold = True
        pq0.font.color.rgb = AMBER
        pq1 = tf_q.add_paragraph()
        pq1.text = q_desc
        pq1.font.size = Pt(9.5)
        pq1.font.color.rgb = TEXT_LIGHT

    # Save to both target file names
    print(f"Saving to {OUT_SLIDES}...")
    prs.save(str(OUT_SLIDES))
    print(f"Saving to {OUT_SLDIES}...")
    prs.save(str(OUT_SLDIES))
    print("Successfully created presentation slides!")

if __name__ == "__main__":
    create_presentation()
