"""Comprehensive 45-slide Graduation Defense Presentation Generator for VlegalAI.
Sourced directly from the Codebase and Final Capstone Report.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(r"f:\VlegalAI")
OUT_SLIDES = ROOT / "slides.pptx"
OUT_SLDIES = ROOT / "sldies.pptx"

# Modern Color Palette
DARK_BG = RGBColor(15, 23, 42)       # #0F172A Slate Dark
CARD_DARK = RGBColor(30, 41, 59)     # #1E293B Slate Navy
TEXT_LIGHT = RGBColor(248, 250, 252) # #F8FAFC
TEXT_MUTED = RGBColor(100, 116, 139) # #64748B
TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A

EMERALD = RGBColor(16, 185, 129)     # #10B981 Vibrant Green
DARK_EMERALD = RGBColor(15, 76, 58)  # #0F4C3A Primary Forest Green
TEAL = RGBColor(14, 165, 233)        # #0EA5E9 Sky Blue
PURPLE = RGBColor(139, 92, 246)      # #8B5CF6 Purple
AMBER = RGBColor(245, 158, 11)       # #F59E0B Amber Gold
ROSE = RGBColor(244, 63, 94)         # #F43F5E Crimson Rose

BG_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC
WHITE = RGBColor(255, 255, 255)
BORDER_LIGHT = RGBColor(226, 232, 240) # #E2E8F0
SOFT_GREEN = RGBColor(236, 253, 245)   # #ECFDF5
SOFT_BLUE = RGBColor(240, 249, 255)    # #F0F9FF
SOFT_AMBER = RGBColor(254, 243, 199)   # #FEF3C7
SOFT_PURPLE = RGBColor(245, 243, 255)  # #F5F3FF

def build_45_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, section: str, title: str, page_num: int, total_pages: int = 45):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.fill.background()

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.45))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = SOFT_GREEN
        top_bar.line.fill.background()

        p_sec = top_bar.text_frame.paragraphs[0]
        p_sec.text = f"  VLEGALAI | {section.upper()}"
        p_sec.font.size = Pt(11)
        p_sec.font.bold = True
        p_sec.font.color.rgb = DARK_EMERALD

        tb_page = slide.shapes.add_textbox(Inches(11.5), Inches(0.05), Inches(1.5), Inches(0.35))
        p_pg = tb_page.text_frame.paragraphs[0]
        p_pg.text = f"{page_num:02d} / {total_pages:02d}"
        p_pg.alignment = PP_ALIGN.RIGHT
        p_pg.font.size = Pt(11)
        p_pg.font.bold = True
        p_pg.font.color.rgb = TEXT_MUTED

        tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12.133), Inches(0.65))
        p_t = tb_title.text_frame.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(19)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.22), Inches(1.2), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = EMERALD
        line.line.fill.background()

        fline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.05), Inches(12.133), Inches(0.02))
        fline.fill.solid()
        fline.fill.fore_color.rgb = BORDER_LIGHT
        fline.line.fill.background()

        tb_f = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(8.0), Inches(0.3))
        pf = tb_f.text_frame.paragraphs[0]
        pf.text = "VlegalAI: Vietnamese Labor Law GraphRAG & AI Assistant | Capstone Defense"
        pf.font.size = Pt(9)
        pf.font.color.rgb = TEXT_MUTED

        tb_fr = slide.shapes.add_textbox(Inches(9.5), Inches(7.1), Inches(3.2), Inches(0.3))
        pfr = tb_fr.text_frame.paragraphs[0]
        pfr.text = f"Slide {page_num:02d}"
        pfr.alignment = PP_ALIGN.RIGHT
        pfr.font.size = Pt(9)
        pfr.font.color.rgb = TEXT_MUTED

    def add_card(slide, title: str, text: str, left, top, width, height, accent_color=EMERALD, bg_color=WHITE, title_size=12, body_size=10):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = BORDER_LIGHT
        box.line.width = Pt(1)

        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent_color
        stripe.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(left + 0.16), Inches(top + 0.1), Inches(width - 0.28), Inches(height - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(title_size)
        p0.font.bold = True
        p0.font.color.rgb = accent_color

        p1 = tf.add_paragraph()
        p1.text = text
        p1.font.size = Pt(body_size)
        p1.font.color.rgb = TEXT_DARK

    def add_stat_box(slide, value: str, label: str, left, top, width, height=0.9, color=DARK_EMERALD):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BORDER_LIGHT
        box.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.06), Inches(width), Inches(height - 0.12))
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = value
        p0.alignment = PP_ALIGN.CENTER
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = color

        p1 = tf.add_paragraph()
        p1.text = label
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(8.8)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # 45 SLIDES DEFINITION
    # =========================================================================

    # 1. Cover
    s = prs.slides.add_slide(blank_layout)
    bg1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_BG
    bg1.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EMERALD
    bar.line.fill.background()

    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.65), Inches(4.5), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_DARK
    pill.line.color.rgb = EMERALD
    p_pill = pill.text_frame.paragraphs[0]
    p_pill.text = "GRADUATION THESIS DEFENSE - CAPSTONE 2026"
    p_pill.alignment = PP_ALIGN.CENTER
    p_pill.font.size = Pt(10.5)
    p_pill.font.bold = True
    p_pill.font.color.rgb = EMERALD

    tb_c1 = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.8))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    p_c1 = tf_c1.paragraphs[0]
    p_c1.text = "VLEGALAI: VIETNAMESE LABOR-LAW GRAPHRAG"
    p_c1.font.size = Pt(30)
    p_c1.font.bold = True
    p_c1.font.color.rgb = WHITE
    p_c2 = tf_c1.add_paragraph()
    p_c2.text = "Hệ thống Trợ lý Pháp lý Lao động Việt Nam ứng dụng GraphRAG và Large Language Models"
    p_c2.font.size = Pt(18)
    p_c2.font.color.rgb = TEAL

    box_sv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.3), Inches(5.6), Inches(2.4))
    box_sv.fill.solid()
    box_sv.fill.fore_color.rgb = CARD_DARK
    box_sv.line.color.rgb = BORDER_LIGHT
    tf_sv = box_sv.text_frame
    tf_sv.word_wrap = True
    p_sv0 = tf_sv.paragraphs[0]
    p_sv0.text = "SINH VIÊN THỰC HIỆN:"
    p_sv0.font.size = Pt(11.5)
    p_sv0.font.bold = True
    p_sv0.font.color.rgb = EMERALD
    p_sv1 = tf_sv.add_paragraph()
    p_sv1.text = "• Trần Tuấn Kiệt (Trưởng nhóm) - QE180152\n• Lê Thanh Đạt - QE170186\n• Phan Bảo Khánh - DE170648"
    p_sv1.font.size = Pt(12.5)
    p_sv1.font.color.rgb = TEXT_LIGHT

    box_gv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.3), Inches(5.6), Inches(2.4))
    box_gv.fill.solid()
    box_gv.fill.fore_color.rgb = CARD_DARK
    box_gv.line.color.rgb = BORDER_LIGHT
    tf_gv = box_gv.text_frame
    tf_gv.word_wrap = True
    p_gv0 = tf_gv.paragraphs[0]
    p_gv0.text = "GIẢNG VIÊN HƯỚNG DẪN:"
    p_gv0.font.size = Pt(11.5)
    p_gv0.font.bold = True
    p_gv0.font.color.rgb = AMBER
    p_gv1 = tf_gv.add_paragraph()
    p_gv1.text = "• GVHD: ThS. Lê Trung Hiếu\n• Đồng GVHD: ThS. Trương Ngọc Hùng\n• Hội đồng: FPT University Quy Nhơn"
    p_gv1.font.size = Pt(12.5)
    p_gv1.font.color.rgb = TEXT_LIGHT

    tb_foot = s.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5))
    p_foot = tb_foot.text_frame.paragraphs[0]
    p_foot.text = "Tháng 08 / 2026 | Chuyên ngành Trí tuệ Nhân tạo (Artificial Intelligence)"
    p_foot.font.size = Pt(10.5)
    p_foot.font.color.rgb = TEXT_MUTED

    # 2. Team Organization & Roles
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Tổ chức Đội ngũ", "Phân công Trách nhiệm & Mô hình Quản trị Đề tài", 2)
    add_card(s, "1. Trần Tuấn Kiệt (Backend & AI Lead)",
             "• Kiến trúc hệ thống, Data pipeline, Deterministic Parser, Chunking & Indexing.\n"
             "• Hiện thực hóa Hybrid RAG, pgvector HNSW, Neo4j GraphRAG, Adaptive Routing.\n"
             "• Tích hợp API FastAPI, Gunicorn, PostgreSQL và Docker GCP Cloud Run.",
             0.6, 1.45, 3.85, 4.4, DARK_EMERALD)
    add_card(s, "2. Lê Thanh Đạt (Product & Frontend Lead)",
             "• Phát triển giao diện người dùng React 18 + Vite SPA, Responsive UX.\n"
             "• Công cụ Soạn thảo, Rà soát và So sánh Hợp đồng Lao động.\n"
             "• Trực quan hóa trích dẫn [S1-Sn], xem toàn văn điều khoản và quản trị phiên chat.",
             4.75, 1.45, 3.85, 4.4, TEAL)
    add_card(s, "3. Phan Bảo Khánh (Quality & Evaluation Lead)",
             "• Xây dựng bộ Benchmark RAGAS 100 câu hỏi luật lao động thực tế.\n"
             "• Thiết kế Test Suite (492 backend unit tests, 16 frontend integration tests).\n"
             "• Thu thập số liệu thực nghiệm, tính toán Bootstrap 95% CI và tài liệu báo cáo.",
             8.9, 1.45, 3.85, 4.4, PURPLE)
    add_card(s, "Nguyên tắc Đóng góp Đội ngũ (Equal Contribution)",
             "Mỗi thành viên chịu trách nhiệm toàn diện trên module của mình, phối hợp liên tục qua CI/CD và kiểm thử hồi quy.",
             0.6, 6.0, 12.1, 0.85, DARK_EMERALD, SOFT_GREEN, 11, 9.5)

    # 3. Capstone Scope & Compliance
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Phạm vi & Chuẩn mực", "Phạm vi Đề tài & Bảng Đối soát Tiêu chí Đánh giá (Checklist Traceability)", 3)
    add_stat_box(s, "74", "VĂN BẢN QUY PHẠM", 0.6, 1.45, 2.2, 0.85, DARK_EMERALD)
    add_stat_box(s, "32,334", "CHUNKS PHÁP LÝ", 3.05, 1.45, 2.2, 0.85, TEAL)
    add_stat_box(s, "29,575", "NODES ĐỒ THỊ", 5.5, 1.45, 2.2, 0.85, PURPLE)
    add_stat_box(s, "108,368", "EDGES QUAN HỆ", 7.95, 1.45, 2.3, 0.85, AMBER)
    add_stat_box(s, "492 / 492", "BACKEND TESTS PASS", 10.5, 1.45, 2.2, 0.85, ROSE)
    add_card(s, "Phạm vi Pháp lý Trọng tâm",
             "• Bộ luật Lao động 45/2019/QH14 và các Luật liên quan (Luật Việc làm, Luật BHXH, Luật ATVSLĐ).\n"
             "• Hệ thống Nghị định hướng dẫn (Nghị định 145/2020/NĐ-CP, Nghị định 293/2025/NĐ-CP về lương tối thiểu vùng...).\n"
             "• Các Thông tư, Quyết định và Văn bản hợp nhất chuyên ngành lao động - tiền lương.",
             0.6, 2.45, 5.9, 3.4, DARK_EMERALD)
    add_card(s, "Tuân thủ Chuẩn mực Capstone (FPT University SE/AI)",
             "• AI1-AI5: Dữ liệu chuẩn hóa, mô hình hóa đồ thị tri thức, đánh giá thực nghiệm RAGAS khách quan.\n"
             "• SE1-SE5: Kiến trúc 4 tầng chuẩn mực, 18 Alembic migrations, bảo mật OIDC PKCE, CI/CD tự động.\n"
             "• PE1-PE3: Mã nguồn 100% tái lập (reproducible), đầy đủ test catalogue, runbooks và tài liệu vận hành.",
             6.8, 2.45, 5.9, 3.4, TEAL)
    add_card(s, "Cam kết Trách nhiệm",
             "Hệ thống là công cụ hỗ trợ ra quyết định (Decision Support), không thay thế tư vấn pháp lý có chứng chỉ hành nghề.",
             0.6, 6.0, 12.1, 0.85, AMBER, SOFT_AMBER, 11, 9.5)

    # 4. Vietnamese Legal Information Challenges
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Thách thức Pháp lý", "Bốn Thách Thức Cốt Tử trong Xử Lý Pháp Luật Lao Động Việt Nam", 4)
    add_card(s, "1. Cấu trúc Phân cấp & Phạm vi",
             "Một quy định không nằm ở 1 câu đơn lẻ mà phụ thuộc vào cấu trúc Văn bản > Chương > Mục > Điều > Khoản > Điểm. Cắt sai làm mất chủ thể và phạm vi áp dụng.",
             0.6, 1.45, 5.9, 2.1, ROSE)
    add_card(s, "2. Dẫn chiếu & Điều kiện Ngoại lệ",
             "Các điều luật thường dẫn chiếu chéo (Ví dụ: 'Trừ trường hợp quy định tại Khoản 2 Điều 156...'). LLM thông thường bỏ qua các ngoại lệ này dẫn đến tư vấn sai.",
             6.8, 1.45, 5.9, 2.1, AMBER)
    add_card(s, "3. Biến động Hiệu lực theo Thời gian",
             "Pháp luật thay đổi liên tục. Sử dụng tri thức cũ từ các Nghị định đã hết hiệu lực có thể gây thiệt hại tài chính và rủi ro pháp lý nghiêm trọng.",
             0.6, 3.7, 5.9, 2.1, PURPLE)
    add_card(s, "4. Ngôn ngữ Người dùng Phức tạp",
             "Người dùng thường hỏi bằng khẩu ngữ, teencode, viết tắt (VD: 'hđlđ', 'bhxh', 'lương net') hoặc đặt câu hỏi phức chứa nhiều vấn đề đồng thời.",
             6.8, 3.7, 5.9, 2.1, TEAL)
    add_card(s, "Kết luận",
             "Hỏi đáp pháp lý đòi hỏi Phục hồi Thông tin Chính xác (Grounded Retrieval) và Tạo sinh có Ràng buộc (Constrained Generation).",
             0.6, 5.95, 12.1, 0.9, DARK_EMERALD, SOFT_GREEN, 11, 9.5)

    # 5. Problem Statement & Research Formulation
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Mục tiêu Nghiên cứu", "Phát biểu Bài toán & Khung Nghiên cứu Thiết kế (Design-Science)", 5)
    add_card(s, "Phát biểu Bài toán (Problem Statement)",
             "Cho một câu hỏi tiếng Việt của người dùng và tệp đính kèm tùy chọn:\n"
             "1. Truy hồi tập bằng chứng pháp luật lao động nhỏ nhất nhưng đầy đủ nhất (Smallest sufficient set).\n"
             "2. Trả lời đầy đủ mọi khía cạnh của câu hỏi mà không bỏ sót ý.\n"
             "3. Cung cấp trích dẫn nguồn luật chính xác, có thể kiểm chứng (Verifiable citations).\n"
             "4. Bảo vệ an toàn tuyệt đối dữ liệu riêng tư của người dùng.",
             0.6, 1.45, 5.9, 4.3, DARK_EMERALD, WHITE, 13, 10.5)
    add_card(s, "Mục tiêu Kỹ thuật Cụ thể",
             "• O1 (Grounded Quality): Đạt Faithfulness >= 0.90 trên bộ dữ liệu kiểm thử chuẩn.\n"
             "• O2 (Latency & Efficiency): Truy hồi < 20ms cho single-hop; E2E latency có thể chấp nhận được.\n"
             "• O3 (Provenance & Safety): 100% trích dẫn [S1-Sn] dẫn đến điều khoản văn bản còn hiệu lực.\n"
             "• O4 (Reproducibility): Mọi kết quả benchmark, index build và test suite đều tái lập được.",
             6.8, 1.45, 5.9, 4.3, TEAL, WHITE, 13, 10.5)
    add_card(s, "Phương pháp Luận Nghiên cứu",
             "Ứng dụng chu trình Design-Science Research: Vấn đề -> Yêu cầu -> Xây dựng Artifact -> Thực nghiệm RAGAS -> Đóng gói Hệ thống.",
             0.6, 5.9, 12.1, 0.95, PURPLE, SOFT_PURPLE, 11, 9.5)

    # 6. Research Questions & Hypotheses
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Câu hỏi Nghiên cứu", "Bốn Câu Hỏi Nghiên Cứu (Research Questions) & Giả Thuyết", 6)
    add_card(s, "RQ1: So sánh Chất lượng giữa các Kiến trúc",
             "Dense RAG, LightRAG và GraphRAG khác biệt như thế nào về Faithfulness, Answer Relevancy và Factual Correctness trên bài toán hỏi đáp luật lao động?",
             0.6, 1.45, 5.9, 2.1, DARK_EMERALD)
    add_card(s, "RQ2: Đánh đổi giữa Độ trễ & Độ phức tạp",
             "Chi phí thời gian (Retrieval & E2E Latency) của việc duyệt đồ thị 2-hop so với tìm kiếm vector đơn thuần là bao nhiêu? Có đáng để áp dụng cho mọi query?",
             6.8, 1.45, 5.9, 2.1, TEAL)
    add_card(s, "RQ3: Tác động của Độ sâu Suy luận",
             "Khi câu hỏi chuyển từ Single-hop sang Multi-hop và Multi-abstract, kiến trúc nào duy trì được độ bao quát thông tin tốt nhất?",
             0.6, 3.7, 5.9, 2.1, PURPLE)
    add_card(s, "RQ4: Tính Minh bạch & Tái lập Kỹ thuật",
             "Làm thế nào để xây dựng một quy trình kiểm thử và benchmark minh bạch, không phóng đại số liệu, có thể tái lập 100% từ mã nguồn?",
             6.8, 3.7, 5.9, 2.1, AMBER)
    add_card(s, "Quyết định Kiến trúc Định hướng",
             "Kết quả giải đáp RQ1-RQ4 trực tiếp dẫn đến thiết kế Adaptive Routing: Phân luồng câu hỏi theo độ phức tạp thay vì ép mọi truy vấn qua đồ thị.",
             0.6, 5.95, 12.1, 0.9, DARK_EMERALD, SOFT_GREEN, 11, 9.5)

    # 7. Related Work & Research Lineage
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Tổng quan Nghiên cứu", "Nghiên cứu Liên quan & Quyết định Kỹ thuật Kế thừa", 7)
    add_card(s, "1. Legal QA & Benchmarking (LegalBench, LexGLUE)",
             "• Guha et al. (2023), Chalkidis et al. (2022): Chỉ ra rằng LLM thuần túy gặp khó khăn lớn trong việc suy luận quy tắc pháp lý và trích dẫn.\n"
             "• Kế thừa: VlegalAI tiếp cận bài toán pháp lý như một bài toán Grounding & Provenance thay vì sinh ngôn ngữ tự do.",
             0.6, 1.45, 5.9, 2.6, DARK_EMERALD)
    add_card(s, "2. Dense Retrieval & Hybrid RAG (DPR, BM25, RRF)",
             "• Karpukhin et al. (2020), Cormack et al. (2009): Vector dày xử lý tốt ngữ nghĩa nhưng yếu ở từ khóa chính xác (số hiệu luật, số tiền).\n"
             "• Kế thừa: Kết hợp Dense Cosine + BM25 FTS qua Weighted Reciprocal Rank Fusion (w_v=0.55, w_b=0.45, K=60).",
             6.8, 1.45, 5.9, 2.6, TEAL)
    add_card(s, "3. GraphRAG & Legal Knowledge Graphs (Microsoft GraphRAG, LightRAG, Vuong et al.)",
             "• Edge et al. (2024), Guo et al. (2024), Vuong et al. (2023 - Đồ thị án lệ Việt Nam): Đồ thị giúp kết nối các thực thể và dẫn chiếu chéo.\n"
             "• Kế thừa: Xây dựng đồ thị 10 tầng có định hướng, giới hạn mở rộng 2-hop và bắt buộc mọi đường đi kết thúc tại LegalChunk có hiệu lực.",
             0.6, 4.2, 12.1, 2.6, PURPLE, SOFT_BLUE, 12, 10)

    # 8. 10-Phase Engineering Evolution
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Tiến trình Phát triển", "Quá trình Tiến Hóa Kỹ Thuật 10 Giai Đoạn (Engineering Evolution)", 8)
    phases_1 = [
        ("Phase 1: LLM thuần túy", "Quan sát thấy ảo giác điều luật -> Chuyển sang bài toán RAG có căn cứ."),
        ("Phase 2: Khảo sát Fine-tuning", "Không đảm bảo trích dẫn, khó cập nhật luật -> Quyết định giữ cố định trọng số LLM."),
        ("Phase 3: Dense RAG cơ bản", "Bỏ sót số hiệu luật, số tiền -> Bổ sung nhánh từ khóa BM25/GIN."),
        ("Phase 4: Hợp nhất Thứ hạng", "Điểm số vector và BM25 không cùng thang đo -> Ứng dụng Weighted RRF (K=60)."),
        ("Phase 5: Cấu trúc Pháp lý", "Cắt văn bản ngẫu nhiên làm gãy câu -> Đề xuất Hierarchy-Aware Chunking.")
    ]
    phases_2 = [
        ("Phase 6: Giới hạn Đồ thị", "Duyệt đồ thị mọi query gây chậm -> Xây dựng Adaptive Routing."),
        ("Phase 7: Xử lý Nhiễu & Kép", "Teencode làm sai lệch query -> Phân tách Facet & chuẩn hóa có chọn lọc."),
        ("Phase 8: Phân biệt Lỗi", "Timeout API bị nhầm là thiếu luật -> Thêm Evidence Gate & Typed Fallback."),
        ("Phase 9: Vận hành Cloud", "Tránh drift dữ liệu -> Đóng gói Container, WIF, 18 Migrations, Cache SHA-256."),
        ("Phase 10: Minh chứng Đánh giá", "Tránh số liệu ảo -> Xây dựng RAGAS Benchmark 100 câu và 492 Unit Tests.")
    ]
    y_pos = 1.45
    for title, desc in phases_1:
        add_card(s, title, desc, 0.6, y_pos, 5.9, 0.95, DARK_EMERALD, WHITE, 11, 9)
        y_pos += 1.05
    y_pos = 1.45
    for title, desc in phases_2:
        add_card(s, title, desc, 6.8, y_pos, 5.9, 0.95, TEAL, WHITE, 11, 9)
        y_pos += 1.05

    # 9. Use Cases & System Boundary
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Phạm vi Nghiệp vụ", "10 Vùng Ranh Giới Hệ Thống & Trường Hợp Sử Dụng (System Boundaries)", 9)
    b1 = [
        ("D.1 Xác thực & Phiên", "Google OIDC, PKCE, HttpOnly Cookie, bảo vệ tài khoản."),
        ("D.2 Cổng Web & API", "FastAPI async, React SPA, quản lý rate limit và routing."),
        ("D.3 Hiểu Câu hỏi & Facet", "Chuẩn hóa teencode, phân tách câu hỏi kép thành các ý độc lập."),
        ("D.4 Truy hồi Lai Thích ứng", "Kết hợp pgvector HNSW, GIN FTS và Neo4j Graph."),
        ("D.5 Cổng Bằng chứng & Trích dẫn", "Kiểm duyệt chunk, gán nhãn [S1-Sn] và hậu kiểm trích dẫn.")
    ]
    b2 = [
        ("D.6 Tài liệu Đính kèm Riêng tư", "Mã hóa AES-GCM, cách ly theo User, xử lý OCR hợp đồng."),
        ("D.7 Quản lý Hội thoại & HITL", "Lưu lịch sử, tóm tắt ngữ cảnh, thu thập đánh giá Good/Bad."),
        ("D.8 Tin tức Pháp lý & Freshness", "Quản trị bài viết, tự động theo dõi văn bản luật mới."),
        ("D.9 Pipeline Đồ thị Tri thức", "Phân tích cú pháp DOCX, sinh vector 1024D, sync PostgreSQL/Neo4j."),
        ("D.10 Triển khai Đám mây GCP", "Cloud Run, Cloud SQL, Celery Worker, CI/CD GitHub Actions.")
    ]
    y_pos = 1.45
    for title, desc in b1:
        add_card(s, title, desc, 0.6, y_pos, 5.9, 0.95, DARK_EMERALD, WHITE, 11, 9)
        y_pos += 1.05
    y_pos = 1.45
    for title, desc in b2:
        add_card(s, title, desc, 6.8, y_pos, 5.9, 0.95, TEAL, WHITE, 11, 9)
        y_pos += 1.05

    # 10. End-to-End User Experience & Web Interface
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Giao diện Người dùng", "Trải Nghiệm Người Dùng & Các Luồng Nghiệp Vụ Chính", 10)
    add_card(s, "1. Hỏi đáp Pháp lý Đa phương thức",
             "• Nhập câu hỏi tự nhiên hoặc tải lên tệp hợp đồng (DOCX/PDF/Ảnh).\n"
             "• Hiển thị câu trả lời định dạng Markdown chuyên nghiệp kèm các thẻ trích dẫn [S1], [S2]...\n"
             "• Nhấp vào thẻ trích dẫn để mở Drawer xem toàn văn điều luật, cơ quan ban hành và tình trạng hiệu lực.",
             0.6, 1.45, 5.9, 2.6, DARK_EMERALD)
    add_card(s, "2. Công cụ Hợp đồng Lao động Toàn diện",
             "• Soạn thảo: Hỗ trợ sinh dự thảo hợp đồng chuẩn luật lao động, tự động áp mức lương tối thiểu vùng.\n"
             "• Rà soát: Quét điều khoản vi phạm pháp luật, phát hiện điều khoản bất lợi cho người lao động.\n"
             "• So sánh: Đối chiếu 2 bản hợp đồng, làm nổi bật nội dung sửa đổi và đánh giá biến động rủi ro.",
             6.8, 1.45, 5.9, 2.6, TEAL)
    add_card(s, "3. Quản lý Phiên & Phản hồi Human-in-the-Loop",
             "• Lưu trữ cây hội thoại đa phiên, tự động tóm tắt ngữ cảnh khi hội thoại kéo dài.\n"
             "• Nút đánh giá Hài lòng (Good) / Không hài lòng (Bad) và yêu cầu Tái sinh câu trả lời (Regenerate).\n"
             "• Toàn bộ dữ liệu phản hồi được ẩn danh hóa và lưu trữ phục vụ đánh giá chất lượng.",
             0.6, 4.2, 12.1, 2.6, PURPLE, SOFT_BLUE, 12.5, 10.5)

    # 11. System Architecture: Runtime & Deployment Views
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Kiến trúc Vận hành", "Mô hình Vận Hành Runtime & Kiến Trúc Triển Khai Cloud", 11)
    add_card(s, "Client Layer (SPA)",
             "React 18 + TypeScript + Vite. Đóng gói SPA tĩnh, giao tiếp qua HTTPS với API Gateway. Quản lý trạng thái bằng Zustand.",
             0.6, 1.45, 3.85, 2.5, DARK_EMERALD)
    add_card(s, "Application Gateway (Cloud Run)",
             "FastAPI async serving với Gunicorn worker pool (8 vCPU, 16GB RAM). Điều phối xác thực, phân luồng routing và tạo prompt.",
             4.75, 1.45, 3.85, 2.5, TEAL)
    add_card(s, "Background Workers (Celery)",
             "Celery Worker + Redis/PostgreSQL Queue xử lý tác vụ nặng: OCR hợp đồng, Re-index đồ thị tri thức, Quét tin tức định kỳ.",
             8.9, 1.45, 3.85, 2.5, PURPLE)
    add_card(s, "Quy trình Xử lý Yêu cầu (Request Lifecycle)",
             "1. Client gửi câu hỏi -> 2. Xác thực Session Cookie -> 3. Phân luồng Single-hop / Multi-hop -> 4. Truy hồi Vector & BM25 & Neo4j -> 5. Evidence Gating -> 6. Gọi Vertex AI Gemini 2.5 Flash -> 7. Citation Post-validation -> 8. Stream kết quả về Client.",
             0.6, 4.15, 12.1, 2.65, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    # 12. Stakeholders & Business Pain Points
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Đối tượng Sử dụng", "Các Nhóm Người Dùng Mục Tiêu & Nỗi Đau Thực Tế (Pain Points)", 12)
    add_card(s, "1. Người Lao Động (Employees)",
             "• Nỗi đau: Không am hiểu thuật ngữ pháp lý, khó tra cứu chế độ thai sản, trợ cấp thôi việc, lương làm thêm giờ.\n"
             "• Giải pháp: Hỏi đáp bằng ngôn ngữ tự nhiên, nhận câu trả lời dễ hiểu kèm căn cứ luật để tự bảo vệ quyền lợi.",
             0.6, 1.45, 5.9, 2.6, DARK_EMERALD)
    add_card(s, "2. Cán bộ Nhân sự / Doanh nghiệp (HR & SMBs)",
             "• Nỗi đau: Tốn nhiều thời gian rà soát hợp đồng, dễ vi phạm quy định mới về mức lương tối thiểu vùng dẫn đến bị phạt.\n"
             "• Giải pháp: Công cụ soạn thảo và rà soát hợp đồng tự động, cập nhật tức thì quy định mới nhất.",
             6.8, 1.45, 5.9, 2.6, TEAL)
    add_card(s, "3. Chuyên viên Pháp chế (Legal Researchers)",
             "• Nỗi đau: Tra cứu dẫn chiếu chéo giữa Luật và các Nghị định hướng dẫn mất nhiều thời gian lật giở tài liệu.\n"
             "• Giải pháp: Đồ thị tri thức 10 tầng mở rộng liên kết điều luật và trích dẫn chuẩn xác trong vài giây.",
             0.6, 4.2, 12.1, 2.6, PURPLE, SOFT_PURPLE, 12.5, 10.5)

    # 13. Knowledge Base Architecture: 10-Layer Model
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Đồ thị Tri thức", "Kiến Trúc Đồ Thị Tri Thức Pháp Lý 10 Tầng (10-Layer Knowledge Graph)", 13)
    add_card(s, "Tầng Nền tảng (Foundation Plane)",
             "• L0 (Nguồn & Hiệu lực): Văn bản gốc, cơ quan ban hành, ngày có hiệu lực, quan hệ sửa đổi/thay thế/bãi bỏ.\n"
             "• L1 (Cấu trúc Văn bản): Document > Chapter > Section > Article > Clause > Point. Quan hệ HAS_PART và DẪN_CHIẾU_ĐẾN.",
             0.6, 1.45, 5.9, 2.5, DARK_EMERALD)
    add_card(s, "Tầng Ngữ nghĩa Pháp lý (Domain Plane)",
             "• L2 (Thuật ngữ): Định nghĩa pháp lý, từ đồng nghĩa.\n"
             "• L3 (Tiền lương): Khoản thu nhập, lương tối thiểu, tỷ lệ %.\n"
             "• L4 (Chủ thể): Quyền và nghĩa vụ người LĐ, NSDLĐ.\n"
             "• L5 (Quy trình): Thủ tục hành chính, hồ sơ, thời hạn giải quyết.\n"
             "• L6 (Thời gian): Mốc thời gian luật định, thời hiệu khiếu nại.\n"
             "• L7 (Chế tài): Hành vi vi phạm, khung phạt tiền, khắc phục.",
             6.8, 1.45, 5.9, 2.5, TEAL)
    add_card(s, "Tầng Ngữ cảnh Đa Văn bản (Cross-Document Plane)",
             "• L8 (Vòng đời Lao động): Chuỗi giai đoạn Giao kết -> Thực hiện -> Tạm hoãn -> Chấm dứt HĐLĐ.\n"
             "• L9 (Án lệ & Thực tiễn): Phán quyết tòa án và giải thích áp dụng luật (sẵn sàng nạp án lệ).\n"
             "• Nguyên tắc cốt tử: Mọi đường đi ngữ nghĩa trên đồ thị bắt buộc phải kết thúc tại một LegalChunk có hiệu lực!",
             0.6, 4.15, 12.1, 2.65, PURPLE, SOFT_BLUE, 12.5, 10.5)

    # 14. Ingestion, Embedding & Sync Pipeline
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Pipeline Dữ liệu", "Quy Trình Ingestion, Vector Embedding & Đồng Bộ Hóa Đa Cơ Sở Dữ Liệu", 14)
    add_card(s, "1. Ingestion & Parse Tất định",
             "Đọc file DOCX theo luồng khối (paragraphs + tables). State machine nhận diện cấu trúc Chương/Điều/Khoản/Điểm, sinh JSON Envelope.",
             0.6, 1.45, 3.85, 2.6, DARK_EMERALD)
    add_card(s, "2. Embedding & Caching",
             "Ghép chuỗi Title + Path + Text. Gọi gemini-embedding-001 (1024D). Cache SHA-256 nội dung, checkpoint batch 640 chunks.",
             4.75, 1.45, 3.85, 2.6, TEAL)
    add_card(s, "3. Đồng bộ & Kích hoạt",
             "Đẩy dữ liệu vào SQLite, PostgreSQL (pgvector + GIN) và Neo4j Graph. Kiểm tra đối soát số lượng và hash trước khi kích hoạt.",
             8.9, 1.45, 3.85, 2.6, PURPLE)
    add_card(s, "Nguyên tắc Bất biến của Index (Index Integrity Contract)",
             "Không bao giờ trộn lẫn vector của các model/revision/dimension khác nhau. Khi một văn bản luật sửa đổi, hệ thống tính lại content hash và cập nhật đồng bộ trên toàn bộ 3 database.",
             0.6, 4.25, 12.1, 2.55, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    # 15. Query-Time Flow: Adaptive Hybrid Retrieval & Generation
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Luồng Truy vấn", "Quy Trình Xử Lý Truy Vấn Thời Gian Thực (Query-Time Flow)", 15)
    add_card(s, "Giai đoạn 1: Hiểu & Lập kế hoạch (Planning)",
             "• Kiểm tra tín hiệu nhiễu: Chuẩn hóa teencode nếu cần.\n"
             "• Phân loại độ phức tạp: Single-hop vs Multi-hop.\n"
             "• Phân tách câu hỏi kép thành danh sách Facet cần trả lời.",
             0.6, 1.45, 5.9, 2.5, DARK_EMERALD)
    add_card(s, "Giai đoạn 2: Truy hồi Lai (Hybrid Retrieval)",
             "• Tìm kiếm Vector Cosine trên pgvector (1024D).\n"
             "• Tìm kiếm từ khóa BM25 trên GIN tsvector.\n"
             "• Hợp nhất thứ hạng qua Weighted RRF (0.55/0.45).\n"
             "• (Nếu Multi-hop): Mở rộng 2-hop trên Neo4j.",
             6.8, 1.45, 5.9, 2.5, TEAL)
    add_card(s, "Giai đoạn 3: Evidence Gate & Tạo sinh có Ràng buộc (Generation)",
             "• Evidence Gate lọc bỏ các chunk không liên quan; nếu không đủ căn cứ sẽ trả về thông báo từ chối suy diễn.\n"
             "• Đóng gói Context Prompt với hợp đồng bằng chứng nghiêm ngặt -> Gọi Gemini 2.5 Flash.\n"
             "• Citation Validator kiểm tra từng khẳng định có đúng với nguồn trích dẫn [S1-Sn] trước khi trả lời.",
             0.6, 4.15, 12.1, 2.65, PURPLE, SOFT_GREEN, 12.5, 10.5)

    # 16. Chat Runtime Sequence & Attachment Handling
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Xử lý Đính kèm", "Trình Tự Runtime & Cơ Chế Xử Lý Tài Liệu Đính Kèm Riêng Tư", 16)
    add_card(s, "Quy trình Xử lý File Đính kèm (Contract / Document Upload)",
             "1. Người dùng upload tệp DOCX / PDF / Ảnh chụp hợp đồng qua API /chat/attachments.\n"
             "2. Hệ thống kiểm tra dung lượng, định dạng và quét mã độc.\n"
             "3. Trích xuất văn bản (bằng python-docx, PyMuPDF hoặc OCR Tesseract/Vision).\n"
             "4. Mã hóa nội dung bằng AES-GCM với khóa bí mật của ứng dụng.\n"
             "5. Cấp phát Token đính kèm có thời hạn (Expiring attachment token) gắn chặt với User ID.",
             0.6, 1.45, 5.9, 4.3, DARK_EMERALD, WHITE, 12.5, 10)
    add_card(s, "Nguyên tắc Bảo mật Cách ly Dữ liệu",
             "• Dữ liệu tệp đính kèm là ngữ cảnh hội thoại riêng tư (Private Conversation Context).\n"
             "• TUYỆT ĐỐI KHÔNG đưa văn bản hợp đồng của người dùng vào Vector Database hay Knowledge Graph dùng chung.\n"
             "• Khi hỏi về hợp đồng, hệ thống ưu tiên tìm kiếm trên ngữ cảnh tệp đính kèm trước khi đối chiếu với luật chung.",
             6.8, 1.45, 5.9, 4.3, ROSE, WHITE, 12.5, 10)
    add_card(s, "Chống Rò rỉ Dữ liệu",
             "Không bao giờ truyền API key hoặc thông tin chứng thực của Cloud về phía trình duyệt Client.",
             0.6, 5.9, 12.1, 0.95, AMBER, SOFT_AMBER, 11, 9.5)

    # 17. Frontend & Backend Package Structure
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Cấu trúc Mã nguồn", "Tổ Chức Gói Mã Nguồn Backend & Frontend (Package Architecture)", 17)
    add_card(s, "Backend Architecture (app/)",
             "• app/api.py: 41 HTTP endpoints (Auth, Chat, Contracts, Laws, Articles).\n"
             "• app/services/: chat_policy.py, retrieval.py, contract_tools.py, embeddings.py.\n"
             "• app/legal_graphrag.py & legal_ontology.py: Parser, Chunking, Neo4j Builder.\n"
             "• app/models.py & schemas.py: 14 SQLAlchemy models, Pydantic schemas.\n"
             "• migrations/: 18 Alembic migration scripts quản lý database schema.",
             0.6, 1.45, 5.9, 4.3, DARK_EMERALD, WHITE, 12.5, 10)
    add_card(s, "Frontend Architecture (frontend/src/)",
             "• components/chat/: ChatContainer, MessageList, CitationDrawer, AttachmentPicker.\n"
             "• components/contracts/: ContractDraftModal, ContractReviewView, ContractCompare.\n"
             "• services/api.ts: Typed API client, Token refresh, Error boundary.\n"
             "• store/: Zustand state stores quản lý Active Conversation, Workspace Artifacts.\n"
             "• Build: Vite TypeScript -> 1.43 kB HTML, 192 kB CSS, 300 kB JS (90 kB gzip).",
             6.8, 1.45, 5.9, 4.3, TEAL, WHITE, 12.5, 10)
    add_card(s, "Đảm bảo Chất lượng Mã nguồn",
             "Ruff linter + TypeScript compiler + Pytest suite (492 tests) chạy tự động trên mỗi commit qua GitHub Actions.",
             0.6, 5.9, 12.1, 0.95, PURPLE, SOFT_PURPLE, 11, 9.5)

    # 18. Cloud SQL Physical Schema (23 Tables + 1 MV)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Cơ sở Dữ liệu", "Mô Hình Cơ Sở Dữ Liệu Vật Lý Cloud SQL: 23 Tables + 1 Materialized View", 18)
    add_stat_box(s, "23", "BASE TABLES", 0.6, 1.45, 2.2, 0.85, DARK_EMERALD)
    add_stat_box(s, "1", "MATERIALIZED VIEW", 3.05, 1.45, 2.2, 0.85, TEAL)
    add_stat_box(s, "18", "ALEMBIC MIGRATIONS", 5.5, 1.45, 2.2, 0.85, PURPLE)
    add_stat_box(s, "1024-D", "PGVECTOR COLUMN", 7.95, 1.45, 2.3, 0.85, DARK_EMERALD)
    add_stat_box(s, "0018", "ALEMBIC HEAD", 10.5, 1.45, 2.2, 0.85, AMBER)

    add_card(s, "1. Identity & Chat (7 tables)",
             "app_user, sso_identity, user_feedback, conversation, chat_message, conversation_summary, chat_answer_feedback",
             0.6, 2.45, 2.85, 2.4, DARK_EMERALD, WHITE, 11, 9.5)
    add_card(s, "2. Content & Catalog (6+1)",
             "article, artifact, signature_packet, legal_document, legal_chunk, legal_answer_cache + legal_catalog_corpus (MV)",
             3.68, 2.45, 2.85, 2.4, TEAL, WHITE, 11, 9.5)
    add_card(s, "3. GraphRAG (4 tables)",
             "graphrag_chunk (vector 1024D), graphrag_embedding_checkpoint, graphrag_index_metadata, graphrag_law_version",
             6.76, 2.45, 2.85, 2.4, PURPLE, WHITE, 11, 9.5)
    add_card(s, "4. Runtime & Queue (6 tables)",
             "guest_rate_limit, kombu_queue, kombu_message, celery_taskmeta, celery_tasksetmeta, alembic_version",
             9.84, 2.45, 2.85, 2.4, AMBER, WHITE, 11, 9.5)
    add_card(s, "Quy tắc Đếm & Quản trị Schema",
             "23 PostgreSQL base tables + 1 materialized view (legal_catalog_corpus). Schema được quản trị bằng 18 file migration đồng bộ.",
             0.6, 5.05, 12.1, 1.8, DARK_EMERALD, SOFT_GREEN, 12, 10)

    # 19. ERD: Identity, Chat & Human-in-the-Loop
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Sơ đồ ERD", "ERD Nhóm Identity, Hội Thoại & Phản Hồi Người Dùng (HITL)", 19)
    add_card(s, "Khóa Chính & Quan Hệ Khóa Ngoại",
             "• app_user: Bảng gốc định danh người dùng qua UUID. Ràng buộc 1-N với sso_identity (Google OAuth).\n"
             "• conversation: Thuộc sở hữu của app_user. Chứa tiêu đề, thời gian tạo, cập nhật.\n"
             "• chat_message: Liên kết với conversation_id. Lưu role (user/assistant), content, metadata trích dẫn.\n"
             "• chat_answer_feedback: Lưu đánh giá GOOD/BAD, lý do, liên kết message_id và user_id phục vụ audit.",
             0.6, 1.45, 5.9, 4.3, DARK_EMERALD, WHITE, 12.5, 10.5)
    add_card(s, "Nguyên tắc Thiết kế Dữ liệu",
             "• Xóa theo tầng (Cascade Delete) có kiểm soát khi người dùng yêu cầu xóa cuộc trò chuyện.\n"
             "• Mọi tin nhắn phản hồi đều giữ lại snapshot trích dẫn [S1-Sn] tại thời điểm tạo ra câu trả lời.\n"
             "• Đảm bảo tính toàn vẹn tham chiếu (Referential Integrity) trên 100% các bảng giao dịch.",
             6.8, 1.45, 5.9, 4.3, TEAL, WHITE, 12.5, 10.5)
    add_card(s, "Minh chứng",
             "Source: diagramv2/04-postgres-erd-identity-chat.mmd | 7 bảng cốt lõi.",
             0.6, 5.9, 12.1, 0.95, PURPLE, SOFT_PURPLE, 11, 9.5)

    # 20. ERD: Content, Legal Catalog & GraphRAG
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Sơ đồ ERD", "ERD Nhóm Văn Bản Pháp Luật, Vector Chunk & Runtime Cache", 20)
    add_card(s, "Cấu trúc Bảng GraphRAG & pgvector",
             "• graphrag_chunk: Bảng lưu trữ 32,334 chunks. Cột embedding kiểu vector(1024), cột tsvector tự sinh, doc_id, node_id, chunk_type.\n"
             "• graphrag_embedding_checkpoint: Lưu trạng thái xử lý embedding theo batch, hỗ trợ resume khi đứt kết nối.\n"
             "• graphrag_index_metadata: Lưu fingerprint, model revision, kích thước vector và ngày build.\n"
             "• graphrag_law_version: Quản lý phiên bản hiệu lực của từng văn bản luật.",
             0.6, 1.45, 5.9, 4.3, PURPLE, WHITE, 12.5, 10.5)
    add_card(s, "Cấu trúc Bảng Nội dung & Cache",
             "• legal_document: Siêu dữ liệu 74 văn bản gốc, cơ quan ban hành, ngày có hiệu lực.\n"
             "• legal_answer_cache: Bộ đệm câu trả lời theo hash câu hỏi và fingerprint index, giảm chi phí gọi LLM.\n"
             "• artifact & signature_packet: Lưu trữ bản thảo hợp đồng và gói chuẩn bị ký số.",
             6.8, 1.45, 5.9, 4.3, AMBER, WHITE, 12.5, 10.5)
    add_card(s, "Minh chứng",
             "Source: diagramv2/05-postgres-erd-content-runtime.mmd | Đảm bảo tính đồng bộ tuyệt đối giữa SQL và Graph.",
             0.6, 5.9, 12.1, 0.95, DARK_EMERALD, SOFT_GREEN, 11, 9.5)

    # 21. Physical Database Guarantees & Constraints
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Ràng buộc CSDL", "Các Cam Đoan & Ràng Buộc Kỹ Thuật của Cơ Sở Dữ Liệu Vật Lý", 21)
    add_card(s, "1. Định danh Khách quan (UUID & Stable IDs)",
             "• Bảng ứng dụng dùng UUID v4 làm khóa chính, chống đoán số thứ tự.\n"
             "• Bảng pháp lý dùng Stable Natural IDs (VD: dieu:bo-luat-45-2019-qh14:91) giúp liên kết tất định với Neo4j và SQLite.",
             0.6, 1.45, 5.9, 2.5, DARK_EMERALD)
    add_card(s, "2. Toàn vẹn Chỉ mục (Index Constraints)",
             "• HNSW Cosine Index trên cột embedding vector(1024).\n"
             "• GIN Index trên tsvector hỗ trợ tìm kiếm toàn văn tiếng Việt.\n"
             "• B-tree Index trên các khóa ngoại (doc_id, node_id, user_id, conversation_id).",
             6.8, 1.45, 5.9, 2.5, TEAL)
    add_card(s, "3. Quản trị Tiến hóa Schema (Alembic)",
             "• 18 file migration kiểm soát phiên bản CSDL nghiêm ngặt.\n"
             "• CI/CD tự động chạy alembic upgrade head trước khi khởi động ứng dụng.\n"
             "• Không bao giờ sửa đổi schema thủ công trên môi trường Production.",
             0.6, 4.15, 12.1, 2.65, PURPLE, SOFT_GREEN, 12.5, 10.5)

    # 22. Class & Module Responsibility View
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Thiết kế Module", "Phân Định Trách Nhiệm Module & Lớp Xử Lý (Class Responsibilities)", 22)
    add_card(s, "ChatOrchestrator (app/services/)",
             "Lớp điều phối trung tâm: Tiếp nhận request -> Phân tích câu hỏi -> Gọi Retrieval Service -> Evidence Gate -> Gọi Gemini LLM -> Xác thực Citation -> Ghi nhận DB.",
             0.6, 1.45, 5.9, 2.5, DARK_EMERALD)
    add_card(s, "LegalGraphBuilder (app/legal_graphrag.py)",
             "Lớp xây dựng đồ thị & chunking: Đọc file DOCX -> Bóc tách cấu trúc -> Sinh Stable IDs -> Tạo 9 loại chunk -> Sinh vector -> Đẩy vào SQLite, Postgres, Neo4j.",
             6.8, 1.45, 5.9, 2.5, TEAL)
    add_card(s, "AdaptiveRetrievalService (app/services/retrieval.py)",
             "Lớp truy hồi thích ứng: Thực thi song song truy vấn pgvector HNSW và GIN BM25 -> Weighted RRF -> Mở rộng Neo4j khi cần -> Khử trùng lặp theo chunk_id.",
             0.6, 4.15, 5.9, 2.65, PURPLE)
    add_card(s, "EvidenceGate & CitationValidator",
             "Lớp kiểm duyệt chất lượng: Kiểm tra độ phù hợp của chunk trước khi đưa vào prompt -> Gán nhãn [S1-Sn] -> Hậu kiểm tra chứng minh pháp lý của câu trả lời.",
             6.8, 4.15, 5.9, 2.65, AMBER)

    # 23. Legal Knowledge Graph: Node & Edge Distributions
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Thống kê Đồ thị", "Phân Bố Thực Thể (Nodes) & Quan Hệ (Edges) trong Đồ Thị Tri Thức", 23)
    add_stat_box(s, "29,575", "TỔNG NODES", 0.6, 1.45, 2.85, 0.85, DARK_EMERALD)
    add_stat_box(s, "108,368", "TỔNG EDGES", 3.68, 1.45, 2.85, 0.85, TEAL)
    add_stat_box(s, "41", "LOẠI NODES", 6.76, 1.45, 2.85, 0.85, PURPLE)
    add_stat_box(s, "43", "LOẠI QUAN HỆ", 9.84, 1.45, 2.85, 0.85, AMBER)

    add_card(s, "Các Loại Node Chiếm Tỷ Trọng Lớn",
             "• Khoản (Clause): 13,714 nodes (đơn vị chứa quy tắc cốt lõi).\n"
             "• Điểm (Point): 9,820 nodes (quy định chi tiết, mức phạt).\n"
             "• Điều (Article): 4,192 nodes (đơn vị trích dẫn chính).\n"
             "• Chương (Chapter): 434 nodes | Mục (Section): 272 nodes.\n"
             "• Thuật ngữ (Term): 249 nodes | Văn bản (Document): 74 nodes.",
             0.6, 2.45, 5.9, 4.3, DARK_EMERALD, WHITE, 12, 10)
    add_card(s, "Các Loại Quan Hệ Phổ Biến",
             "• THUỘC_VỀ: 28,727 edges (quan hệ phân cấp con -> cha).\n"
             "• QUY_ĐỊNH_TẠI: 19,249 edges (nghĩa vụ gắn với điều khoản).\n"
             "• CÓ_KHOẢN: 13,714 edges | CÓ_ĐIỂM: 9,820 edges.\n"
             "• DẪN_CHIẾU_ĐẾN: 9,904 edges (liên kết chéo giữa các điều luật).\n"
             "• THUỘC_CHỦ_ĐỀ: 7,401 edges | CÓ_ĐIỀU: 4,192 edges.",
             6.8, 2.45, 5.9, 4.3, TEAL, WHITE, 12, 10)

    # 24. 10-Layer Legal Knowledge Graph Breakdown (L0-L4)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Chi tiết Đồ thị", "Phân Tích Chi Tiết 5 Tầng Tri Thức Đầu Tiên (L0 - L4)", 24)
    add_card(s, "L0: Nguồn & Hiệu lực (Provenance & Validity)",
             "Lưu trữ cơ quan ban hành, ngày ban hành, ngày có hiệu lực, trạng thái (Còn hiệu lực, Hết hiệu lực, Bị sửa đổi). Lọc bỏ ngay lập tức các văn bản đã bị bãi bỏ.",
             0.6, 1.45, 12.1, 0.95, DARK_EMERALD, WHITE, 11.5, 9.5)
    add_card(s, "L1: Cấu trúc Văn bản (Document Structure)",
             "Cấu trúc cây Document > Chapter > Section > Article > Clause > Point. Hỗ trợ tra cứu chính xác theo số điều, số khoản và mở rộng cha-con.",
             0.6, 2.55, 12.1, 0.95, TEAL, WHITE, 11.5, 9.5)
    add_card(s, "L2: Thuật ngữ & Chủ đề (Terminology & Topics)",
             "Trích xuất định nghĩa pháp lý tự động ('Người lao động là...', 'Hợp đồng lao động là...'), liên kết từ đồng nghĩa để bắc cầu giữa câu hỏi người dùng và luật.",
             0.6, 3.65, 12.1, 0.95, PURPLE, WHITE, 11.5, 9.5)
    add_card(s, "L3: Tiền lương & Tiền thưởng (Wages & Rewards)",
             "Cấu thành thu nhập, lương tối thiểu vùng (Vùng I-IV), hình thức trả lương, phụ cấp, tỷ lệ hưởng làm thêm giờ (150%, 200%, 300%), công thức tính lương.",
             0.6, 4.75, 12.1, 0.95, AMBER, WHITE, 11.5, 9.5)
    add_card(s, "L4: Chủ thể & Quan hệ Lao động (Actors & Relations)",
             "Xác định chủ thể hành vi: Người lao động, Người sử dụng lao động, Tổ chức công đoàn, Cơ quan nhà nước; quyền hạn và nghĩa vụ tương ứng.",
             0.6, 5.85, 12.1, 0.95, ROSE, WHITE, 11.5, 9.5)

    # 25. 10-Layer Legal Knowledge Graph Breakdown (L5-L9)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Chi tiết Đồ thị", "Phân Tích Chi Tiết 5 Tầng Tri Thức Tiếp Theo (L5 - L9)", 25)
    add_card(s, "L5: Quy trình & Thủ tục (Procedures)",
             "Trình tự các bước: Thông báo trước, nộp hồ sơ, tham vấn ý kiến công đoàn, thời hạn giải quyết khiếu nại, cơ quan có thẩm quyền xử lý.",
             0.6, 1.45, 12.1, 0.95, DARK_EMERALD, WHITE, 11.5, 9.5)
    add_card(s, "L6: Thời gian & Thời hiệu (Time & Limitations)",
             "Mốc thời gian luật định: Thời hạn thử việc (30 ngày, 60 ngày, 180 ngày), thời hạn báo trước đơn phương chấm dứt HĐ (45 ngày, 30 ngày, 3 ngày), thời hiệu xử lý kỷ luật.",
             0.6, 2.55, 12.1, 0.95, TEAL, WHITE, 11.5, 9.5)
    add_card(s, "L7: Chế tài & Rủi ro (Sanctions & Risks)",
             "Hành vi vi phạm, khung tiền phạt vi phạm hành chính (Nghị định 12/2022/NĐ-CP), biện pháp khắc phục hậu quả, bồi thường thiệt hại.",
             0.6, 3.65, 12.1, 0.95, ROSE, WHITE, 11.5, 9.5)
    add_card(s, "L8: Vòng đời Hợp đồng Lao động (Legal Lifecycle)",
             "Chuỗi giai đoạn: Giao kết HĐ -> Thực hiện HĐ -> Thay đổi/Phụ lục HĐ -> Tạm hoãn HĐ -> Chấm dứt HĐLĐ -> Giải quyết quyền lợi sau chấm dứt.",
             0.6, 4.75, 12.1, 0.95, PURPLE, WHITE, 11.5, 9.5)
    add_card(s, "L9: Án lệ & Thực tiễn Xét xử (Case Law)",
             "Phán quyết tòa án, tình tiết cốt lõi và án lệ hướng dẫn áp dụng luật (sẵn sàng kết nối khi bổ sung kho dữ liệu bản án).",
             0.6, 5.85, 12.1, 0.95, AMBER, WHITE, 11.5, 9.5)

    # 26. Adaptive Routing & Complexity-Aware Retrieval
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Phân luồng Thích ứng", "Cơ Chế Phân Luồng Câu Hỏi Thích Ứng (Complexity-Aware Routing)", 26)
    add_card(s, "1. Luồng Single-Hop (Câu hỏi Đơn tầng)",
             "• Đặc điểm: Hỏi về 1 điều luật cụ thể, 1 con số (mức lương tối thiểu Vùng I, thời hạn thử việc).\n"
             "• Hành vi hệ thống: Bỏ qua đồ thị Neo4j, chạy thẳng qua Hybrid RAG (pgvector + GIN tsvector).\n"
             "• Hiệu quả: Tốc độ truy hồi cực nhanh (~15.8 ms), độ chính xác ngữ cảnh Context Precision đạt 0.7783.",
             0.6, 1.45, 5.9, 2.6, DARK_EMERALD)
    add_card(s, "2. Luồng Multi-Hop / Abstract (Câu hỏi Đa tầng & Tổng hợp)",
             "• Đặc điểm: Câu hỏi so sánh, câu hỏi tình huống có dẫn chiếu chéo giữa nhiều văn bản.\n"
             "• Hành vi hệ thống: Kích hoạt duyệt đồ thị Neo4j 2-hop từ các seed chunks để gom đầy đủ quan hệ.\n"
             "• Hiệu quả: Độ trung thực Faithfulness đạt 0.9213, hạn chế tối đa việc bỏ sót quy định chéo.",
             6.8, 1.45, 5.9, 2.6, TEAL)
    add_card(s, "Quy tắc Đánh giá Trade-off",
             "Không ép mọi câu hỏi phải đi qua đồ thị Neo4j. Adaptive Routing giúp cân bằng hoàn hảo giữa Tốc độ phục vụ (Latency) và Khả năng suy luận đa tầng (Multi-hop Reasoning).",
             0.6, 4.2, 12.1, 2.6, PURPLE, SOFT_GREEN, 12.5, 10.5)

    # 27. Corpus & Index Baseline
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Dữ liệu & Chỉ mục", "Hiện Trạng Kho Dữ Liệu & Chỉ Mục Pháp Luật (Corpus Baseline)", 27)
    add_stat_box(s, "74", "VĂN BẢN (73 DOCX, 1 DOC)", 0.6, 1.45, 2.85, 0.85, DARK_EMERALD)
    add_stat_box(s, "32,334", "CHUNKS TRUY HỒI", 3.68, 1.45, 2.85, 0.85, TEAL)
    add_stat_box(s, "3,309,530", "TỔNG SỐ TỪ (TOKENS)", 6.76, 1.45, 2.85, 0.85, PURPLE)
    add_stat_box(s, "102.35", "SỐ TỪ TB / CHUNK", 9.84, 1.45, 2.85, 0.85, AMBER)

    add_card(s, "Quy mô & Phân loại Văn bản",
             "• 1 Bộ luật Lao động (45/2019/QH14).\n"
             "• 7 Luật chuyên ngành (Luật BHXH, Luật Việc làm, Luật An toàn vệ sinh lao động...).\n"
             "• 42 Nghị định của Chính phủ hướng dẫn chi tiết.\n"
             "• 20 Thông tư của Bộ LĐ-TBXH và các Bộ ngành liên quan.\n"
             "• 4 Văn bản hợp nhất quy chuẩn.",
             0.6, 2.45, 5.9, 4.3, DARK_EMERALD, WHITE, 12, 10)
    add_card(s, "Đặc tính Kích thước Chunk",
             "• Min tokens: 4 tokens (đã loại bỏ rác/phân cách).\n"
             "• Max tokens: 440 tokens (giữ nguyên đơn vị pháp lý trước khi trượt).\n"
             "• Tỷ lệ Khoản + Điểm chiếm 72.56% tổng số chunks -> Phản ánh đúng bản chất chi tiết của văn bản pháp luật.",
             6.8, 2.45, 5.9, 4.3, TEAL, WHITE, 12, 10)

    # 28. Document & Chunk Distribution Analysis
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Phân bố Chunk", "Phân Tích Chi Tiết Phân Bố 9 Loại Chunk trong Hệ Thống", 28)
    chunk_dist = [
        ("Khoản (Clause Chunks)", "13,700 chunks (42.37%)", "Đơn vị chứa điều kiện, nghĩa vụ, quyền lợi cốt lõi."),
        ("Điểm (Point Chunks)", "9,762 chunks (30.19%)", "Quy định chi tiết, danh sách hành vi vi phạm, mức phạt."),
        ("Điều (Article Chunks)", "4,192 chunks (12.96%)", "Toàn văn điều luật cho tra cứu tổng quát và trích dẫn."),
        ("Bảng biểu (Table Chunks)", "1,248 chunks (3.86%)", "Bảng lương tối thiểu vùng, phụ cấp, hệ số đóng BHXH."),
        ("Khác (Intro, Structure, Sliding...)", "3,432 chunks (10.62%)", "Cửa sổ trượt bổ trợ cho điều luật dài, tiêu đề chương mục.")
    ]
    y_pos = 1.45
    for title, count, desc in chunk_dist:
        add_card(s, f"{title} - {count}", desc, 0.6, y_pos, 12.1, 0.95, DARK_EMERALD, WHITE, 11.5, 9.5)
        y_pos += 1.05

    # 29. [KEY QUESTION 5] Deterministic Document Parsing & Envelope Schema
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 5: PARSING", "Cấu Trúc Document khi Parsing & Phương Pháp Tất Định (Deterministic Parsing)", 29)
    add_card(s, "1. Tại sao Parsing Tất định (Không dùng LLM)?",
             "• Dùng LLM để parse dễ bị mất dòng, ảo giác số điều khoản và không bảo đảm tính nhất quán (reproducibility).\n"
             "• VlegalAI sử dụng State Machine kết hợp Regular Expressions đặc thù pháp lý Việt Nam: Duyệt tuần tự qua luồng khối (ordered block stream).\n"
             "• Trạng thái Parser: σ_i = (d, c_i, s_i, a_i, q_i, p_i) tương ứng Document > Chapter > Section > Article > Clause > Point.",
             0.6, 1.45, 5.9, 2.6, DARK_EMERALD)
    add_card(s, "2. Quy chuẩn Nhận dạng Header Regex",
             "• Chương: ^Chương\\s+([IVXLCDM]+|\\d+)\n"
             "• Mục: ^Mục\\s+([IVXLCDM]+|\\d+)\n"
             "• Điều: ^Điều\\s+(\\d+[a-zA-Z]?)\\s*[\\.:]\\s*(.+)\n"
             "• Khoản: ^(\\d{1,3})\\.\\s+(.+)\n"
             "• Điểm: ^([a-zđ](\\d+)?)\\)\\s+(.+)\n"
             "• Bảng biểu (Table): Giữ nguyên vị trí ngay dưới điều khoản chứa nó.",
             6.8, 1.45, 5.9, 2.6, TEAL)
    add_card(s, "3. Cấu trúc Document Envelope Output khi Parsing (JSON Schema 1.0)",
             "Mỗi văn bản được chuẩn hóa thành 1 Envelope chứa 4 mảng thực thể độc lập:\n"
             "• source: { path, filename, size_bytes, sha256 } -> Khóa bất biến ràng buộc toàn vẹn dữ liệu.\n"
             "• document: { doc_id, filename, title, code (45/2019/QH14), doc_type, issuer, text } -> Toàn văn và siêu dữ liệu.\n"
             "• nodes[]: Mảng các thực thể pháp lý { node_id, label, parent_id, path_label, text, ordinal, child_count }.\n"
             "• edges[]: Mảng quan hệ có hướng { edge_id, source_id, target_id, relation, evidence }.\n"
             "• chunks[]: Mảng đơn vị truy hồi { chunk_id, doc_id, node_id, chunk_type, title, citation, text, token_count }.",
             0.6, 4.2, 12.1, 2.6, PURPLE, SOFT_BLUE, 12.5, 10.5)

    # 30. [KEY QUESTION 1] Hierarchy-Aware Chunking & String Splitting
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 1: CHUNKING", "Chiến Lược Hierarchy-Aware Chunking & Thuật Toán Cắt Chuỗi (Sliding Window)", 30)
    add_stat_box(s, "W = 360", "ĐỘ DÀI CỬA SỔ (TỪ)", 0.6, 1.45, 2.85, 0.85, DARK_EMERALD)
    add_stat_box(s, "O = 70", "ĐỘ GỐI ĐẦU OVERLAP", 3.68, 1.45, 2.85, 0.85, TEAL)
    add_stat_box(s, "Δ = 290", "BƯỚC NHẢY STRIDE", 6.76, 1.45, 2.85, 0.85, PURPLE)
    add_stat_box(s, "N ≤ 440", "NGƯỠNG GIỮ NGUYÊN 1 CHUNK", 9.84, 1.45, 2.85, 0.85, AMBER)

    add_card(s, "1. Nguyên lý Phân đoạn theo Cấp bậc Pháp lý",
             "• Không cắt theo số ký tự cố định. Ranh giới phân đoạn đầu tiên là Đơn vị Cấu trúc Pháp lý (Legal Unit: Điều/Khoản/Điểm).\n"
             "• Mỗi chunk luôn gắn liền với mã định danh node (node_id) và đường dẫn phân cấp (path_label).",
             0.6, 2.45, 5.9, 2.1, DARK_EMERALD)
    add_card(s, "2. Thuật toán Cắt chuỗi & Cửa sổ trượt",
             "• Regex đếm từ: T = [0-9A-Za-zÀ-ỹĐđ]+ (an toàn tiếng Việt).\n"
             "• Nếu N(x) <= 440 từ: Giữ nguyên 1 chunk duy nhất.\n"
             "• Nếu N(x) > 440 từ: C_j = x[j*Δ : j*Δ + W] với bước nhảy Δ = 290 từ. Cửa sổ sau gắn nhãn sliding. Đuôi < 80 từ bị hủy.",
             6.8, 2.45, 5.9, 2.1, TEAL)
    add_card(s, "3. Hợp đồng Văn bản Embedding (Text Contract)",
             "Chuỗi đưa vào Embedding: X_c = title(c) + '\\n' + path_label(c) + '\\n' + chunk.text\n"
             "-> Đảm bảo một Khoản/Điểm ngắn vẫn mang đầy đủ ngữ cảnh của Điều luật và Tên văn bản chứa nó!",
             0.6, 4.7, 12.1, 2.15, PURPLE, SOFT_GREEN, 12.5, 10)

    # 31. [KEY QUESTION 4] Data Processing, Embedding Generation & Caching
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 4: EMBEDDING", "Quá Trình Xử Lý Dữ Liệu, Vector Embedding (1024D) & Caching SHA-256", 31)
    add_card(s, "1. Model & Chiều không gian Vector",
             "• Model: Vertex AI gemini-embedding-001\n"
             "• Chiều không gian: m = 1,024 chiều (float32).\n"
             "• Task Type: RETRIEVAL_DOCUMENT cho tài liệu, RETRIEVAL_QUERY cho câu hỏi người dùng.\n"
             "• Chuẩn hóa L2: e(x) = z(x) / ||z(x)||_2 -> Khoảng cách Cosine d_cos = 1 - e(q)^T * e(c).",
             0.6, 1.45, 5.9, 2.6, DARK_EMERALD)
    add_card(s, "2. Quản lý Dung lượng & Checkpoint Batch",
             "• Mỗi vector 1024D float32 = 4,096 bytes/chunk.\n"
             "• Tổng 32,334 chunks ≈ 126.3 MiB raw vector payload.\n"
             "• Checkpoint lưu theo batch 640 chunks vào PostgreSQL: Khi gặp sự cố mạng có thể resume ngay lập tức.",
             6.8, 1.45, 5.9, 2.6, TEAL)
    add_card(s, "3. Caching theo Mã băm Nội dung (Content Hash SHA-256)",
             "• Hash nội dung: h_c = SHA256(X_c).\n"
             "• Vector chỉ được tái sử dụng khi khớp toàn bộ 7 yếu tố: chunk_id + h_c + provider + model + revision + task_type + dimension (1024).\n"
             "• Bất kỳ sửa đổi nào trong văn bản luật sẽ kích hoạt tính lại vector tự động, ngăn ngừa stale embedding tuyệt đối.",
             0.6, 4.2, 12.1, 2.6, PURPLE, SOFT_BLUE, 12.5, 10.5)

    # 32. [KEY QUESTION 1] Multi-Store Indexing & Hybrid RRF
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 1: INDEXING", "Hệ Thống Indexing Đa Tầng Đồng Bộ & Hợp Nhất Thứ Hạng Weighted RRF", 32)
    add_card(s, "1. pgvector HNSW Index",
             "Cột vector(1024). Thuật toán HNSW (vector_cosine_ops, M=16, ef_construction=64) cho truy hồi tương đồng ngữ nghĩa cực nhanh.",
             0.6, 1.45, 3.85, 2.6, DARK_EMERALD)
    add_card(s, "2. PostgreSQL GIN tsvector",
             "Cột generated tsvector trên [title, citation, text]. Khớp số hiệu văn bản (45/2019/QH14), số điều luật, thuật ngữ chuyên ngành.",
             4.75, 1.45, 3.85, 2.6, TEAL)
    add_card(s, "3. Neo4j Graph Index",
             "29,575 nodes & 108,368 edges. Unique Constraint trên LegalChunk(chunk_id), Fulltext Index trên [title, citation, text].",
             8.9, 1.45, 3.85, 2.6, PURPLE)
    add_card(s, "Hợp nhất Thứ hạng RRF (Reciprocal Rank Fusion) tại Query Time",
             "• Kết hợp điểm số Vector và BM25 bằng công thức Weighted RRF (K = 60):\n"
             "  S_RRF(c) = w_v * (K + 1) / (K + r_v(c)) + w_b * (K + 1) / (K + r_b(c)) với trọng số thực nghiệm w_v = 0.55 (Vector) và w_b = 0.45 (BM25).\n"
             "• Điều chỉnh điểm số cơ sở theo mức độ bao phủ từ khóa: S_0(c) = S_RRF(c) / (r_f(c)^0.35) + 0.9 * m_c / min(|T_q|, 10) + B(c, q).\n"
             "• Atomic Activation: Chỉ kích hoạt phục vụ khi tổng số chunks và hash đồng bộ 100% giữa PostgreSQL và Neo4j.",
             0.6, 4.2, 12.1, 2.65, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    # 33. [KEY QUESTION 3] Model Strategy: Why VlegalAI Does NOT Fine-Tune LLM
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 3: MÔ HÌNH", "Có Train / Fine-Tune LLM hay không? Quyết Định Kỹ Thuật Cốt Lõi (ADR-04)", 33)
    box_no = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.45), Inches(3.2), Inches(5.3))
    box_no.fill.solid()
    box_no.fill.fore_color.rgb = DARK_BG
    box_no.line.fill.background()
    tf_no = box_no.text_frame
    tf_no.word_wrap = True
    p_no0 = tf_no.paragraphs[0]
    p_no0.text = "KHÔNG"
    p_no0.alignment = PP_ALIGN.CENTER
    p_no0.font.size = Pt(36)
    p_no0.font.bold = True
    p_no0.font.color.rgb = ROSE
    p_no1 = tf_no.add_paragraph()
    p_no1.text = "Fine-tune hay Train lại Trọng số LLM\n(Trọng số mô hình giữ cố định: θ' = θ)"
    p_no1.alignment = PP_ALIGN.CENTER
    p_no1.font.size = Pt(12)
    p_no1.font.bold = True
    p_no1.font.color.rgb = WHITE
    p_no2 = tf_no.add_paragraph()
    p_no2.text = "\n• Không Pretraining\n• Không SFT\n• Không LoRA / QLoRA\n• Dùng Hosted Gemini 2.5 Flash qua Vertex AI"
    p_no2.font.size = Pt(10.5)
    p_no2.font.color.rgb = TEXT_LIGHT

    add_card(s, "1. Khả năng Dẫn chứng Nguồn gốc (Provenance)",
             "LLM fine-tuned lưu kiến thức vào 'hộp đen' trọng số, không thể bảo đảm trích dẫn chính xác tuyệt đối từng Điều/Khoản. RAG giữ nguyên vẹn chuỗi bằng chứng tường minh: Answer > Evidence > Chunk > Provision > Instrument.",
             4.0, 1.45, 4.3, 2.55, DARK_EMERALD)
    add_card(s, "2. Tính Cập nhật & Luật Hết hiệu lực (Freshness)",
             "Pháp luật thay đổi liên tục (Nghị định mới thay thế cũ). Nếu fine-tune, mỗi lần đổi luật phải thu thập dữ liệu & train lại rất tốn kém và dễ lẫn lộn. Với GraphRAG, chỉ cần re-index lại DB trong vài phút.",
             8.5, 1.45, 4.2, 2.55, TEAL)
    add_card(s, "3. Kiểm soát Ảo giác Tuyệt đối (Hallucination Control)",
             "Lĩnh vực pháp lý đòi hỏi sự chính xác tuyệt đối. Việc ép mô hình chỉ sinh câu trả lời dựa trên tập bằng chứng được truy hồi (Evidence-first prompt contract) giúp triệt tiêu hoàn toàn nguy cơ bịa đặt số điều luật.",
             4.0, 4.2, 4.3, 2.55, PURPLE)
    add_card(s, "4. Tối ưu hóa Trọng tâm Kỹ thuật",
             "Thay vì phân tán tài nguyên để train model, nhóm tập trung tối ưu: Cấu trúc đồ thị 10 tầng, Chiến lược Chunking phân cấp, Thuật toán Hybrid RRF, Prompt Engineering và Bộ lọc Kiểm chứng Trích dẫn (Evidence Gate).",
             8.5, 4.2, 4.2, 2.55, AMBER)

    # 34. Experimental Architecture Definitions
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Thiết kế Thực nghiệm", "Bốn Kiến Trúc So Sánh Đối Đầu trong Thực Nghiệm RAGAS", 34)
    add_card(s, "1. Dense RAG (Baseline Vector)",
             "• Truy hồi thuần túy bằng vector embedding qua gemini-embedding-001 (1024D) trên pgvector.\n"
             "• Lấy Top-k = 4 chunks có khoảng cách cosine nhỏ nhất đưa thẳng vào Gemini 2.5 Flash.\n"
             "• Đóng vai trò làm baseline chuẩn so sánh tốc độ và độ chính xác ngữ nghĩa.",
             0.6, 1.45, 5.9, 2.55, DARK_EMERALD)
    add_card(s, "2. LightRAG (Hybrid + 1-hop Graph)",
             "• Kết hợp Dense Vector + FTS5 BM25 qua Reciprocal Rank Fusion.\n"
             "• Mở rộng 1-hop trên đồ thị tri thức để bổ sung các thực thể lân cận trực tiếp.\n"
             "• Đóng vai trò là kiến trúc lai nhẹ cân bằng giữa từ khóa và quan hệ.",
             6.8, 1.45, 5.9, 2.55, TEAL)
    add_card(s, "3. GraphRAG (Hybrid Seed + 2-hop Graph)",
             "• Lấy hạt giống từ Hybrid Retrieval, sau đó mở rộng đồ thị 2-hop trên Neo4j.\n"
             "• Duyệt qua các quan hệ ngữ nghĩa: DẪN_CHIẾU_ĐẾN, QUY_ĐỊNH_TẠI, CÓ_KHOẢN...\n"
             "• Đóng vai trò là kiến trúc suy luận sâu cho các câu hỏi quan hệ phức tạp.",
             0.6, 4.15, 5.9, 2.65, PURPLE)
    add_card(s, "4. RAG + GraphRAG (Combined Pipeline)",
             "• Kết hợp bằng chứng hạt giống của RAG với đồ thị mở rộng có trọng số.\n"
             "• Ghi nhận 100 dòng request latency và checkpoint đánh giá từng phần.",
             6.8, 4.15, 5.9, 2.65, AMBER)

    # 35. [KEY QUESTION 2] Evaluation Methodology: Accuracy as an 8-D Vector
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 2: ĐÁNH GIÁ", "Phương Pháp Đánh Giá Độ Chính Xác: Vector Chất Lượng 8 Chiều & RAGAS", 35)
    add_card(s, "Độ chính xác trong Pháp lý là một Vector Đa chiều",
             "Q = (Q_ret, Q_faith, Q_rel, Q_fact, Q_cite, Q_facet, Q_temp, Q_lat)\n"
             "Exact match chuỗi không phù hợp cho pháp lý vì 2 câu trả lời khác từ ngữ vẫn có thể cùng biểu đạt đúng 1 quy định luật.",
             0.6, 1.45, 12.1, 1.1, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    add_card(s, "1. FAITHFULNESS (Tính trung thực)",
             "Tỷ lệ các phát biểu trong câu trả lời được chứng minh trực tiếp bởi ngữ cảnh trích dẫn (|S_q| / |A_q|). Triệt tiêu ảo giác.",
             0.6, 2.7, 2.85, 2.2, DARK_EMERALD)
    add_card(s, "2. ANSWER RELEVANCY",
             "Mức độ câu trả lời giải quyết đúng và trúng câu hỏi của người dùng, không trả lời lan man hoặc lạc đề.",
             3.68, 2.7, 2.85, 2.2, TEAL)
    add_card(s, "3. CONTEXT PRECISION & RECALL",
             "Độ chính xác và độ bao phủ của tập chunk trích xuất so với tập bằng chứng chuẩn (Ground Truth).",
             6.76, 2.7, 2.85, 2.2, PURPLE)
    add_card(s, "4. FACTUAL CORRECTNESS",
             "Độ chính xác về mặt sự thật và kết luận pháp lý so với câu trả lời mẫu của chuyên gia luật.",
             9.84, 2.7, 2.85, 2.2, AMBER)

    add_card(s, "Context ID Precision & Recall",
             "Đo lường mức độ trùng khớp chính xác mã định danh điều khoản pháp lý: R_ID = |R_k ∩ G| / |G| và P_ID = |R_k ∩ G| / |R_k|.",
             0.6, 5.05, 5.9, 1.75, DARK_EMERALD)
    add_card(s, "Bộ Benchmark RAGAS 100 Câu hỏi Thực tế",
             "100 câu hỏi luật lao động thực tế chia làm 3 nhóm: 50 Single-hop, 25 Multi-hop Specific, 25 Multi-hop Abstract. So sánh trực diện 4 kiến trúc: Dense RAG, LightRAG, GraphRAG, RAG+GraphRAG.",
             6.8, 5.05, 5.9, 1.75, TEAL)

    # 36. [KEY QUESTION 2] Aggregate Quality Results Table
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 2: KẾT QUẢ", "Bảng Kết Quả So Sánh Thực Nghiệm 4 Kiến Trúc (RAGAS Benchmark 100)", 36)
    table_shape = s.shapes.add_table(5, 9, Inches(0.6), Inches(1.45), Inches(12.1), Inches(2.2))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    for i in range(1, 9):
        table.columns[i].width = Inches(1.28)

    headers = ["Kiến trúc", "Faithfulness", "Ans. Rel.", "Ctx. Prec.", "Ctx. Rec.", "Factual", "ID Recall", "Latency", "Overall"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE

    rows_data = [
        ["Dense RAG", "0.9111", "0.9126", "0.7783", "0.8154", "0.3985", "0.6800", "15.8 ms", "0.7619*"],
        ["LightRAG", "0.9160", "0.8951", "0.6950", "0.8124", "0.4271", "0.6650", "148.8 ms", "0.7481"],
        ["GraphRAG", "0.9213*", "0.8641", "0.5558", "0.7820", "0.4003", "0.5750", "146.7 ms", "0.7041"],
        ["RAG+GraphRAG", "0.9403#", "--", "--", "--", "--", "0.6850*", "162.6 ms", "Checkpoint"]
    ]
    for row_idx, r in enumerate(rows_data):
        for col_idx, val in enumerate(r):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT_GREEN if row_idx == 0 else (SOFT_BLUE if row_idx == 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.bold = (col_idx == 0 or "*" in val)
            p.font.color.rgb = DARK_EMERALD if "*" in val else TEXT_DARK

    add_card(s, "Nhận xét & Đánh giá Khoa học từ Dữ liệu Thực nghiệm",
             "• GraphRAG đạt độ trung thực Faithfulness cao nhất (0.9213) nhờ khả năng bao quát ngữ cảnh cấu trúc quan hệ, hạn chế tối đa ảo giác.\n"
             "• Dense RAG đạt Context Precision (0.7783) và Tốc độ truy hồi tốt nhất (15.8ms vs ~147ms của Graph) cho các câu hỏi tra cứu đơn tầng (Single-hop).\n"
             "• Minh chứng cho quyết định thiết kế Adaptive Routing: Dùng Hybrid Dense RAG làm mặc định cho câu hỏi đơn giản và chỉ kích hoạt GraphRAG cho câu hỏi phức tạp.",
             0.6, 3.9, 12.1, 2.9, DARK_EMERALD, SOFT_GREEN, 12.5, 10.5)

    # 37. [KEY QUESTION 2] Performance Stratification by Question Complexity
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "TRỌNG TÂM 2: PHÂN TÍCH", "Hiệu Năng theo Phân Loại Độ Phức Tạp Câu Hỏi (Single-hop vs Multi-hop)", 37)
    add_card(s, "1. Single-Hop Specific (50 câu)",
             "Điểm RAGAS: Dense 0.8114 | Light 0.8039 | Graph 0.7458\n\n"
             "• Tra cứu trực tiếp 1 điều luật, 1 mức lương tối thiểu, 1 thời hạn cụ thể.\n"
             "• Dense RAG chiếm ưu thế tuyệt đối về tốc độ và độ chính xác ngữ cảnh do không bị nhiễu bởi các node quan hệ xa.",
             0.6, 1.45, 3.85, 3.6, DARK_EMERALD)
    add_card(s, "2. Multi-Hop Specific (25 câu)",
             "Điểm RAGAS: Dense 0.7625 | Light 0.7283 | Graph 0.6914\n\n"
             "• Câu hỏi kết hợp 2-3 điều luật (VD: Điều kiện hưởng lương + Trách nhiệm người sử dụng lao động).\n"
             "• Đồ thị hỗ trợ tìm kiếm đường dẫn chiếu chéo hiệu quả.",
             4.75, 1.45, 3.85, 3.6, TEAL)
    add_card(s, "3. Multi-Hop Abstract (25 câu)",
             "Điểm RAGAS: Dense 0.6624 | Light 0.6563 | Graph 0.6334\n\n"
             "• Câu hỏi tình huống tổng hợp, so sánh quyền lợi giữa các nhóm lao động.\n"
             "• GraphRAG có mức độ suy giảm điểm số ít nhất (-0.112 vs -0.149 của Dense), chứng minh khả năng duy trì thông tin đa tầng.",
             8.9, 1.45, 3.85, 3.6, PURPLE)
    add_card(s, "Kết luận về Bài toán Đánh đổi (Trade-Off): Độ chính xác vs Thời gian phản hồi",
             "• Đồ thị tri thức (Neo4j) giúp mở rộng quan hệ pháp lý rất tốt nhưng làm tăng độ trễ truy hồi từ 15.8ms lên ~147ms.\n"
             "• Thời gian sinh câu trả lời của LLM chiếm đa số (5.7s - 8.3s).\n"
             "• Thiết kế VlegalAI kết hợp: Seed Hybrid Retrieval + Bounded Graph Expansion (2-hop) là tối ưu nhất cho môi trường sản phẩm thực tế.",
             0.6, 5.25, 12.1, 1.6, AMBER, SOFT_AMBER, 12, 10)

    # 38. Request-Level Latency Baseline
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Độ trễ Hệ thống", "Đo Lường Thực Nghiệm Độ Trễ (Retrieval & End-to-End Latency)", 38)
    add_stat_box(s, "15.8 ms", "DENSE RETRIEVAL (MEAN)", 0.6, 1.45, 2.85, 0.85, DARK_EMERALD)
    add_stat_box(s, "146.7 ms", "GRAPH RETRIEVAL (MEAN)", 3.68, 1.45, 2.85, 0.85, TEAL)
    add_stat_box(s, "5.72 s", "DENSE E2E (MEAN)", 6.76, 1.45, 2.85, 0.85, PURPLE)
    add_stat_box(s, "6.13 s", "GRAPH E2E (MEAN)", 9.84, 1.45, 2.85, 0.85, AMBER)

    add_card(s, "Phân tích Thành phần Độ trễ",
             "• Độ trễ truy hồi (Retrieval Latency): Dense RAG (15.8ms) nhanh hơn ~9.3 lần so với GraphRAG (146.7ms) và LightRAG (148.8ms) do không tốn thời gian duyệt network graph qua Neo4j.\n"
             "• Độ trễ tạo sinh (Generation Latency): Chiếm 95-98% tổng thời gian E2E (5.7s - 8.3s) do phụ thuộc vào tốc độ sinh token của Gemini 2.5 Flash trên Vertex AI Cloud.\n"
             "• Độ trễ đuôi (p95 Tail Latency): Đạt ~25.4s - 48.7s ở các câu hỏi phức tạp có số lượng token dài.",
             0.6, 2.45, 12.1, 4.4, DARK_EMERALD, WHITE, 12.5, 10.5)

    # 39. Paired Bootstrap 95% Confidence Intervals
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Phân tích Thống kê", "Khoảng Tin Cậy Ghép Cặp Bootstrap 95% (Paired Bootstrap 95% CI)", 39)
    add_card(s, "1. Dense RAG vs LightRAG",
             "• Chênh lệch trung bình: +0.0138\n• 95% CI: [-0.0144, +0.0423]\n• Kết luận: Khoảng tin cậy chứa 0 -> Chênh lệch chưa có ý nghĩa thống kê rõ rệt.",
             0.6, 1.45, 3.85, 3.0, DARK_EMERALD)
    add_card(s, "2. Dense RAG vs GraphRAG",
             "• Chênh lệch trung bình: +0.0578\n• 95% CI: [+0.0270, +0.0916]\n• Kết luận: Khoảng tin cậy hoàn toàn dương -> Dense RAG cao hơn có ý nghĩa thống kê về điểm Overall.",
             4.75, 1.45, 3.85, 3.0, TEAL)
    add_card(s, "3. LightRAG vs GraphRAG",
             "• Chênh lệch trung bình: +0.0440\n• 95% CI: [+0.0092, +0.0816]\n• Kết luận: Khoảng tin cậy hoàn toàn dương -> LightRAG cao hơn có ý nghĩa thống kê về điểm Overall.",
             8.9, 1.45, 3.85, 3.0, PURPLE)
    add_card(s, "Ý nghĩa Khoa học",
             "Phương pháp Bootstrap bảo toàn cấu trúc ghép cặp từng câu hỏi (paired question indices), mang lại bằng chứng định lượng khách quan và đáng tin cậy.",
             0.6, 4.65, 12.1, 2.2, DARK_EMERALD, SOFT_GREEN, 12, 10)

    # 40. Checkpoint Completeness & Missing Value Handling
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Tính Toàn vẹn", "Tính Toàn Vẹn Dữ Liệu Thực Nghiệm & Xử Lý Giá Trị Khuyết", 40)
    add_card(s, "Nguyên tắc Trung thực trong Báo cáo Kết quả (Honest Evaluation)",
             "• 3 kiến trúc chính (Dense RAG, LightRAG, GraphRAG): Hoàn thành 100/100 câu hỏi trên toàn bộ các chỉ số.\n"
             "• Nhánh thực nghiệm RAG+GraphRAG: Checkpoint ghi nhận 100 dòng request latency và ID Recall/Precision, nhưng chỉ có 45 dòng Faithfulness (nhóm single-hop) do lỗi quota evaluator.\n"
             "• Xử lý giá trị khuyết: TUYỆT ĐỐI KHÔNG tự điền số 0 hoặc số giả lập. Báo cáo minh bạch dữ liệu thực tế và loại nhánh RAG+GraphRAG khỏi bảng xếp hạng tổng thể.",
             0.6, 1.45, 12.1, 5.4, DARK_EMERALD, SOFT_GREEN, 13, 11)

    # 41. Benchmark Error Analysis & Observed Symptoms
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Phân tích Lỗi", "Phân Tích Lỗi Thực Nghiệm & Các Hiện Tượng Quan Sát Được", 41)
    errs = [
        ("Identifier Precision thấp (0.210 - 0.243)", "Tập chunk truy hồi thường bao gồm thêm các điều khoản bối cảnh hữu ích nhưng không nằm trong Ground Truth tối giản."),
        ("Graph Context Precision giảm khi mở rộng", "Duyệt 2-hop trên đồ thị đưa vào một số node lân cận không trực tiếp liên quan đến câu hỏi đơn tầng."),
        ("Điểm Multi-hop Abstract thấp nhất", "Các câu hỏi tổng hợp trừu tượng đòi hỏi gom nhiều điều luật từ nhiều nghị định khác nhau, dễ bị loãng context."),
        ("Độ trễ p95 đuôi dài (25s - 48s)", "Do phụ thuộc vào thời gian suy luận và sinh token dài của Cloud API Vertex AI khi gặp câu hỏi phức.")
    ]
    y_pos = 1.45
    for title, desc in errs:
        add_card(s, title, desc, 0.6, y_pos, 12.1, 1.15, ROSE, WHITE, 11.5, 9.5)
        y_pos += 1.3

    # 42. System Verification Evidence (492 Backend Tests)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Minh chứng Kiểm thử", "Hệ Thống Kiểm Thử Tự Động: 492 Backend Tests & 16 Frontend Tests", 42)
    add_stat_box(s, "492 / 492", "BACKEND TESTS PASS", 0.6, 1.45, 2.85, 0.85, DARK_EMERALD)
    add_stat_box(s, "16 / 16", "FRONTEND TESTS PASS", 3.68, 1.45, 2.85, 0.85, TEAL)
    add_stat_box(s, "18", "ALEMBIC MIGRATIONS", 6.76, 1.45, 2.85, 0.85, PURPLE)
    add_stat_box(s, "100%", "BUILD PASS (VITE & DOCKER)", 9.84, 1.45, 2.85, 0.85, AMBER)

    add_card(s, "Phạm vi Bao phủ của Test Suite",
             "• Unit Tests: Kiểm tra Deterministic Parser, Thuật toán Chunking W=360 O=70, Trích xuất Header Regex, Tính toán RRF.\n"
             "• Integration Tests: Kiểm tra API endpoints (Auth, Chat, Contracts, Laws), pgvector HNSW search, Neo4j graph traversal.\n"
             "• Security Tests: Kiểm tra mã hóa AES-GCM, Token đính kèm, xác thực Google OIDC và cách ly dữ liệu người dùng.",
             0.6, 2.45, 12.1, 4.4, DARK_EMERALD, WHITE, 12.5, 10.5)

    # 43. Answers to Research Questions (RQ1 - RQ4)
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Giải đáp RQ", "Tổng Kết Lời Giải Cho 4 Câu Hỏi Nghiên Cứu (RQ1 - RQ4)", 43)
    add_card(s, "RQ1 (Chất lượng)",
             "GraphRAG đạt Faithfulness cao nhất (0.9213) giúp chống ảo giác tốt nhất. Dense RAG dẫn đầu về điểm Overall RAGAS (0.7619) nhờ Context Precision cao ở câu hỏi đơn tầng.",
             0.6, 1.45, 5.9, 2.55, DARK_EMERALD)
    add_card(s, "RQ2 (Độ trễ)",
             "Dense RAG nhanh hơn ~9.3 lần so với GraphRAG ở khâu truy hồi (15.8ms vs 146.7ms). Điều này chứng minh thiết kế Adaptive Routing là hoàn toàn đúng đắn.",
             6.8, 1.45, 5.9, 2.55, TEAL)
    add_card(s, "RQ3 (Độ sâu suy luận)",
             "GraphRAG có mức độ suy giảm điểm số ít nhất (-0.112) khi chuyển từ Single-hop sang Multi-abstract, chứng minh ưu thế của đồ thị khi câu hỏi cần tổng hợp đa điều luật.",
             0.6, 4.15, 5.9, 2.65, PURPLE)
    add_card(s, "RQ4 (Tính tái lập)",
             "Đạt tính tái lập 100% ở cấp độ phần mềm (1 commit, 492 tests pass, 1 migration head, Docker image bất biến, mã nguồn và benchmark mở).",
             6.8, 4.15, 5.9, 2.65, AMBER)

    # 44. Threats to Validity, Limitations & Roadmap
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Hạn chế & Hướng mở", "Các Mối Đe Dọa Giá Trị, Hạn Chế Hiện Tại & Kế Hoạch Tương Lai", 44)
    add_card(s, "Hạn chế Hiện tại",
             "• Độ trễ đuôi p95 còn phụ thuộc vào Cloud API của Vertex AI.\n"
             "• Bộ dữ liệu đánh giá 100 câu hỏi cần được mở rộng thêm thẩm định chéo của nhiều chuyên gia luật độc lập.\n"
             "• Tầng Án lệ (Layer 9) mới ở dạng khung thiết kế, chưa nạp toàn bộ kho bản án thực tế.",
             0.6, 1.45, 5.9, 4.3, AMBER, WHITE, 12, 10)
    add_card(s, "Kế hoạch Phát triển Tương lai (Roadmap)",
             "1. Huấn luyện mô hình Reranker cục bộ để giảm độ trễ truy hồi và lọc bớt nhiễu đồ thị.\n"
             "2. Nạp thêm kho Án lệ Tòa án và Quyết định xử phạt vi phạm hành chính thực tế vào Layer 9.\n"
             "3. Mở rộng kho tri thức sang Luật Doanh nghiệp, Luật Đầu tư và Luật Thuế.",
             6.8, 1.45, 5.9, 4.3, TEAL, WHITE, 12, 10)
    add_card(s, "Tuyên bố Giá trị",
             "VlegalAI đã chứng minh tính khả thi và hiệu quả vượt trội của việc kết hợp Đồ thị Tri thức và Hybrid RAG trong bài toán pháp lý chuyên sâu.",
             0.6, 5.9, 12.1, 0.95, DARK_EMERALD, SOFT_GREEN, 11, 9.5)

    # 45. Conclusion & Defense Q&A
    s = prs.slides.add_slide(blank_layout)
    bg45 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg45.fill.solid()
    bg45.fill.fore_color.rgb = DARK_BG
    bg45.line.fill.background()

    tb_end = s.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.3))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True
    p_end0 = tf_end.paragraphs[0]
    p_end0.text = "CẢM ƠN QUÝ THẦY CÔ TRONG HỘI ĐỒNG"
    p_end0.alignment = PP_ALIGN.CENTER
    p_end0.font.size = Pt(26)
    p_end0.font.bold = True
    p_end0.font.color.rgb = WHITE
    p_end1 = tf_end.add_paragraph()
    p_end1.text = "NHÓM VLEGALAI SẴN SÀNG LẮNG NGHE Ý KIẾN ĐÓNG GÓP & TRẢ LỜI CÂU HỎI"
    p_end1.alignment = PP_ALIGN.CENTER
    p_end1.font.size = Pt(13.5)
    p_end1.font.bold = True
    p_end1.font.color.rgb = TEAL

    qa_list = [
        ("1. Chunking & Cắt chuỗi", "Hierarchy-Aware theo Điều/Khoản/Điểm. Cửa sổ trượt W=360, Overlap=70 từ; đếm từ bằng regex tiếng Việt."),
        ("2. Có Train LLM Không?", "KHÔNG train/fine-tune weight. Dùng Hosted Gemini 2.5 Flash + RAG/GraphRAG để giữ 100% provenance & cập nhật luật mới."),
        ("3. Đánh Giá Độ Chính Xác", "Vector 8 chiều qua RAGAS (100 câu hỏi luật thực tế): Faithfulness đạt 0.9213, ID Recall đạt 0.6800."),
        ("4. Xử Lý & Vector Embedding", "Input contract: Title + Path + Text. Model gemini-embedding-001 (1024D), chuẩn hóa L2, caching SHA-256."),
        ("5. Cấu Trúc Document Parsing", "Deterministic State Machine phân tích cú pháp DOCX thành Envelope 4 mảng: source, document, nodes, edges, chunks.")
    ]

    x_positions = [0.8, 4.8, 8.8, 2.8, 6.8]
    y_positions = [2.4, 2.4, 2.4, 4.5, 4.5]
    for idx, (q_title, q_desc) in enumerate(qa_list):
        box_qa = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(y_positions[idx]), Inches(3.7), Inches(1.8))
        box_qa.fill.solid()
        box_qa.fill.fore_color.rgb = CARD_DARK
        box_qa.line.color.rgb = EMERALD
        tf_q = box_qa.text_frame
        tf_q.word_wrap = True
        tf_q.margin_left = tf_q.margin_top = tf_q.margin_right = tf_q.margin_bottom = Inches(0.12)
        pq0 = tf_q.paragraphs[0]
        pq0.text = q_title
        pq0.font.size = Pt(11.5)
        pq0.font.bold = True
        pq0.font.color.rgb = AMBER
        pq1 = tf_q.add_paragraph()
        pq1.text = q_desc
        pq1.font.size = Pt(9.2)
        pq1.font.color.rgb = TEXT_LIGHT

    tb_foot = s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4))
    pf = tb_foot.text_frame.paragraphs[0]
    pf.text = "VlegalAI: Trợ lý Pháp lý Lao động Việt Nam | FPT University Quy Nhơn 2026"
    pf.alignment = PP_ALIGN.CENTER
    pf.font.size = Pt(9.5)
    pf.font.color.rgb = TEXT_MUTED

    # Save to slides.pptx and sldies.pptx
    print(f"Saving {len(prs.slides)} slides to {OUT_SLIDES}...")
    prs.save(str(OUT_SLIDES))
    print(f"Saving {len(prs.slides)} slides to {OUT_SLDIES}...")
    prs.save(str(OUT_SLDIES))
    print("Successfully created full 45-slide presentation deck!")

if __name__ == "__main__":
    build_45_slides()
